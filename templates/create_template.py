#!/usr/bin/env python3
"""Create a base PPTX template with sample slides for each layout."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return shape

def add_shape(slide, shape_type, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

blank_layout = prs.slide_layouts[6]

# === 1. Cover ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "cover"
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
          prs.slide_width, prs.slide_height, DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.2),
          Inches(0.15), Inches(1.8), ACCENT_BLUE)
add_textbox(slide, Inches(1.1), Inches(3.0), Inches(10), Inches(1.2),
            "__TITLE__", font_size=44, bold=True, color=WHITE)
add_textbox(slide, Inches(1.1), Inches(4.3), Inches(10), Inches(0.6),
            "__SUBTITLE__", font_size=20, color=RGBColor(0xAA, 0xBB, 0xCC))
add_textbox(slide, Inches(1.1), Inches(5.0), Inches(10), Inches(0.5),
            "__DATE__", font_size=14, color=RGBColor(0x88, 0x88, 0x88))

# === 2. TOC ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "toc"
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
            "__TITLE__", font_size=32, bold=True, color=DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4),
          Inches(2), Inches(0.06), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
            "__ITEMS__", font_size=18, color=DARK_TEXT)

# === 3. Title + Content ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "title_content"
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            "__TITLE__", font_size=28, bold=True, color=DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
          Inches(1.5), Inches(0.05), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
            "__CONTENT__", font_size=16, color=DARK_TEXT)

# === 4. Two Column ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "two_column"
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            "__TITLE__", font_size=28, bold=True, color=DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
          Inches(1.5), Inches(0.05), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.3), Inches(5.5),
            "__LEFT__", font_size=16, color=DARK_TEXT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.5),
          Inches(0.02), Inches(5.5), LIGHT_GRAY)
add_textbox(slide, Inches(6.7), Inches(1.5), Inches(5.3), Inches(5.5),
            "__RIGHT__", font_size=16, color=DARK_TEXT)

# === 5. Chart ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "chart"
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            "__TITLE__", font_size=28, bold=True, color=DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
          Inches(1.5), Inches(0.05), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
            "__DESC__", font_size=14, color=MID_GRAY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.4),
          Inches(11.5), Inches(4.5), LIGHT_GRAY)
add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(4.5),
            "__CHART__", font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA),
            align=PP_ALIGN.CENTER)

# === 6. Closing ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "closing"
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
          prs.slide_width, prs.slide_height, DARK_NAVY)
add_textbox(slide, Inches(0), Inches(3.0), prs.slide_width, Inches(1.0),
            "__TITLE__", font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(4.2), prs.slide_width, Inches(0.6),
            "__SUBTITLE__", font_size=20, color=RGBColor(0xAA, 0xBB, 0xCC),
            align=PP_ALIGN.CENTER)

# === 7. Table ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "table"
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            "__TITLE__", font_size=28, bold=True, color=DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
          Inches(1.5), Inches(0.05), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
            "__DESC__", font_size=14, color=MID_GRAY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.4),
          Inches(11.5), Inches(4.5), LIGHT_GRAY)
add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(4.5),
            "__TABLE__", font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA),
            align=PP_ALIGN.CENTER)

# === 8. Timeline ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "timeline"
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            "__TITLE__", font_size=28, bold=True, color=DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
          Inches(1.5), Inches(0.05), ACCENT_BLUE)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.0),
          Inches(11.5), Inches(0.04), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0),
            "__TIMELINE__", font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA),
            align=PP_ALIGN.CENTER)

# === 9. Image + Content ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "image_content"
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            "__TITLE__", font_size=28, bold=True, color=DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
          Inches(1.5), Inches(0.05), ACCENT_BLUE)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5),
          Inches(5.5), Inches(5.2), LIGHT_GRAY)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2),
            "__IMAGE__", font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA),
            align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(6.7), Inches(1.5), Inches(5.6), Inches(5.2),
            "__CONTENT__", font_size=16, color=DARK_TEXT)

# === 10. Quote ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "quote"
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
          prs.slide_width, prs.slide_height, DARK_NAVY)
add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
            "__QUOTE__", font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
            "__AUTHOR__", font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC),
            align=PP_ALIGN.CENTER)

# === 11. Team ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "team"
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            "__TITLE__", font_size=28, bold=True, color=DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
          Inches(1.5), Inches(0.05), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.2),
            "__TEAM__", font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA),
            align=PP_ALIGN.CENTER)

# === 12. Data Highlight ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "data_highlight"
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            "__TITLE__", font_size=28, bold=True, color=DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
          Inches(1.5), Inches(0.05), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.0),
            "__BIG_NUMBER__", font_size=72, bold=True, color=ACCENT_BLUE,
            align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.0),
            "__LABEL__", font_size=20, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# === 13. Process ===
slide = prs.slides.add_slide(blank_layout)
slide.name = "process"
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            "__TITLE__", font_size=28, bold=True, color=DARK_NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
          Inches(1.5), Inches(0.05), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.2),
            "__PROCESS__", font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA),
            align=PP_ALIGN.CENTER)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
prs.save(os.path.join(BASE_DIR, 'base_template.pptx'))
print("Template saved with 13 layouts:")
for s in prs.slides:
    print(f"  - {s.name}")
