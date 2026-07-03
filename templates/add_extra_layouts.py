#!/usr/bin/env python3
"""Add extra layout slides to existing template."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
prs = Presentation(os.path.join(BASE_DIR, 'base_template.pptx'))

DARK_NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return shape

def add_shape(slide, shape_type, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

blank_layout = prs.slide_layouts[6]

# Extra layouts that reuse chart area placeholder
for name in ["waterfall", "funnel", "gantt", "wordcloud"]:
    slide = prs.slides.add_slide(blank_layout)
    slide.name = name
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                "__TITLE__", font_size=28, bold=True, color=DARK_NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
              Inches(1.5), Inches(0.05), ACCENT_BLUE)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
                "__DESC__", font_size=14, color=MID_GRAY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.1),
              Inches(11.5), Inches(4.8), LIGHT_GRAY)
    add_textbox(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.8),
                "__CHART__", font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA),
                align=PP_ALIGN.CENTER)

prs.save(os.path.join(BASE_DIR, 'base_template.pptx'))
print("Extra layouts added. Total slides:", len(prs.slides))
