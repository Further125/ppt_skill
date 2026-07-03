"""
Plugin: Watermark

Adds a subtle watermark text to every slide.
Configure via deck_spec["watermark"] = "CONFIDENTIAL"
"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from pptx.dml.color import RGBColor


def post_slide(slide_spec, slide, context):
    """Add watermark if configured."""
    deck = context.get("deck_spec", {})
    watermark_text = deck.get("watermark") or slide_spec.get("watermark")
    if not watermark_text:
        return slide

    prs = context.get("prs")
    if prs is None:
        return slide
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Center watermark, rotated
    box_size = min(slide_w, slide_h) * 0.6
    left = (slide_w - box_size) // 2
    top = (slide_h - box_size) // 2

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, box_size, box_size
    )
    shape.fill.background()
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = watermark_text
    p.font.size = Pt(48)
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)  # light gray
    p.alignment = 2  # center

    # Rotate the shape
    shape.rotation = -45

    # Send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)

    return slide
