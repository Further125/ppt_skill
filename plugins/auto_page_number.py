"""
Plugin: Auto Page Number

Automatically adds page numbers to all slides except cover and closing.
"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt


def post_slide(slide_spec, slide, context):
    """Add page number to bottom-right corner."""
    layout = slide_spec.get("layout", "")
    if layout in ("cover", "closing"):
        return slide

    slide_num = context.get("slide_number", 0)
    total = context.get("total_slides", 0)

    # Add a small text box at bottom-right
    prs = context.get("prs")
    if prs is None:
        return slide
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    box_w = 300000  # ~0.33 inch
    box_h = 200000  # ~0.22 inch
    left = slide_w - box_w - 200000  # right margin
    top = slide_h - box_h - 150000   # bottom margin

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, box_w, box_h
    )
    shape.fill.background()
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_num} / {total}"
    p.font.size = Pt(10)
    p.font.name = "Microsoft YaHei"
    p.alignment = 3  # right align

    return slide
