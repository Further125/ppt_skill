import os
import random
import sys
import tempfile

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

try:
    from PIL import Image, ImageDraw, ImageFont
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import numpy as np
    HAS_MPL = True
except ImportError:
    HAS_MPL = False


def _temp_path(suffix):
    fd, path = tempfile.mkstemp(suffix=suffix, prefix='pptskill_')
    os.close(fd)
    return path


# ── Waterfall ──

def add_waterfall_chart(slide, data, placeholder_shape):
    if not HAS_MPL:
        return False
    categories = data.get("categories", [])
    values = data.get("values", [])
    if not categories or not values:
        return False

    path = _temp_path('.png')
    fig, ax = plt.subplots(figsize=(8, 4.5))
    running = 0
    for i, (cat, val) in enumerate(zip(categories, values)):
        if i == 0:
            running = val
            ax.bar(i, val, color="#3B82F6")
        elif i == len(categories) - 1:
            ax.bar(i, val, color="#10B981")
        else:
            bottom = running
            running += val
            ax.bar(i, val, bottom=bottom, color="#3B82F6" if val >= 0 else "#EF4444")
    ax.set_xticks(range(len(categories)))
    ax.set_xticklabels(categories, fontsize=10)
    ax.set_ylabel("Value")
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    plt.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)

    slide.shapes.add_picture(path, placeholder_shape.left, placeholder_shape.top,
                             placeholder_shape.width, placeholder_shape.height)
    return True


# ── Funnel ──

def add_funnel_chart(slide, data, placeholder_shape):
    if not HAS_MPL:
        return False
    categories = data.get("categories", [])
    values = data.get("values", [])
    if not categories or not values:
        return False

    path = _temp_path('.png')
    fig, ax = plt.subplots(figsize=(8, 4.5))
    colors = ["#3B82F6", "#10B981", "#F59E0B", "#EF4444", "#8B5CF6", "#EC4899", "#06B6D4"]
    for i, (cat, val) in enumerate(zip(categories, values)):
        width = val / max(values) * 0.9
        left = (1 - width) / 2
        ax.barh(i, width, left=left, height=0.7, color=colors[i % len(colors)], edgecolor='white')
        ax.text(0.5, i, f"{cat}: {val:,}", ha='center', va='center', color='white', fontsize=10, fontweight='bold')
    ax.set_yticks([])
    ax.set_xticks([])
    ax.set_xlim(0, 1)
    ax.invert_yaxis()
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_visible(False)
    ax.spines['left'].set_visible(False)
    plt.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)

    slide.shapes.add_picture(path, placeholder_shape.left, placeholder_shape.top,
                             placeholder_shape.width, placeholder_shape.height)
    return True


# ── Gantt ──

def add_gantt_chart(slide, data, placeholder_shape):
    if not HAS_MPL:
        return False
    categories = data.get("categories", [])
    values = data.get("values", [])
    if not categories or not values:
        return False

    path = _temp_path('.png')
    fig, ax = plt.subplots(figsize=(8, 4.5))
    colors = ["#3B82F6", "#10B981", "#F59E0B", "#EF4444", "#8B5CF6", "#EC4899"]
    start = 0
    for i, (cat, val) in enumerate(zip(categories, values)):
        ax.barh(i, val, left=start, height=0.5, color=colors[i % len(colors)], edgecolor='white')
        ax.text(start + val / 2, i, cat, ha='center', va='center', color='white', fontsize=9, fontweight='bold')
        start += val
    ax.set_yticks([])
    ax.set_xlabel("Days")
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    plt.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)

    slide.shapes.add_picture(path, placeholder_shape.left, placeholder_shape.top,
                             placeholder_shape.width, placeholder_shape.height)
    return True


# ── Radar ──

def add_radar_chart(slide, data, placeholder_shape):
    """Render radar chart as image for consistent cross-app display."""
    if not HAS_MPL:
        return False
    categories = data.get("categories", [])
    values = data.get("values", [])
    series_name = data.get("series_name", "Data")
    if not categories or not values:
        return False

    path = _temp_path('.png')

    # Configure matplotlib for Chinese text
    plt.rcParams['font.sans-serif'] = ['Noto Sans CJK JP', 'DejaVu Sans']
    plt.rcParams['axes.unicode_minus'] = False

    # Close the polygon
    angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
    values_plot = values + [values[0]]
    angles += angles[:1]

    fig, ax = plt.subplots(figsize=(6, 5.5), subplot_kw=dict(polar=True))
    # No fill, only line and markers
    ax.plot(angles, values_plot, color='#3B82F6', linewidth=2.5, marker='o', markersize=7)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories, fontsize=12)
    ax.set_ylim(0, max(values) * 1.1)
    ax.set_title(series_name, fontsize=14, pad=20)
    ax.grid(True, linestyle='--', alpha=0.5)
    plt.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)

    slide.shapes.add_picture(path, placeholder_shape.left, placeholder_shape.top,
                             placeholder_shape.width, placeholder_shape.height)
    return True


# ── Wordcloud ──

def generate_wordcloud(text_list, output_path, width=800, height=400):
    if not text_list:
        return None

    freq = {}
    for item in text_list:
        if isinstance(item, dict):
            word = item.get("word") or item.get("text", "")
            count = item.get("count") or item.get("weight", 1)
        else:
            word = str(item)
            count = 1
        freq[word] = freq.get(word, 0) + count

    if not freq:
        return None

    max_count = max(freq.values())
    img = Image.new("RGB", (width, height), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    try:
        font_base = ImageFont.truetype("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", 20)
    except Exception:
        font_base = ImageFont.load_default()

    positions = []
    colors = [(0x3B, 0x82, 0xF6), (0x10, 0xB9, 0x81), (0xF5, 0x9E, 0x0B),
              (0xEF, 0x44, 0x44), (0x8B, 0x5C, 0xF6), (0xEC, 0x48, 0x99),
              (0x06, 0xB6, 0xD4), (0x84, 0xCC, 0x16)]

    for word, count in sorted(freq.items(), key=lambda x: -x[1]):
        size = int(16 + (count / max_count) * 48)
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", size)
        except Exception:
            font = font_base

        bbox = draw.textbbox((0, 0), word, font=font)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]

        placed = False
        for _ in range(100):
            px = random.randint(10, width - tw - 10)
            py = random.randint(10, height - th - 10)
            rect = (px, py, px + tw, py + th)
            overlap = any(px < r[2] and px + tw > r[0] and py < r[3] and py + th > r[1] for r in positions)
            if not overlap:
                color = random.choice(colors)
                draw.text((px, py), word, fill=color, font=font)
                positions.append(rect)
                placed = True
                break

    img.save(output_path)
    return output_path
