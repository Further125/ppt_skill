#!/usr/bin/env python3
"""
Animation utilities for PPTX slides.
Supports: fade, fly_in, wipe, zoom, bounce, appear
Usage:
  from animator import add_slide_animation
  add_slide_animation(slide, effect="fade", delay_idx=0)
"""
import os, sys
from lxml import etree
from pptx.oxml.ns import qn

# Preset IDs for entrance effects (p:anim/p:set with presetClass="entr")
PRESETS = {
    "appear": 0,
    "fade": 1,
    "fly_in": 2,
    "wipe": 3,
    "zoom": 8,
    "bounce": 11,
    "swivel": 19,
    "float": 21,
}

def _ensure_timing(slide):
    """Ensure slide has p:timing element."""
    cSld = slide._element
    timing = cSld.find(qn('p:timing'))
    if timing is None:
        timing = etree.SubElement(cSld, qn('p:timing'))
    return timing

def _ensure_tnLst(timing):
    """Ensure timing has p:tnLst."""
    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = etree.SubElement(timing, qn('p:tnLst'))
    return tnLst

def _ensure_seq(tnLst):
    """Ensure tnLst has a sequence for click-triggered animations."""
    par = tnLst.find(qn('p:par'))
    if par is None:
        par = etree.SubElement(tnLst, qn('p:par'))
    seq = par.find(qn('p:seq'))
    if seq is None:
        seq = etree.SubElement(par, qn('p:seq'))
        seq.set('concurrent', '1')
        seq.set('nextAc', 'seek')
        # Add trigger (click on slide)
        cTn = etree.SubElement(seq, qn('p:cTn'))
        cTn.set('id', '1')
        cTn.set('dur', 'indefinite')
        cTn.set('nodeType', 'mainSeq')
        prev = etree.SubElement(cTn, qn('p:prevCondLst'))
        cond = etree.SubElement(prev, qn('p:cond'))
        cond.set('evt', 'onPrev')
        cond.set('delay', '0')
        tg = etree.SubElement(cond, qn('p:tn'))
        tg.set('val', '1')
        next_el = etree.SubElement(cTn, qn('p:nextCondLst'))
        cond2 = etree.SubElement(next_el, qn('p:cond'))
        cond2.set('evt', 'onNext')
        cond2.set('delay', '0')
        tg2 = etree.SubElement(cond2, qn('p:tn'))
        tg2.set('val', '1')
    return seq

def _get_next_anim_id(slide):
    """Find next available animation node id."""
    cSld = slide._element
    ids = []
    for elem in cSld.iter():
        for attr in ['id']:
            val = elem.get(qn('p:' + attr)) or elem.get(qn('a:' + attr))
            if val and val.isdigit():
                ids.append(int(val))
    return max(ids, default=1) + 1

def add_animation_to_shape(slide, shape, effect="fade", delay_idx=0, duration=500):
    """Add an entrance animation to a single shape."""
    preset_id = PRESETS.get(effect, PRESETS["fade"])
    timing = _ensure_timing(slide)
    tnLst = _ensure_tnLst(timing)
    seq = _ensure_seq(tnLst)
    next_id = _get_next_anim_id(slide)

    # Build animation par
    par = etree.SubElement(seq, qn('p:par'))
    cTn = etree.SubElement(par, qn('p:cTn'))
    cTn.set('id', str(next_id))
    cTn.set('fill', 'hold')

    st = etree.SubElement(cTn, qn('p:stCondLst'))
    cond = etree.SubElement(st, qn('p:cond'))
    cond.set('delay', str(delay_idx * 200))  # stagger by 200ms

    child_par = etree.SubElement(cTn, qn('p:childTnLst'))

    # The set/anim element
    set_el = etree.SubElement(child_par, qn('p:set'))
    cBhvr = etree.SubElement(set_el, qn('p:cBhvr'))
    cBhvr.set('additive', 'base')

    cB = etree.SubElement(cBhvr, qn('p:cB'))
    cB.set('rId', str(next_id + 1))

    tgtEl = etree.SubElement(cBhvr, qn('p:tgtEl'))
    spTgt = etree.SubElement(tgtEl, qn('p:spTgt'))
    spTgt.set('spid', str(shape.shape_id))

    attrName = etree.SubElement(cBhvr, qn('p:attrNameLst'))
    attr_name = etree.SubElement(attrName, qn('p:attrName'))
    attr_name.text = "style.visibility"

    to_el = etree.SubElement(set_el, qn('p:to'))
    to_el.set('val', 'visible')

    # Also add a preset filter for visual effect (works in PowerPoint)
    anim = etree.SubElement(child_par, qn('p:anim'))
    anim.set('calcmode', 'auto')
    anim.set('valueType', 'num')

    cBhvr2 = etree.SubElement(anim, qn('p:cBhvr'))
    cBhvr2.set('additive', 'sum')
    tgtEl2 = etree.SubElement(cBhvr2, qn('p:tgtEl'))
    spTgt2 = etree.SubElement(tgtEl2, qn('p:spTgt'))
    spTgt2.set('spid', str(shape.shape_id))
    attrName2 = etree.SubElement(cBhvr2, qn('p:attrNameLst'))
    attr_name2 = etree.SubElement(attrName2, qn('p:attrName'))
    attr_name2.text = "ppt_x"

    # This simplified version sets visibility; for full visual presets,
    # PowerPoint renders the presetClass="entr" effects when present.
    # For robust preset animation, we inject a preset filter element.
    child_tn = etree.SubElement(child_par, qn('p:childTnLst'))
    animEffect = etree.SubElement(child_par, qn('p:animEffect'))
    animEffect.set('transition', 'in')
    animEffect.set('filter', _effect_filter_name(effect))
    tgtEl3 = etree.SubElement(animEffect, qn('p:tgtEl'))
    spTgt3 = etree.SubElement(tgtEl3, qn('p:spTgt'))
    spTgt3.set('spid', str(shape.shape_id))

def _effect_filter_name(effect):
    mapping = {
        "fade": "fade",
        "fly_in": "fly",
        "wipe": "wipe",
        "zoom": "zoom",
        "bounce": "bounce",
        "appear": "appear",
        "swivel": "swivel",
        "float": "float",
    }
    return mapping.get(effect, "fade")

def add_slide_animation(slide, effect="fade", stagger=True, skip_shapes=None):
    """Add entrance animation to all shapes on a slide."""
    if skip_shapes is None:
        skip_shapes = []
    for idx, shape in enumerate(slide.shapes):
        if shape.shape_id in skip_shapes:
            continue
        delay = idx if stagger else 0
        add_animation_to_shape(slide, shape, effect=effect, delay_idx=delay)

def animate_presentation(prs, effect="fade", stagger=True):
    """Add animations to all slides in a presentation."""
    for slide in prs.slides:
        add_slide_animation(slide, effect=effect, stagger=stagger)

