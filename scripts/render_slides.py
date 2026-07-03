#!/usr/bin/env python3
"""
Render PPTX slides to PNG preview images.
Supports two modes:
  1. LibreOffice + pdftoppm (best quality, requires system deps)
  2. Pure-Python soft preview (fallback, no external deps)

Usage:
  python render_slides.py input.pptx output_dir [--dpi 144]
"""
import sys
import os
import io
import argparse
import subprocess
import shutil


from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont


def has_libreoffice():
    return shutil.which("soffice") is not None


def has_pdftoppm():
    return shutil.which("pdftoppm") is not None


def _pdf_to_pngs(pdf_path, output_dir, dpi=144, prefix="slide"):
    """Convert PDF pages to PNGs using PyMuPDF (fitz)."""
    try:
        import fitz
    except ImportError:
        return False
    doc = fitz.open(pdf_path)
    scale = dpi / 72
    for i, page in enumerate(doc):
        mat = fitz.Matrix(scale, scale)
        pix = page.get_pixmap(matrix=mat)
        out_path = os.path.join(output_dir, f"{prefix}_{i + 1:03d}.png")
        pix.save(out_path)
    doc.close()
    return True


def render_with_libreoffice(pptx_path, output_dir, dpi=144):
    """Use LibreOffice to convert PPTX -> PDF -> PNG."""
    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]

    # LibreOffice needs a writable HOME + user profile. In read-only sandboxes
    # only /tmp is writable (tmpfs), so put the profile there.
    # Using output_dir could fail when output_dir is on a read-only mount.
    lo_profile = os.path.join("/tmp", f"lo_profile_{os.getpid()}")
    os.makedirs(lo_profile, exist_ok=True)

    env = os.environ.copy()
    env["HOME"] = "/tmp"  # give LO a writable home (tmpfs)

    # PPTX -> PDF
    cmd = [
        "soffice", "--headless",
        "-env:UserInstallation=file://" + lo_profile,
        "--convert-to", "pdf",
        "--outdir", output_dir, pptx_path
    ]
    print(f"Running: {' '.join(cmd)}")
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, env=env, timeout=120)
    except subprocess.TimeoutExpired:
        print("LibreOffice conversion timed out (120s)")
        return False
    if result.returncode != 0:
        print(f"LibreOffice error: {result.stderr}")
        return False

    # Find generated PDF
    pdf_file = os.path.join(output_dir, os.path.basename(pptx_path).replace(".pptx", ".pdf"))
    if not os.path.exists(pdf_file):
        for f in os.listdir(output_dir):
            if f.endswith(".pdf"):
                pdf_file = os.path.join(output_dir, f)
                break

    if not os.path.exists(pdf_file):
        print("PDF not found after conversion")
        return False

    # PDF -> PNG (prefer pdftoppm, fallback to PyMuPDF)
    if has_pdftoppm():
        cmd = [
            "pdftoppm", "-png", "-r", str(dpi),
            "-cropbox", pdf_file,
            os.path.join(output_dir, "slide")
        ]
        print(f"Running: {' '.join(cmd)}")
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        except subprocess.TimeoutExpired:
            print("pdftoppm timed out")
            return False
        if result.returncode != 0:
            print(f"pdftoppm error: {result.stderr}")
            return False
        # Rename outputs
        for f in sorted(os.listdir(output_dir)):
            if f.startswith("slide-") and f.endswith(".png"):
                num = f.replace("slide-", "").replace(".png", "")
                new_name = f"slide_{int(num):03d}.png"
                os.rename(os.path.join(output_dir, f), os.path.join(output_dir, new_name))
    else:
        print("pdftoppm not found, using PyMuPDF fallback for PDF->PNG...")
        ok = _pdf_to_pngs(pdf_file, output_dir, dpi=dpi)
        if not ok:
            print("PyMuPDF fallback failed")
            return False

    print(f"Rendered to: {output_dir}")
    return True


def _get_rgb(color_obj, default=(0, 0, 0)):
    """Safely extract RGB tuple from a pptx color object."""
    try:
        if color_obj and color_obj.rgb:
            return (color_obj.rgb[0], color_obj.rgb[1], color_obj.rgb[2])
    except Exception:
        pass
    return default


def _shape_fill_color(shape, default=None):
    """Extract solid fill color from a shape, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            c = fill.fore_color
            if c and c.rgb:
                return (c.rgb[0], c.rgb[1], c.rgb[2])
    except Exception:
        pass
    return default


def _shape_line_color(shape, default=None):
    """Extract line color from a shape."""
    try:
        line = shape.line
        if line.fill.type is not None and line.fill.type == 1:  # SOLID
            c = line.fill.fore_color
            if c and c.rgb:
                return (c.rgb[0], c.rgb[1], c.rgb[2])
    except Exception:
        pass
    return default


def _emu_to_px(emu, dpi):
    return round(emu / 914400 * dpi)


def _pt_to_px(pt, dpi):
    """Convert points to pixels at given DPI (72 pt = 1 inch)."""
    return int(pt * dpi / 72)


def _load_fonts():
    """Load project-local or system CJK fonts."""
    project_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    candidates = [
        os.path.join(project_dir, "fonts", "NotoSansCJKsc-Regular.otf"),
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for fp in candidates:
        if os.path.exists(fp):
            return {'base': fp}
    return None


def _get_font(base_path, px_size):
    try:
        return ImageFont.truetype(base_path, px_size)
    except Exception:
        return ImageFont.load_default()


def _wrap_text(draw, text, font, max_width):
    """Simple greedy word-wrap for CJK text."""
    if not text:
        return []
    # Pillow textbbox with anchor doesn't support embedded newlines
    text = text.replace('\n', ' ')
    lines = []
    current_line = ""
    for ch in text:
        test = current_line + ch
        bbox = draw.textbbox((0, 0), test, font=font, anchor='lt')
        w = bbox[2] - bbox[0]
        if w > max_width and current_line:
            lines.append(current_line)
            current_line = ch
        else:
            current_line = test
    if current_line:
        lines.append(current_line)
    return lines if lines else [text]


def render_soft_preview(pptx_path, output_dir, dpi=144):
    """Pure-Python soft preview using Pillow. Renders shapes, fills, text, images."""
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)

    slide_w_px = _emu_to_px(prs.slide_width, dpi)
    slide_h_px = _emu_to_px(prs.slide_height, dpi)

    fonts = _load_fonts()
    base_font_path = fonts['base'] if fonts else None

    for idx, slide in enumerate(prs.slides):
        img = Image.new("RGB", (slide_w_px, slide_h_px), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        shapes = list(slide.shapes)

        # --- Pass 1: draw fills / backgrounds ---
        for shape in shapes:
            x = _emu_to_px(shape.left, dpi)
            y = _emu_to_px(shape.top, dpi)
            w = _emu_to_px(shape.width, dpi)
            h = _emu_to_px(shape.height, dpi)

            if w <= 0 or h <= 0:
                continue

            st = shape.shape_type

            if st == MSO_SHAPE_TYPE.AUTO_SHAPE or st == MSO_SHAPE_TYPE.TEXT_BOX:
                fill_color = _shape_fill_color(shape)
                if fill_color:
                    draw.rectangle([x, y, x + w, y + h], fill=fill_color)
                line_color = _shape_line_color(shape)
                if line_color:
                    draw.rectangle([x, y, x + w, y + h], outline=line_color, width=1)

            elif st == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pic_bytes = shape.image.blob
                    pic_img = Image.open(io.BytesIO(pic_bytes)).convert("RGB")
                    pic_img = pic_img.resize((w, h), Image.LANCZOS)
                    img.paste(pic_img, (x, y))
                except Exception:
                    draw.rectangle([x, y, x + w, y + h], outline=(180, 180, 180), width=1)
                    draw.rectangle([x + 1, y + 1, x + w - 1, y + h - 1], fill=(240, 240, 240))
                    mid_y = y + h // 2
                    draw.line([(x + 10, mid_y - 10), (x + w - 10, mid_y + 10)], fill=(180, 180, 180), width=2)
                    draw.line([(x + 10, mid_y + 10), (x + w - 10, mid_y - 10)], fill=(180, 180, 180), width=2)

            elif st == MSO_SHAPE_TYPE.CHART:
                draw.rectangle([x, y, x + w, y + h], outline=(100, 100, 100), width=1)
                draw.rectangle([x + 1, y + 1, x + w - 1, y + h - 1], fill=(250, 250, 250))
                chart_font = _get_font(base_font_path, _pt_to_px(14, dpi))
                label = "[Chart]"
                bbox = draw.textbbox((0, 0), label, font=chart_font, anchor='lt')
                tw = bbox[2] - bbox[0]
                th = bbox[3] - bbox[1]
                draw.text((x + (w - tw) // 2, y + (h - th) // 2), label, fill=(120, 120, 120), font=chart_font, anchor='lt')

            elif st == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                num_rows = len(table.rows)
                num_cols = len(table.columns)
                if num_rows == 0 or num_cols == 0:
                    continue
                # Compute column widths and row heights to avoid gaps from integer division
                col_widths = [w // num_cols] * num_cols
                for i in range(w % num_cols):
                    col_widths[i] += 1
                row_heights = [h // num_rows] * num_rows
                for i in range(h % num_rows):
                    row_heights[i] += 1
                for r_idx in range(num_rows):
                    for c_idx in range(num_cols):
                        cell = table.cell(r_idx, c_idx)
                        cx = x + sum(col_widths[:c_idx])
                        cy = y + sum(row_heights[:r_idx])
                        cw = col_widths[c_idx]
                        ch = row_heights[r_idx]
                        # Cell background
                        try:
                            if cell.fill.type is not None and cell.fill.type == 1:
                                ccol = cell.fill.fore_color
                                if ccol and ccol.rgb:
                                    draw.rectangle([cx, cy, cx + cw, cy + ch],
                                                   fill=(ccol.rgb[0], ccol.rgb[1], ccol.rgb[2]))
                        except Exception:
                            pass
                        # Cell border
                        draw.rectangle([cx, cy, cx + cw, cy + ch],
                                       outline=(180, 180, 180), width=1)
                        # Cell text (rendered in pass 2 via text_frame, but do a quick pass here)
                        cell_text = cell.text.strip()
                        if cell_text:
                            cfont = _get_font(base_font_path, _pt_to_px(12, dpi))
                            try:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        if run.font.size:
                                            cpt = run.font.size.pt if hasattr(run.font.size, 'pt') else int(run.font.size) // 12700
                                            cfont = _get_font(base_font_path, _pt_to_px(cpt, dpi))
                                            break
                            except Exception:
                                pass
                            # Auto-determine text color from cell background brightness
                            ccolor = (31, 41, 55)
                            try:
                                if cell.fill.type is not None and cell.fill.type == 1:
                                    ccol = cell.fill.fore_color
                                    if ccol and ccol.rgb:
                                        brightness = (ccol.rgb[0]*299 + ccol.rgb[1]*587 + ccol.rgb[2]*114) / 1000
                                        if brightness < 128:
                                            ccolor = (255, 255, 255)
                            except Exception:
                                pass
                            try:
                                bbox = draw.textbbox((0, 0), cell_text, font=cfont, anchor='lt')
                                tw = bbox[2] - bbox[0]
                                th = bbox[3] - bbox[1]
                                tx = cx + (cw - tw) // 2
                                ty = cy + (ch - th) // 2
                                draw.text((tx, ty), cell_text, fill=ccolor, font=cfont, anchor='lt')
                            except Exception:
                                pass

        # --- Pass 2: draw text on top ---
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            x = _emu_to_px(shape.left, dpi)
            y = _emu_to_px(shape.top, dpi)
            w = _emu_to_px(shape.width, dpi)
            h = _emu_to_px(shape.height, dpi)

            tf = shape.text_frame
            if not tf.text.strip():
                continue

            max_pt = 16
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pt = run.font.size.pt if hasattr(run.font.size, 'pt') else int(run.font.size) // 12700
                        max_pt = max(max_pt, pt)

            if max_pt >= 32:
                px = _pt_to_px(max_pt, dpi)
                font = _get_font(base_font_path, px)
                color = (255, 255, 255)
            elif max_pt >= 20:
                px = _pt_to_px(max_pt, dpi)
                font = _get_font(base_font_path, px)
                color = (31, 41, 55)
            else:
                px = _pt_to_px(max_pt, dpi)
                font = _get_font(base_font_path, max(px, 12))
                color = (31, 41, 55)

            fill_color = _shape_fill_color(shape)
            if fill_color:
                brightness = (fill_color[0] * 299 + fill_color[1] * 587 + fill_color[2] * 114) / 1000
                if brightness < 128:
                    color = (255, 255, 255)

            # Check slide-wide dark background (cover / closing / quote)
            if slide.name.lower() in ("cover", "closing", "quote"):
                for bg_shape in shapes:
                    if bg_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        bg_fill = _shape_fill_color(bg_shape)
                        if bg_fill:
                            brightness = (bg_fill[0] * 299 + bg_fill[1] * 587 + bg_fill[2] * 114) / 1000
                            if brightness < 128:
                                bx = _emu_to_px(bg_shape.left, dpi)
                                by = _emu_to_px(bg_shape.top, dpi)
                                bw = _emu_to_px(bg_shape.width, dpi)
                                bh = _emu_to_px(bg_shape.height, dpi)
                                if x >= bx and y >= by and x + w <= bx + bw and y + h <= by + bh:
                                    color = (255, 255, 255)
                                    break

            pad = 8
            cur_y = y + pad
            for para in tf.paragraphs:
                para_text = para.text
                if not para_text.strip():
                    cur_y += _pt_to_px(max_pt, dpi) // 2
                    continue

                para_pt = max_pt
                for run in para.runs:
                    if run.font.size:
                        pt = run.font.size.pt if hasattr(run.font.size, 'pt') else int(run.font.size) // 12700
                        para_pt = max(para_pt, pt)
                para_px = _pt_to_px(para_pt, dpi)
                para_font = _get_font(base_font_path, para_px)

                para_color = color
                if para.runs:
                    run_color = _get_rgb(para.runs[0].font.color, default=None)
                    if run_color:
                        para_color = run_color

                lines = _wrap_text(draw, para_text, para_font, max(1, w - pad * 2))
                align = para.alignment
                for line in lines:
                    bbox = draw.textbbox((0, 0), line, font=para_font, anchor='lt')
                    tw = bbox[2] - bbox[0]
                    th = bbox[3] - bbox[1]
                    if cur_y + th > y + h - pad:
                        break
                    if align == PP_ALIGN.CENTER:
                        tx = x + (w - tw) // 2
                    elif align == PP_ALIGN.RIGHT:
                        tx = x + w - tw - pad
                    else:
                        tx = x + pad
                    draw.text((tx, cur_y), line, fill=para_color, font=para_font, anchor='lt')
                    cur_y += th + int(para_px * 0.3)
                cur_y += int(para_px * 0.3)

        out_path = os.path.join(output_dir, f"slide_{idx + 1:03d}.png")
        img.save(out_path)
        print(f"  -> {out_path}")

    print(f"Soft preview rendered to: {output_dir}")
    return True


def main():
    parser = argparse.ArgumentParser(description="Render PPTX to PNG previews")
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument("output_dir", help="Output directory for PNGs")
    parser.add_argument("--dpi", type=int, default=144, help="DPI for rendering")
    parser.add_argument("--soft", action="store_true", help="Force soft preview mode")
    args = parser.parse_args()

    if not args.soft and has_libreoffice():
        print("Using LibreOffice for high-quality rendering...")
        success = render_with_libreoffice(args.pptx, args.output_dir, args.dpi)
        if success:
            return
        print("LibreOffice rendering failed, falling back to soft preview...")

    print("Using pure-Python soft preview...")
    render_soft_preview(args.pptx, args.output_dir, args.dpi)


if __name__ == "__main__":
    main()
