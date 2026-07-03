#!/usr/bin/env python3
"""
Theme engine for PPTX — apply color/font themes to generated presentations.
Usage:
  from theme_engine import apply_theme
  apply_theme(prs, theme_path_or_name)
"""
import sys, json, os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

THEMES_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "references", "themes")

BUILT_IN_THEMES = {
    "default": {
        "name": "Default",
        "primary": "#3B82F6",
        "secondary": "#0B1F3A",
        "accent": ["#3B82F6", "#10B981", "#F59E0B", "#EF4444", "#8B5CF6", "#EC4899"],
        "background": "#FFFFFF",
        "text": "#1F2937",
        "light_bg": "#F3F4F6",
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei"
    },
    "dark": {
        "name": "Dark Navy",
        "primary": "#60A5FA",
        "secondary": "#0F172A",
        "accent": ["#60A5FA", "#34D399", "#FBBF24", "#F87171", "#A78BFA", "#F472B6"],
        "background": "#0F172A",
        "text": "#F8FAFC",
        "light_bg": "#1E293B",
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei"
    },
    "warm": {
        "name": "Warm Sunset",
        "primary": "#EA580C",
        "secondary": "#431407",
        "accent": ["#EA580C", "#CA8A04", "#DC2626", "#16A34A", "#0891B2", "#7C3AED"],
        "background": "#FFFBEB",
        "text": "#292524",
        "light_bg": "#FEF3C7",
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei"
    },
    "forest": {
        "name": "Forest Green",
        "primary": "#16A34A",
        "secondary": "#052E16",
        "accent": ["#16A34A", "#059669", "#CA8A04", "#0891B2", "#7C3AED", "#BE123C"],
        "background": "#F0FDF4",
        "text": "#14532D",
        "light_bg": "#DCFCE7",
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei"
    },
    "minimal": {
        "name": "Minimal Gray",
        "primary": "#374151",
        "secondary": "#111827",
        "accent": ["#374151", "#6B7280", "#9CA3AF", "#4B5563", "#1F2937", "#111827"],
        "background": "#FAFAFA",
        "text": "#1F2937",
        "light_bg": "#F3F4F6",
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei"
    }
}

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def load_theme(theme_input):
    """Load theme from built-in name, JSON file path, or dict."""
    if isinstance(theme_input, dict):
        return theme_input
    if theme_input in BUILT_IN_THEMES:
        return BUILT_IN_THEMES[theme_input]
    if os.path.exists(theme_input):
        with open(theme_input, 'r', encoding='utf-8') as f:
            return json.load(f)
    # Try themes dir
    theme_path = os.path.join(THEMES_DIR, f"{theme_input}.json")
    if os.path.exists(theme_path):
        with open(theme_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    print(f"Warning: theme '{theme_input}' not found, using default")
    return BUILT_IN_THEMES["default"]

def _rgb_brightness(rgb):
    """Return relative brightness 0-1 for an RGBColor."""
    r, g, b = rgb[0], rgb[1], rgb[2]
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255.0

def _luminance(rgb):
    """Calculate relative luminance per WCAG 2.1."""
    def _channel(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * _channel(rgb[0]) + 0.7152 * _channel(rgb[1]) + 0.0722 * _channel(rgb[2])

def contrast_ratio(c1, c2):
    """WCAG contrast ratio between two RGBColors."""
    l1, l2 = _luminance(c1), _luminance(c2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)

# Legacy palette references used by build_pptx.py
OLD_DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
OLD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OLD_ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
OLD_TIMELINE_BAR = RGBColor(0x93, 0xC5, 0xFD)
OLD_NODE_BLUE = RGBColor(0x25, 0x63, 0xEB)
OLD_LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
OLD_TEMPLATE_GRAY_1 = RGBColor(0xAA, 0xBB, 0xCC)
OLD_TEMPLATE_GRAY_2 = RGBColor(0x88, 0x88, 0x88)

def _slide_has_dark_background(slide, prs):
    """Heuristic: does this slide already have a dark background?"""
    slide_area = prs.slide_width * prs.slide_height
    for shape in slide.shapes:
        shape_area = shape.width * shape.height
        if shape_area > slide_area * 0.8:
            # Large background shape
            try:
                if shape.fill.type is not None and shape.fill.fore_color and shape.fill.fore_color.rgb:
                    return _rgb_brightness(shape.fill.fore_color.rgb) < 0.5
            except Exception:
                pass
    return False

def _ensure_slide_background(slide, prs, color):
    """Add a full-slide background rectangle if none exists."""
    slide_area = prs.slide_width * prs.slide_height
    for shape in slide.shapes:
        if shape.width * shape.height > slide_area * 0.9:
            try:
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = color
                    return
            except Exception:
                pass
    # No background shape found — add one at the very back
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # Move to the back so it doesn't cover other shapes
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def apply_theme(prs, theme_input):
    """Apply theme colors and fonts to all slides in a presentation."""
    theme = load_theme(theme_input)
    primary = hex_to_rgb(theme["primary"])
    secondary = hex_to_rgb(theme["secondary"])
    bg = hex_to_rgb(theme["background"])
    text = hex_to_rgb(theme["text"])
    light_bg = hex_to_rgb(theme["light_bg"])
    font_title = theme.get("font_title", "Microsoft YaHei")
    font_body = theme.get("font_body", "Microsoft YaHei")

    bg_brightness = _rgb_brightness(bg)
    is_dark_theme = bg_brightness < 0.5

    # dark_text = readable on light backgrounds
    # light_text = readable on dark backgrounds
    dark_text_for_theme = text if not is_dark_theme else RGBColor(0x1F, 0x29, 0x37)
    # For very dark backgrounds, use slightly dimmed white to avoid excessive
    # contrast / halation (e.g. #E2E8F0 instead of #FFFFFF). This keeps contrast
    # in the 10-15:1 range rather than 21:1.
    if is_dark_theme and bg_brightness < 0.15:
        light_text_for_theme = RGBColor(0xE2, 0xE8, 0xF0)
    else:
        light_text_for_theme = OLD_WHITE if is_dark_theme else text

    for slide in prs.slides:
        slide_area = prs.slide_width * prs.slide_height

        # ── Ensure consistent slide background ───────────────
        # Dark themes: use background (deep dark) for all slides so white text has
        # strong contrast.  Secondary is reserved for accent bars / highlights.
        if is_dark_theme:
            _ensure_slide_background(slide, prs, bg)
        # For light themes, we leave default white background alone

        has_dark_bg = _slide_has_dark_background(slide, prs) or is_dark_theme

        for shape in slide.shapes:
            shape_area = shape.width * shape.height
            is_bg = shape_area > slide_area * 0.8

            # ── Background shapes ──────────────────────────────
            if is_bg and hasattr(shape, 'fill') and shape.fill.type is not None:
                try:
                    shape.fill.solid()
                    if is_dark_theme:
                        shape.fill.fore_color.rgb = bg
                    else:
                        shape.fill.fore_color.rgb = bg
                except Exception:
                    pass

            # ── Shape fills (accent bars, placeholders, etc.) ──
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                try:
                    if shape.fill.fore_color and shape.fill.fore_color.rgb:
                        old = shape.fill.fore_color.rgb
                        if old == OLD_LIGHT_GRAY:
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = light_bg
                        elif old == OLD_ACCENT_BLUE:
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = primary
                        elif old == OLD_TIMELINE_BAR:
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = primary
                        elif old == OLD_NODE_BLUE:
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = primary
                except Exception:
                    pass

            # ── Tables ─────────────────────────────────────────
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                try:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            cell_bg_brightness = 255
                            cell_bg = None
                            try:
                                if cell.fill.type is not None and cell.fill.fore_color and cell.fill.fore_color.rgb:
                                    c = cell.fill.fore_color.rgb
                                    cell_bg = c
                                    cell_bg_brightness = _rgb_brightness(c)
                            except Exception:
                                pass

                            # Keep table cells at their original light colors; do
                            # NOT force them to dark in dark themes.  Text is set
                            # to dark for readability on light cells.  color_guard
                            # will warn about high background contrast if needed.
                            cell_text_color = dark_text_for_theme if cell_bg_brightness >= 0.5 else light_text_for_theme
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    run.font.name = font_body
                                    try:
                                        current_color = run.font.color.rgb
                                    except AttributeError:
                                        current_color = None
                                    # Only override legacy colors; preserve custom colors
                                    if current_color in (OLD_DARK_TEXT, None, OLD_WHITE, OLD_TEMPLATE_GRAY_1, OLD_TEMPLATE_GRAY_2):
                                        run.font.color.rgb = cell_text_color
                except Exception:
                    pass

            # ── Charts ─────────────────────────────────────────
            if shape.shape_type == 3:  # MSO_SHAPE_TYPE.CHART
                try:
                    chart = shape.chart
                    from lxml import etree
                    C_NS = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
                    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

                    # Chart background: explicitly set to light so it stays light
                    # across all renderers (including LibreOffice dark mode).
                    # For dark themes this creates a light island on dark slide —
                    # color_guard will warn about this background-background contrast.
                    chart_bg = OLD_WHITE  # explicit light background
                    chart_bg_brightness = _rgb_brightness(chart_bg)
                    chart_text_color = dark_text_for_theme  # dark text on light chart

                    def _set_solid_fill(parent_elem, color):
                        """Set solid fill on a chart spPr element."""
                        spPr = parent_elem.find('{' + C_NS + '}spPr')
                        if spPr is None:
                            spPr = etree.SubElement(parent_elem, '{' + C_NS + '}spPr')
                        solidFill = spPr.find('{' + A_NS + '}solidFill')
                        if solidFill is None:
                            solidFill = etree.SubElement(spPr, '{' + A_NS + '}solidFill')
                        else:
                            for child in list(solidFill):
                                solidFill.remove(child)
                        srgbClr = etree.SubElement(solidFill, '{' + A_NS + '}srgbClr')
                        srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))

                    # Explicitly set chartSpace + plotArea to white/light
                    _set_solid_fill(chart._element, chart_bg)
                    plot_area = chart._element.find('.//' + '{' + C_NS + '}plotArea')
                    if plot_area is not None:
                        _set_solid_fill(plot_area, chart_bg)

                    def _set_txPr_color(parent_elem, color):
                        """Set default text color on a chart txPr element."""
                        txPr = parent_elem.find('{' + C_NS + '}txPr')
                        if txPr is None:
                            txPr = etree.SubElement(parent_elem, '{' + C_NS + '}txPr')
                            bodyPr = txPr.find('{' + A_NS + '}bodyPr')
                            if bodyPr is None:
                                etree.SubElement(txPr, '{' + A_NS + '}bodyPr')
                            lstStyle = txPr.find('{' + A_NS + '}lstStyle')
                            if lstStyle is None:
                                etree.SubElement(txPr, '{' + A_NS + '}lstStyle')
                        p = txPr.find('{' + A_NS + '}p')
                        if p is None:
                            p = etree.SubElement(txPr, '{' + A_NS + '}p')
                        pPr = p.find('{' + A_NS + '}pPr')
                        if pPr is None:
                            pPr = etree.SubElement(p, '{' + A_NS + '}pPr')
                        defRPr = pPr.find('{' + A_NS + '}defRPr')
                        if defRPr is None:
                            defRPr = etree.SubElement(pPr, '{' + A_NS + '}defRPr')
                        solidFill = defRPr.find('{' + A_NS + '}solidFill')
                        if solidFill is None:
                            solidFill = etree.SubElement(defRPr, '{' + A_NS + '}solidFill')
                        else:
                            for child in list(solidFill):
                                solidFill.remove(child)
                        srgbClr = etree.SubElement(solidFill, '{' + A_NS + '}srgbClr')
                        srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))

                    # 1. Chart-level default text color (affects series names,
                    #    data labels, and any text without explicit override)
                    _set_txPr_color(chart._element, chart_text_color)

                    # 2. Plot-area default text color
                    plot_area = chart._element.find('.//' + '{' + C_NS + '}plotArea')
                    if plot_area is not None:
                        _set_txPr_color(plot_area, chart_text_color)

                    # 3. Axis tick labels
                    try:
                        if chart.category_axis and chart.category_axis.tick_labels:
                            chart.category_axis.tick_labels.font.color.rgb = chart_text_color
                    except ValueError:
                        pass
                    try:
                        if chart.value_axis and chart.value_axis.tick_labels:
                            chart.value_axis.tick_labels.font.color.rgb = chart_text_color
                    except ValueError:
                        pass

                    # 4. Legend
                    if chart.has_legend and chart.legend:
                        try:
                            chart.legend.font.color.rgb = chart_text_color
                        except Exception:
                            pass
                        _set_txPr_color(chart.legend._element, chart_text_color)

                    # 5. Chart title
                    if chart.has_title and chart.chart_title:
                        for para in chart.chart_title.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.color.rgb = chart_text_color

                    # 6. Series names (tx elements inside each ser)
                    if plot_area is not None:
                        for chart_elem in plot_area:
                            if chart_elem.tag.endswith('}doughnutChart') or chart_elem.tag.endswith('}pieChart') or chart_elem.tag.endswith('}barChart') or chart_elem.tag.endswith('}lineChart') or chart_elem.tag.endswith('}columnChart') or chart_elem.tag.endswith('}areaChart'):
                                for ser in chart_elem.findall('{' + C_NS + '}ser'):
                                    tx = ser.find('{' + C_NS + '}tx')
                                    if tx is not None:
                                        _set_txPr_color(tx, chart_text_color)
                except Exception:
                    pass

            # ── Text ───────────────────────────────────────────
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_body

                    try:
                        current_color = run.font.color.rgb
                    except AttributeError:
                        current_color = None

                    # Determine target text color based on slide background
                    target_color = light_text_for_theme if has_dark_bg else dark_text_for_theme

                    is_title_size = run.font.size and run.font.size.pt >= 28
                    if is_title_size:
                        run.font.name = font_title

                    # Helper: detect fill of underlying auto_shape (for textboxes
                    # sitting on top of colored shapes like process boxes)
                    def _detect_underlying_fill(shp, sld):
                        try:
                            if shp.fill.type is not None and shp.fill.fore_color and shp.fill.fore_color.rgb:
                                return shp.fill.fore_color.rgb
                        except Exception:
                            pass
                        # Textbox with no fill: check overlapping auto_shapes
                        sx, sy, sw, shv = shp.left, shp.top, shp.width, shp.height
                        best_fill = None
                        best_area = 0
                        for other in sld.shapes:
                            if other.shape_type != 1:  # only AUTO_SHAPE
                                continue
                            try:
                                if other.fill.type is None or not other.fill.fore_color or not other.fill.fore_color.rgb:
                                    continue
                                ox, oy, ow, oh = other.left, other.top, other.width, other.height
                                if sx < ox + ow and sx + sw > ox and sy < oy + oh and sy + shv > oy:
                                    area = other.width * other.height
                                    if area > best_area:
                                        best_area = area
                                        best_fill = other.fill.fore_color.rgb
                            except Exception:
                                pass
                        return best_fill

                    shape_fill = _detect_underlying_fill(shape, slide)

                    # Map legacy colors to theme-aware colors
                    if current_color in (OLD_DARK_TEXT, None):
                        if shape_fill is not None and shape_fill != bg:
                            fill_l = _luminance(shape_fill)
                            white_l = _luminance(OLD_WHITE)
                            dark_l = _luminance(OLD_DARK_TEXT)
                            white_ratio = (max(fill_l, white_l) + 0.05) / (min(fill_l, white_l) + 0.05)
                            dark_ratio = (max(fill_l, dark_l) + 0.05) / (min(fill_l, dark_l) + 0.05)
                            chosen = OLD_WHITE if white_ratio > dark_ratio else OLD_DARK_TEXT
                            run.font.color.rgb = chosen
                        else:
                            run.font.color.rgb = target_color
                    elif current_color == OLD_WHITE:
                        if shape_fill is not None and shape_fill != bg:
                            run.font.color.rgb = OLD_WHITE
                        else:
                            run.font.color.rgb = light_text_for_theme
                    elif current_color == OLD_ACCENT_BLUE:
                        run.font.color.rgb = primary
                    elif current_color in (OLD_TEMPLATE_GRAY_1, OLD_TEMPLATE_GRAY_2):
                        # Template default grays: on dark bg use light text,
                        # on light bg use dark text (or a slightly muted dark)
                        run.font.color.rgb = light_text_for_theme if has_dark_bg else dark_text_for_theme
                    else:
                        # Preserve explicitly-set custom colors (user may have chosen them intentionally)
                        pass

    print(f"Theme applied: {theme.get('name', 'Custom')} (dark={is_dark_theme})")
    return prs

def audit_contrast(prs, theme_input):
    """Audit every text shape against its local background and report low-contrast issues.
    Returns a list of warning dicts."""
    theme = load_theme(theme_input)
    bg = hex_to_rgb(theme["background"])
    text = hex_to_rgb(theme["text"])
    warnings = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_area = prs.slide_width * prs.slide_height
        # Determine slide-level background color
        slide_bg = bg
        for shape in slide.shapes:
            if shape.width * shape.height > slide_area * 0.8:
                try:
                    if shape.fill.type is not None and shape.fill.fore_color and shape.fill.fore_color.rgb:
                        slide_bg = shape.fill.fore_color.rgb
                except Exception:
                    pass

        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            # Sample background at this shape's location
            local_bg = _sample_background_at_shape(slide, shape, slide_bg, slide_area)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        c = run.font.color.rgb
                    except AttributeError:
                        c = text
                    ratio = contrast_ratio(c, local_bg)
                    if ratio < 3.0:
                        warnings.append({
                            "slide": slide_idx,
                            "shape_text": shape.text_frame.text[:60].replace("\n", " "),
                            "text_color": f"#{c[0]:02X}{c[1]:02X}{c[2]:02X}",
                            "bg_color": f"#{local_bg[0]:02X}{local_bg[1]:02X}{local_bg[2]:02X}",
                            "contrast": round(ratio, 2),
                            "severity": "CRITICAL" if ratio < 2.0 else "WARNING"
                        })
    return warnings

def _sample_background_at_shape(slide, text_shape, slide_bg, slide_area):
    """Return the background color at the text shape's bounding box.
    If an image or another shape covers the text area, return that shape's fill color.
    Falls back to slide_bg."""
    tx, ty, tw, th = text_shape.left, text_shape.top, text_shape.width, text_shape.height
    cx, cy = tx + tw // 2, ty + th // 2

    for shape in slide.shapes:
        if shape == text_shape:
            continue
        sx, sy, sw, sh = shape.left, shape.top, shape.width, shape.height
        # Check if text shape center lies inside this shape
        if sx <= cx <= sx + sw and sy <= cy <= sy + sh:
            # If shape has a solid fill, that's our local background
            try:
                if shape.fill.type is not None and shape.fill.fore_color and shape.fill.fore_color.rgb:
                    return shape.fill.fore_color.rgb
            except Exception:
                pass
            # If it's a picture, we can't easily sample; fall through
    return slide_bg

def list_themes():
    """List available built-in themes."""
    return list(BUILT_IN_THEMES.keys())

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", help="Input PPTX")
    parser.add_argument("--theme", default="default", help="Theme name or JSON path")
    parser.add_argument("--output", "-o", required=True, help="Output PPTX")
    parser.add_argument("--list", action="store_true", help="List built-in themes")
    parser.add_argument("--audit", action="store_true", help="Run contrast audit after applying theme")
    args = parser.parse_args()

    if args.list:
        for t in list_themes():
            print(f"  {t}: {BUILT_IN_THEMES[t]['name']}")
        sys.exit(0)

    prs = Presentation(args.pptx)
    apply_theme(prs, args.theme)

    if args.audit:
        warnings = audit_contrast(prs, args.theme)
        if warnings:
            print(f"\n⚠️  Contrast audit found {len(warnings)} issue(s):")
            for w in warnings:
                print(f"  Slide {w['slide']} [{w['severity']}] ratio={w['contrast']} "
                      f"text={w['text_color']} on bg={w['bg_color']} "
                      f"\"{w['shape_text']}\"")
        else:
            print("\n✅ Contrast audit passed — no low-contrast issues found.")

    prs.save(args.output)
    print(f"Saved: {args.output}")
