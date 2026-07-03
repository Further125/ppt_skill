#!/usr/bin/env python3
"""
Adaptive layout strategies for ppt-skill.

Handles overflow, density, and responsive adjustments.
"""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    from pptx.util import Pt, Emu
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def _shape_actual_text_height(shape, safety_factor=1.15):
    """Estimate rendered text height using PowerPoint font metrics."""
    if not shape.has_text_frame:
        return 0
    tf = shape.text_frame
    if not tf.text.strip():
        return 0
    try:
        from text_fitter import _pptx_wrap_text, _pptx_text_height, _emu_to_px
    except ImportError:
        return shape.height

    width_px = _emu_to_px(shape.width)
    paragraphs = []
    for para in tf.paragraphs:
        text = para.text
        if not text.strip():
            continue
        font_size = 16
        for run in para.runs:
            if run.font.size:
                pt = run.font.size.pt if hasattr(run.font.size, 'pt') else int(run.font.size) // 12700
                font_size = pt
                break
        lines = _pptx_wrap_text(text, font_size, width_px * 0.92)
        paragraphs.append((lines, font_size))

    total_h = sum(_pptx_text_height(lines, fs) for lines, fs in paragraphs)
    total_h *= safety_factor
    return int(total_h * 12700)  # px to EMU


def apply_adaptive_strategy(shape, strategy="auto", slide=None, min_font=8, max_font=72):
    """Apply adaptive strategy to a shape.

    Strategies:
    - auto: try shrink -> expand -> truncate in sequence
    - shrink: reduce font size to fit
    - expand: increase shape height, push down overlapping shapes
    - truncate: truncate text with ... if it overflows
    - multi_column: split bullet list into 2 columns
    """
    if not shape.has_text_frame:
        return False

    tf = shape.text_frame
    text = tf.text.strip()
    if not text:
        return False

    actual_h = _shape_actual_text_height(shape)
    if actual_h <= shape.height:
        return True  # already fits

    # Get current font size
    current_size = 16
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size:
                pt = run.font.size.pt if hasattr(run.font.size, 'pt') else int(run.font.size) // 12700
                current_size = pt
                break
        break

    if strategy in ("auto", "shrink"):
        # Try shrinking first
        try:
            from text_fitter import fit_shape_text
            fit_shape_text(shape, max_size=current_size, min_size=min_font, word_wrap=True)
            actual_h = _shape_actual_text_height(shape)
            if actual_h <= shape.height:
                return True
        except Exception:
            pass

    if strategy in ("auto", "multi_column"):
        # Try multi-column for bullet lists
        if _try_multi_column(shape):
            return True

    if strategy in ("auto", "expand"):
        # Try expanding height
        if slide and _try_expand(shape, slide):
            return True

    if strategy in ("auto", "truncate"):
        # Last resort: truncate
        _truncate_shape_text(shape)
        return True

    return False


def _try_multi_column(shape):
    """Split long bullet lists into 2 columns within the same shape."""
    if not shape.has_text_frame:
        return False

    tf = shape.text_frame
    # Check if we have a long bullet list
    bullet_texts = []
    for para in tf.paragraphs:
        text = para.text.strip()
        if text:
            bullet_texts.append(text)

    if len(bullet_texts) < 6:
        return False

    # Use 2 columns
    tf.word_wrap = True
    try:
        # python-pptx doesn't expose columns directly on text_frame,
        # but we can simulate by adjusting paragraph spacing and using tab stops
        # For now, we'll reformat the text into two columns using tabs
        mid = (len(bullet_texts) + 1) // 2
        left_col = bullet_texts[:mid]
        right_col = bullet_texts[mid:]

        # Build two-column text
        max_len = max(len(left_col), len(right_col))
        lines = []
        for i in range(max_len):
            left = left_col[i] if i < len(left_col) else ""
            right = right_col[i] if i < len(right_col) else ""
            if right:
                lines.append(f"{left}\t{right}")
            else:
                lines.append(left)

        new_text = "\n".join(lines)
        # Set text on first paragraph, clear others
        if tf.paragraphs:
            tf.paragraphs[0].text = new_text
            for para in tf.paragraphs[1:]:
                para.text = ""
        return True
    except Exception:
        return False


def _try_expand(shape, slide):
    """Expand shape height and push down overlapping shapes."""
    actual_h = _shape_actual_text_height(shape)
    if actual_h <= 0:
        return False

    padding = 50000
    required_h = actual_h + padding
    current_h = shape.height

    if required_h <= current_h:
        return True

    # Get slide dimensions
    slide_height = 6858000  # default 16:9
    try:
        spTree = slide.shapes._spTree
        sld = spTree.getparent()
        if sld is not None:
            cy = sld.attrib.get('cy')
            if cy:
                slide_height = int(cy)
    except Exception:
        pass

    max_allowed_bottom = slide_height - 100000
    old_bottom = shape.top + current_h
    new_bottom = min(shape.top + required_h, max_allowed_bottom)
    new_height = new_bottom - shape.top

    if new_height <= current_h:
        return False  # can't expand

    shape.height = new_height
    shift = new_bottom - old_bottom

    # Push down overlapping shapes
    for other in slide.shapes:
        if other is shape:
            continue
        other_top = other.top
        if other_top >= old_bottom - 10000:
            shape_left = shape.left
            shape_right = shape.left + shape.width
            other_left = other.left
            other_right = other.left + other.width
            h_overlap = not (shape_right < other_left or shape_left > other_right)
            if h_overlap:
                new_other_top = other.top + shift
                if new_other_top + other.height <= max_allowed_bottom + 200000:
                    other.top = new_other_top

    return True


def _truncate_shape_text(shape, suffix=" ..."):
    """Truncate text to fit shape, adding suffix."""
    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    text = tf.text.strip()
    if not text:
        return

    # Simple binary search for truncation point
    low, high = 0, len(text)
    best_len = 0

    while low <= high:
        mid = (low + high) // 2
        truncated = text[:mid] + suffix

        # Temporarily set text and measure
        original_texts = []
        for para in tf.paragraphs:
            original_texts.append(para.text)

        if tf.paragraphs:
            tf.paragraphs[0].text = truncated
            for para in tf.paragraphs[1:]:
                para.text = ""

        actual_h = _shape_actual_text_height(shape)

        if actual_h <= shape.height:
            best_len = mid
            low = mid + 1
        else:
            high = mid - 1

        # Restore original
        for i, para in enumerate(tf.paragraphs):
            if i < len(original_texts):
                para.text = original_texts[i]

    if best_len > 0:
        final_text = text[:best_len] + suffix
        if tf.paragraphs:
            tf.paragraphs[0].text = final_text
            for para in tf.paragraphs[1:]:
                para.text = ""


def adjust_slide_for_overflow(slide, strategy="auto"):
    """Enhanced overflow handling for a slide."""
    text_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            if shape._element.getparent().tag.endswith('}tc'):
                continue
        except Exception:
            pass
        text_shapes.append(shape)

    if not text_shapes:
        return []

    text_shapes.sort(key=lambda s: s.top)
    results = []

    for shape in text_shapes:
        success = apply_adaptive_strategy(shape, strategy=strategy, slide=slide)
        results.append({
            "shape": shape.name,
            "success": success,
            "height": shape.height,
        })

    return results


def auto_paginate_deck(deck_spec, max_items_per_slide=10):
    """Automatically split slides with too many content items into multiple slides.

    Returns a new deck_spec with paginated slides.
    """
    import copy
    new_slides = []
    for slide in copy.deepcopy(deck_spec).get("slides", []):
        content = slide.get("content", [])
        if isinstance(content, str):
            content = [content]

        if len(content) > max_items_per_slide and slide.get("layout") in ("title_content", "two_column"):
            # Split content into chunks
            chunks = [content[i:i + max_items_per_slide] for i in range(0, len(content), max_items_per_slide)]
            total = len(chunks)
            for idx, chunk in enumerate(chunks):
                new_slide = copy.deepcopy(slide)
                new_slide["content"] = chunk
                if total > 1:
                    title = slide.get("title", "")
                    new_slide["title"] = f"{title} ({idx + 1}/{total})"
                new_slides.append(new_slide)
        else:
            new_slides.append(slide)

    result = copy.deepcopy(deck_spec)
    result["slides"] = new_slides
    return result


def suggest_slide_split(slides_spec):
    """Analyze deck spec and suggest slides that might need splitting."""
    suggestions = []
    for i, slide in enumerate(slides_spec.get("slides", [])):
        content = slide.get("content", [])
        if isinstance(content, str):
            content = [content]
        item_count = len(content)

        if item_count > 12:
            suggestions.append({
                "slide": i + 1,
                "layout": slide.get("layout"),
                "issue": f"Too many items ({item_count}), consider splitting into 2 slides",
                "suggested_action": "split"
            })
        elif item_count > 8:
            suggestions.append({
                "slide": i + 1,
                "layout": slide.get("layout"),
                "issue": f"High density ({item_count} items), may need multi-column or font shrink",
                "suggested_action": "multi_column"
            })

        # Check chart data density
        chart_data = slide.get("chart_data", {})
        categories = chart_data.get("categories", [])
        if len(categories) > 12:
            suggestions.append({
                "slide": i + 1,
                "layout": "chart",
                "issue": f"Chart has {len(categories)} categories, may be hard to read",
                "suggested_action": "group_categories"
            })

    return suggestions
