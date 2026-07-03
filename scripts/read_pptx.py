#!/usr/bin/env python3
"""
Read an existing PPTX and output a JSON deck representation.

Usage:
    python scripts/read_pptx.py input.pptx --output deck.json
    python scripts/read_pptx.py input.pptx --output deck.json --with-layouts
"""

import argparse
import json
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def _rgb_to_hex(rgb_color):
    """Convert RGBColor to hex string."""
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color[0]:02X}{rgb_color[1]:02X}{rgb_color[2]:02X}"
    except Exception:
        return None


def _extract_text(shape):
    """Extract text from a shape, preserving paragraph structure."""
    if not shape.has_text_frame:
        return None
    paragraphs = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    if not paragraphs:
        return None
    if len(paragraphs) == 1:
        return paragraphs[0]
    return paragraphs


def _extract_shape_data(shape):
    """Extract data about a single shape."""
    data = {
        "name": shape.name,
        "shape_type": str(shape.shape_type).split(".")[-1].lower() if shape.shape_type else "unknown",
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }

    # Text
    text = _extract_text(shape)
    if text:
        data["text"] = text

    # Fill
    try:
        if shape.fill.type is not None:
            if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                hex_color = _rgb_to_hex(shape.fill.fore_color.rgb)
                if hex_color:
                    data["fill"] = hex_color
    except Exception:
        pass

    # Line/border
    try:
        if shape.line and shape.line.fill.type is not None:
            hex_color = _rgb_to_hex(shape.line.fill.fore_color.rgb)
            if hex_color:
                data["border_color"] = hex_color
            if shape.line.width:
                data["border_width"] = shape.line.width
    except Exception:
        pass

    # Image
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            data["image_path"] = "[embedded_image]"
        except Exception:
            pass

    # Chart
    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        try:
            chart = shape.chart
            chart_data = {"type": str(chart.chart_type).split(".")[-1].lower()}
            # Try to extract categories and series
            if hasattr(chart, 'plots') and chart.plots:
                plot = chart.plots[0]
                if hasattr(plot, 'categories') and plot.categories:
                    chart_data["categories"] = [cat.label if hasattr(cat, 'label') else str(cat) for cat in plot.categories]
                if hasattr(plot, 'series') and plot.series:
                    series_list = []
                    for s in plot.series:
                        series_list.append({
                            "name": s.name if s.name else "Series",
                            "values": list(s.values) if hasattr(s, 'values') else []
                        })
                    chart_data["series"] = series_list
            data["chart_data"] = chart_data
        except Exception:
            pass

    # Table
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            table = shape.table
            rows = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                rows.append(row_data)
            if rows:
                data["table_data"] = {
                    "headers": rows[0] if rows else [],
                    "rows": rows[1:] if len(rows) > 1 else []
                }
        except Exception:
            pass

    return data


def _infer_layout(slide_data):
    """Infer layout type from slide content."""
    shapes = slide_data.get("shapes", [])
    has_chart = any("chart_data" in s for s in shapes)
    has_table = any("table_data" in s for s in shapes)
    has_image = any(s.get("shape_type") == "picture" for s in shapes)

    text_shapes = [s for s in shapes if s.get("text")]
    texts = [str(s.get("text", "")).lower() for s in text_shapes]

    # Check for title-like text
    title_shapes = [s for s in text_shapes if s.get("top", 0) < 1500000]

    if has_chart:
        return "chart"
    if has_table:
        return "table"
    if has_image and len(text_shapes) > 1:
        return "image_content"

    # Check for cover/closing patterns
    for t in texts:
        if any(word in t for word in ["thank you", "谢谢", "questions", "q & a"]):
            return "closing"

    if len(title_shapes) == 1 and len(text_shapes) == 1:
        # Single title shape - might be cover or closing
        if len(shapes) <= 3:
            return "cover"

    if len(text_shapes) >= 2:
        return "title_content"

    return "title_content"


def read_pptx(pptx_path, include_layouts=False):
    """Read a PPTX file and return a deck dict."""
    if not HAS_PPTX:
        raise ImportError("python-pptx is required")

    prs = Presentation(pptx_path)

    deck = {
        "title": os.path.splitext(os.path.basename(pptx_path))[0],
        "source": pptx_path,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": []
    }

    if include_layouts:
        deck["layouts"] = []
        for layout in prs.slide_layouts:
            layout_info = {
                "name": layout.name,
                "shapes": []
            }
            for shape in layout.placeholders:
                layout_info["shapes"].append({
                    "name": shape.name,
                    "type": str(shape.placeholder_format.type).split(".")[-1] if shape.placeholder_format else "unknown",
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                })
            deck["layouts"].append(layout_info)

    for idx, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": idx,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "unknown",
            "shapes": []
        }

        for shape in slide.shapes:
            shape_data = _extract_shape_data(shape)
            if shape_data:
                slide_data["shapes"].append(shape_data)

        # Try to infer layout
        inferred = _infer_layout(slide_data)
        if inferred:
            slide_data["inferred_layout"] = inferred

        deck["slides"].append(slide_data)

    return deck


def main():
    parser = argparse.ArgumentParser(description="Read a PPTX and output JSON representation")
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument("--output", "-o", default="-", help="Output JSON file (default: stdout)")
    parser.add_argument("--with-layouts", action="store_true", help="Include slide layouts")
    parser.add_argument("--indent", type=int, default=2, help="JSON indent")
    args = parser.parse_args()

    deck = read_pptx(args.pptx, include_layouts=args.with_layouts)
    json_str = json.dumps(deck, ensure_ascii=False, indent=args.indent)

    if args.output == "-":
        print(json_str)
    else:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(json_str)
        print(f"Read {len(deck['slides'])} slides from {args.pptx}")
        print(f"Output written to {args.output}")


if __name__ == "__main__":
    main()
