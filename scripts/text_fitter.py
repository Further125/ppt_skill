#!/usr/bin/env python3
"""Text auto-fit: compute best font size to fit text inside a shape/cell.

Uses empirical character-width models calibrated for PowerPoint's
Microsoft YaHei rendering, plus Pillow for fallback measurement.
"""

import os
from PIL import ImageFont, ImageDraw, Image
from pptx.util import Pt

# Default font search paths (CJK-first)
_DEFAULT_FONT_CANDIDATES = [
    os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                 "fonts", "NotoSansCJKsc-Regular.otf"),
    "/usr/share/fonts/opentype/noto/NotoSansCJKsc-Regular.otf",
    "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
]

_font_cache = {}


def _resolve_font(font_path=None):
    if font_path and os.path.exists(font_path):
        return font_path
    for fp in _DEFAULT_FONT_CANDIDATES:
        if os.path.exists(fp):
            return fp
    return None


def _get_font(font_path, point_size):
    key = (font_path, point_size)
    if key not in _font_cache:
        _font_cache[key] = ImageFont.truetype(font_path, point_size)
    return _font_cache[key]


def _measure_text(text, font_path, point_size):
    """Return (width_px, height_px) for *text* rendered at *point_size*."""
    font = _get_font(font_path, point_size)
    draw = ImageDraw.Draw(Image.new("RGB", (1, 1)))
    try:
        bbox = draw.textbbox((0, 0), text, font=font, anchor='lt')
        return bbox[2] - bbox[0], bbox[3] - bbox[1]
    except Exception:
        return font.getsize(text)


def _px_to_emu(px):
    return int(px * 914400 / 96)


def _emu_to_px(emu):
    return emu * 96 / 914400


# ── PowerPoint empirical width model (Microsoft YaHei) ──

def _pptx_char_width(char, font_size):
    """Estimated width of a single char in PowerPoint (pt units)."""
    cp = ord(char)
    # CJK unified ideographs and punctuation blocks
    is_cjk = (0x4E00 <= cp <= 0x9FFF or
              0x3400 <= cp <= 0x4DBF or
              0x2E80 <= cp <= 0x2EFF or
              0x3000 <= cp <= 0x303F or
              0x31C0 <= cp <= 0x31EF or
              0xF900 <= cp <= 0xFAFF or
              0xFE30 <= cp <= 0xFE4F or
              0xFF00 <= cp <= 0xFFEF)
    if is_cjk:
        return font_size * 0.98
    if 'a' <= char <= 'z':
        return font_size * 0.50
    if 'A' <= char <= 'Z':
        return font_size * 0.62
    if '0' <= char <= '9':
        return font_size * 0.55
    if char in ' ,.;:!?()[]{}"\'<>/\\|-_=+@#$%&*':
        return font_size * 0.32
    if char == ' ':
        return font_size * 0.28
    return font_size * 0.50


def _pptx_wrap_text(text, font_size, max_width_pt):
    """Word-wrap *text* at *font_size* into lines fitting *max_width_pt*.

    Uses the empirical PowerPoint width model.
    """
    if not text or not text.strip():
        return []
    lines = []
    current = ""
    current_w = 0.0
    for ch in text:
        if ch == '\n':
            if current:
                lines.append(current)
            current = ""
            current_w = 0.0
            continue
        ch_w = _pptx_char_width(ch, font_size)
        # If this single char is wider than the whole line, force it on
        if ch_w > max_width_pt and not current:
            lines.append(ch)
            current = ""
            current_w = 0.0
            continue
        if current_w + ch_w <= max_width_pt:
            current += ch
            current_w += ch_w
        else:
            lines.append(current)
            current = ch
            current_w = ch_w
    if current:
        lines.append(current)
    return lines


def _pptx_text_height(lines, font_size):
    """Estimated total height (pt) of wrapped lines in PowerPoint."""
    if not lines:
        return 0
    line_h = font_size * 1.40  # Microsoft YaHei line height approx
    # Space between paragraphs / bullet lines
    para_gap = font_size * 0.35
    return line_h * len(lines) + para_gap * max(0, len(lines) - 1)


# ── Public API ──

def best_fit_font_size(text, max_width_emu, max_height_emu,
                       font_path=None, max_size=72, min_size=8,
                       word_wrap=True):
    """Return the largest point size <= max_size that lets *text* fit.

    Uses the empirical PowerPoint width model rather than Pillow metrics,
    because Pillow's Noto Sans CJK differs from PowerPoint's Microsoft YaHei.
    """
    if not text or not text.strip():
        return max_size

    max_w_pt = _emu_to_px(max_width_emu) * 72 / 96  # px→pt, but _emu_to_px already gives px
    # Actually _emu_to_px returns px at 96 DPI; 1 pt = 1 px at 72 DPI
    # So px * (72/96) = pt
    max_w_pt = _emu_to_px(max_width_emu) * 0.75
    max_h_pt = _emu_to_px(max_height_emu) * 0.75

    def _fits(size):
        if word_wrap:
            lines = _pptx_wrap_text(text, size, max_w_pt)
            if not lines:
                return False
            h = _pptx_text_height(lines, size)
            return h <= max_h_pt
        else:
            total_w = sum(_pptx_char_width(ch, size) for ch in text)
            return total_w <= max_w_pt

    lo, hi = min_size, int(max_size)
    best = min_size
    while lo <= hi:
        mid = (lo + hi) // 2
        if _fits(mid):
            best = mid
            lo = mid + 1
        else:
            hi = mid - 1
    return best


def fit_shape_text(shape, max_size=72, min_size=8, font_path=None,
                   honor_margins=True, word_wrap=True):
    """Auto-fit all text in *shape* to fit within its bounding box."""
    if not shape.has_text_frame:
        return None
    tf = shape.text_frame
    text = tf.text.strip()
    if not text:
        return None

    width = shape.width
    height = shape.height
    if honor_margins:
        width -= (tf.margin_left or 0) + (tf.margin_right or 0)
        height -= (tf.margin_top or 0) + (tf.margin_bottom or 0)
    width = max(width, 100000)
    height = max(height, 100000)

    best = best_fit_font_size(
        text, width, height,
        font_path=font_path, max_size=max_size, min_size=min_size,
        word_wrap=word_wrap,
    )
    if best is None:
        return None

    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(best)
        if para.font.size is None or para.font.size > Pt(best):
            para.font.size = Pt(best)

    return best


def fit_table_cell(cell, cell_width, cell_height, max_size=16, min_size=8,
                   font_path=None):
    font_path = _resolve_font(font_path)
    if not font_path:
        return None

    tf = cell.text_frame
    text = tf.text.strip()
    if not text:
        return max_size

    margin = getattr(tf, 'margin_left', 100800) or 100800
    width = cell_width - margin * 2
    height = cell_height - margin * 2
    width = max(width, 50000)
    height = max(height, 50000)

    best = best_fit_font_size(
        text, width, height,
        font_path=font_path, max_size=max_size, min_size=min_size,
        word_wrap=True,
    )
    if best is None:
        return min_size

    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(best)
        para.font.size = Pt(best)
    return best


def fit_table(table, max_size=16, min_size=8, font_path=None,
              uniform=True):
    font_path = _resolve_font(font_path)
    if not font_path:
        return None

    col_widths = [col.width for col in table.columns]
    row_heights = [row.height for row in table.rows]

    sizes = []
    for r_idx, row in enumerate(table.rows):
        row_sizes = []
        for c_idx, cell in enumerate(row.cells):
            s = fit_table_cell(
                cell, col_widths[c_idx], row_heights[r_idx],
                max_size=max_size, min_size=min_size,
                font_path=font_path,
            )
            row_sizes.append(s)
        sizes.append(row_sizes)

    if uniform:
        flat = [s for row in sizes for s in row if s is not None]
        if not flat:
            return None
        chosen = min(flat)
        for row in table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(chosen)
                    para.font.size = Pt(chosen)
        return chosen
    return sizes


def auto_fit_table_columns(table, headers, rows, header_pt=16, data_pt=14,
                           font_path=None, min_col_width_emu=300000):
    """Adjust table column widths based on content."""
    font_path = _resolve_font(font_path)
    if not font_path:
        return

    num_cols = len(headers)
    cell_margin = 100800
    padding_px = 8

    def _text_width_emu(text, pt_size, is_bold=False):
        effective_size = pt_size + 2 if is_bold else pt_size
        total_w = sum(_pptx_char_width(ch, effective_size) for ch in str(text))
        return _px_to_emu(total_w * 96 / 72) + _px_to_emu(padding_px)

    col_widths = []
    for c in range(num_cols):
        header_w = _text_width_emu(headers[c], header_pt, is_bold=True)
        data_w = max(
            (_text_width_emu(row[c], data_pt) for row in rows),
            default=0,
        )
        col_widths.append(max(header_w, data_w) + cell_margin * 2)

    total_measured = sum(col_widths)
    table_width = sum(col.width for col in table.columns)

    available = table_width - min_col_width_emu * num_cols
    if available <= 0:
        return

    scaled = []
    for w in col_widths:
        scaled.append(min_col_width_emu + int(w / total_measured * available))

    diff = table_width - sum(scaled)
    if diff != 0:
        widest = scaled.index(max(scaled))
        scaled[widest] += diff

    for i, w in enumerate(scaled):
        table.columns[i].width = w
