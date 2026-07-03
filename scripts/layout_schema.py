#!/usr/bin/env python3
import os, sys, json, copy
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
PREDEFINED_REGIONS = {"full": {"left": "0%", "top": "0%", "width": "100%", "height": "100%"}, "top-header": {"left": "0%", "top": "0%", "width": "100%", "height": "15%"}, "left-half": {"left": "0%", "top": "0%", "width": "50%", "height": "100%"}, "right-half": {"left": "50%", "top": "0%", "width": "50%", "height": "100%"}, "left-third": {"left": "0%", "top": "0%", "width": "33%", "height": "100%"}, "center-third": {"left": "33%", "top": "0%", "width": "34%", "height": "100%"}, "right-third": {"left": "67%", "top": "0%", "width": "33%", "height": "100%"}, "top-third": {"left": "0%", "top": "0%", "width": "100%", "height": "33%"}, "middle-third": {"left": "0%", "top": "33%", "width": "100%", "height": "34%"}, "bottom-third": {"left": "0%", "top": "67%", "width": "100%", "height": "33%"}, "top-left": {"left": "0%", "top": "0%", "width": "50%", "height": "50%"}, "top-right": {"left": "50%", "top": "0%", "width": "50%", "height": "50%"}, "bottom-left": {"left": "0%", "top": "50%", "width": "50%", "height": "50%"}, "bottom-right": {"left": "50%", "top": "50%", "width": "50%", "height": "50%"}, "center": {"left": "20%", "top": "30%", "width": "60%", "height": "40%"}}
def _parse_relative(value, total):
    if isinstance(value, (int, float)): return int(value)
    v = str(value).strip()
    if v.endswith("%"): return int(total * float(v[:-1]) / 100.0)
    if v.endswith("px"): return int(float(v[:-2]) * 9525)
    if v.endswith("in"): return int(float(v[:-2]) * 914400)
    if v.endswith("cm"): return int(float(v[:-2]) * 360000)
    if v.endswith("mm"): return int(float(v[:-2]) * 36000)
    return int(float(v))
def resolve_region(region_spec, slide_width, slide_height):
    if isinstance(region_spec, str): region_spec = PREDEFINED_REGIONS.get(region_spec, PREDEFINED_REGIONS["full"])
    region = {"left": _parse_relative(region_spec.get("left", "0%"), slide_width), "top": _parse_relative(region_spec.get("top", "0%"), slide_height), "width": _parse_relative(region_spec.get("width", "100%"), slide_width), "height": _parse_relative(region_spec.get("height", "100%"), slide_height)}
    if "min_height" in region_spec:
        region["min_height"] = _parse_relative(region_spec["min_height"], slide_height)
    return region
def _hex_to_rgb(hex_str): hex_str = hex_str.lstrip("#"); return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
DEFAULT_STYLES = {"heading-1": {"font_size": 48, "bold": True, "color": "#1F2937", "font_name": "Microsoft YaHei", "alignment": "left"}, "body": {"font_size": 30, "bold": False, "color": "#1F2937", "font_name": "Microsoft YaHei", "alignment": "left"}, "inverse-heading": {"font_size": 48, "bold": True, "color": "#FFFFFF", "font_name": "Microsoft YaHei", "alignment": "left"}, "inverse-body": {"font_size": 30, "bold": False, "color": "#FFFFFF", "font_name": "Microsoft YaHei", "alignment": "left"}, "accent-line": {"fill": "#3B82F6"}, "big-number": {"font_size": 80, "bold": True, "color": "#3B82F6", "font_name": "Microsoft YaHei", "alignment": "center"}, "label": {"font_size": 24, "bold": False, "color": "#1F2937", "font_name": "Microsoft YaHei", "alignment": "center"}}
def resolve_style(style_name, theme=None):
    if isinstance(style_name, dict):
        style = copy.deepcopy(style_name)
    else:
        style = copy.deepcopy(DEFAULT_STYLES.get(style_name, {}))
        if theme and "styles" in theme:
            style.update(theme["styles"].get(style_name, {}))
    return style
def _apply_text_style(paragraph, text, style):
    run = paragraph.add_run(); run.text = str(text) if text else ""
    if "font_size" in style: run.font.size = Pt(style["font_size"])
    if "bold" in style: run.font.bold = style["bold"]
    if "color" in style: run.font.color.rgb = _hex_to_rgb(style["color"])
    if "font_name" in style: run.font.name = style["font_name"]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    if "alignment" in style: paragraph.alignment = align_map.get(style["alignment"], PP_ALIGN.LEFT)
    return run
import re

_URL_RE = re.compile(r'https?://[^\s<>"\')\]]+')


def _add_hyperlink_runs_lschema(p, text, base_style):
    """Add runs to paragraph, converting URLs to hyperlinks."""
    pos = 0
    base_size = base_style.get("font_size", 18)
    base_color = base_style.get("color", "#1F2937")
    base_font = base_style.get("font_name", "Microsoft YaHei")
    for m in _URL_RE.finditer(text):
        if m.start() > pos:
            run = p.add_run()
            run.text = text[pos:m.start()]
            run.font.size = Pt(base_size)
            run.font.color.rgb = _hex_to_rgb(base_color)
            run.font.name = base_font
        run = p.add_run()
        run.text = m.group(0)
        run.font.size = Pt(base_size)
        run.font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)
        run.font.underline = True
        run.font.name = base_font
        run.hyperlink.address = m.group(0)
        pos = m.end()
    if pos < len(text):
        run = p.add_run()
        run.text = text[pos:]
        run.font.size = Pt(base_size)
        run.font.color.rgb = _hex_to_rgb(base_color)
        run.font.name = base_font


def render_text_shape(slide, region, content, style, adaptive=None):
    left, top, width, height = region["left"], region["top"], region["width"], region["height"]
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame; tf.word_wrap = True; tf.clear()
    if isinstance(content, list) and content and isinstance(content[0], dict):
        p = tf.paragraphs[0]
        for run_spec in content:
            if isinstance(run_spec, str):
                run = p.add_run(); run.text = run_spec
                if "font_size" in style: run.font.size = Pt(style["font_size"])
                if style.get("bold", False): run.font.bold = True
                if "color" in style: run.font.color.rgb = _hex_to_rgb(style["color"])
                if "font_name" in style: run.font.name = style["font_name"]
                continue
            run = p.add_run(); run.text = str(run_spec.get("text", ""))
            size = run_spec.get("size", style.get("font_size", 18))
            if size: run.font.size = Pt(size)
            if run_spec.get("bold", style.get("bold", False)): run.font.bold = True
            color = run_spec.get("color", style.get("color", "#1F2937"))
            if color: run.font.color.rgb = _hex_to_rgb(color)
            font_name = run_spec.get("font_name", style.get("font_name", "Microsoft YaHei"))
            if font_name: run.font.name = font_name
        align = style.get("alignment", "left")
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    elif isinstance(content, list):
        in_code = False
        code_lang = ""
        code_start_para_idx = -1
        base_size = style.get("font_size", 18)
        base_color = style.get("color", "#1F2937")
        base_font = style.get("font_name", "Microsoft YaHei")
        code_font_size = max(base_size - 2, 10)
        code_blocks = []  # list of (start_para_idx, end_para_idx, lang, first_code_para)
        para_idx = 0
        first_code_para = None
        for i, item in enumerate(content):
            raw = str(item)
            stripped = raw.strip()
            # Code block start marker
            if stripped.startswith('[') and stripped.endswith(']'):
                lang = stripped[1:-1]
                if lang in ('code', 'python', 'js', 'java', 'cpp', 'json', 'yaml', 'xml', 'sql', 'bash', 'sh'):
                    in_code = True
                    code_lang = lang
                    code_start_para_idx = para_idx
                    continue
                elif lang == '/code':
                    in_code = False
                    if code_start_para_idx >= 0:
                        code_blocks.append((code_start_para_idx, para_idx - 1, code_lang, first_code_para))
                        code_start_para_idx = -1
                    code_lang = ""
                    first_code_para = None
                    continue
            # Normal content paragraph
            p = tf.paragraphs[0] if para_idx == 0 else tf.add_paragraph()
            para_idx += 1
            # Preserve leading spaces for code, strip for normal content
            s = raw.rstrip() if in_code else stripped
            # Indent prefix (only for non-code)
            indent_match = None if in_code else re.match(r'^indent:(\d+)\|(.*)$', s)
            if indent_match:
                p.level = min(int(indent_match.group(1)), 8)
                s = indent_match.group(2)
            else:
                p.level = 0
            bullet = "" if in_code else "•"
            display = f"{bullet}{s}" if in_code else f"{bullet} {s}"
            if in_code:
                p.text = display
                for run in p.runs:
                    run.font.name = "Courier New"
                    run.font.size = Pt(code_font_size)
                    run.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
                if first_code_para is None:
                    first_code_para = p
            elif _URL_RE.search(display):
                p.text = ""
                _add_hyperlink_runs_lschema(p, display, style)
            else:
                _apply_text_style(p, display, style)
            # Consistent line spacing for constraint layout accuracy
            p.line_spacing = 1.4
            p.space_before = Pt(0)
            p.space_after = Pt(0)
        # Add gray background shapes + language labels for code blocks
        line_height = int(code_font_size * 1.5 * 12700)
        for start_para_idx, end_para_idx, lang, fcp in code_blocks:
            block_top = top + int(start_para_idx * line_height) + 30000
            block_height = int((end_para_idx - start_para_idx + 1) * line_height) + 30000
            bg_left = left + 20000
            bg_width = width - 40000
            # Background shape — rectangle (no rounded corners)
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                bg_left, block_top, bg_width, block_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
            bg.line.fill.background()
            # Label bar: same-width rectangle flush with top of background
            bar_height = int(line_height * 0.82)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                bg_left, block_top, bg_width, bar_height
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
            bar.line.fill.background()
            bar_tf = bar.text_frame
            bar_tf.word_wrap = False
            bar_p = bar_tf.paragraphs[0]
            bar_p.text = lang
            bar_p.alignment = PP_ALIGN.LEFT
            bar_p.font.name = "Courier New"
            bar_p.font.size = Pt(code_font_size)
            bar_p.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
            bar_p.font.bold = True
            # Small gap between label bar and code: add space_after to spacer row
            spacer_para = None
            if start_para_idx < len(tf.paragraphs):
                spacer_para = tf.paragraphs[start_para_idx]
            if spacer_para is not None:
                spacer_para.space_after = Pt(2)
            # Send background to back; bar stays on top
            spTree = slide.shapes._spTree
            spTree.remove(bg._element)
            spTree.insert(2, bg._element)
    else:
        p = tf.paragraphs[0]; _apply_text_style(p, str(content) if content else "", style)
    if adaptive and adaptive.get("strategy") == "shrink":
        try:
            from text_fitter import fit_shape_text
            max_size = adaptive.get("max_size", style.get("font_size", 72))
            fit_shape_text(shape, max_size=max_size, min_size=adaptive.get("min_size", 8))
        except Exception as e: print(f"  Warning: auto-fit failed: {e}")
    return shape
def render_shape_shape(slide, region, shape_type, style, content=None):
    left, top, width, height = region["left"], region["top"], region["width"], region["height"]
    mso_map = {"rectangle": MSO_SHAPE.RECTANGLE, "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE, "oval": MSO_SHAPE.OVAL, "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE, "diamond": MSO_SHAPE.DIAMOND, "pentagon": MSO_SHAPE.PENTAGON, "hexagon": MSO_SHAPE.HEXAGON, "star": MSO_SHAPE.STAR_5_POINT, "arrow_right": MSO_SHAPE.RIGHT_ARROW, "arrow_left": MSO_SHAPE.LEFT_ARROW, "arrow_up": MSO_SHAPE.UP_ARROW, "arrow_down": MSO_SHAPE.DOWN_ARROW, "chevron": MSO_SHAPE.CHEVRON, "parallelogram": MSO_SHAPE.PARALLELOGRAM, "trapezoid": MSO_SHAPE.TRAPEZOID, "donut": MSO_SHAPE.DONUT}
    mso_type = mso_map.get(shape_type, MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(mso_type, left, top, width, height)
    if "fill" in style: shape.fill.solid(); shape.fill.fore_color.rgb = _hex_to_rgb(style["fill"])
    if "border_color" in style: shape.line.color.rgb = _hex_to_rgb(style["border_color"])
    if "border_width" in style: shape.line.width = Pt(style["border_width"])
    else: shape.line.fill.background()
    if content:
        tf = shape.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        _apply_text_style(p, str(content), style)
        align = style.get("alignment", "center")
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
    return shape
def render_image_shape(slide, region, image_path):
    left, top, width, height = region["left"], region["top"], region["width"], region["height"]
    if not os.path.exists(image_path):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
        tf = shape.text_frame; p = tf.paragraphs[0]
        p.text = f"[Image: {os.path.basename(image_path)}]"; p.alignment = PP_ALIGN.CENTER
        return shape
    return slide.shapes.add_picture(image_path, left, top, width, height)
SCHEMAS_DIR = os.path.join(os.path.dirname(SCRIPT_DIR), "references", "layout_schemas")
def load_schema(name):
    path = os.path.join(SCHEMAS_DIR, f"{name}.json")
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f: return json.load(f)
    return None
def list_schemas():
    if not os.path.isdir(SCHEMAS_DIR): return []
    return [f.replace(".json", "") for f in os.listdir(SCHEMAS_DIR) if f.endswith(".json")]
def _get_shape_font_size(shape):
    """Extract the rendered font size (in points) from a text shape."""
    try:
        if shape.has_text_frame and shape.text_frame.paragraphs:
            for p in shape.text_frame.paragraphs:
                if p.runs and p.runs[0].font.size:
                    return p.runs[0].font.size.pt
    except Exception:
        pass
    return None


def render_slide_from_schema(slide, schema, slide_spec, prs, theme=None):
    sw, sh = prs.slide_width, prs.slide_height
    rendered, shape_registry = [], {}
    font_size_registry = {}  # role -> pt
    shapes_defs = schema.get("shapes", [])

    # ── Constraint layout mode ──
    if schema.get("layout_mode") == "constraint":
        try:
            sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
            from constraint_layout import solve_constraints
            solved, truncated_contents = solve_constraints(shapes_defs, slide_spec, sw, sh,
                                                           resolve_region, resolve_style, theme)
            for shape_def, region, style in solved:
                role = shape_def.get("role", "unknown")
                shape_type = shape_def.get("type", "textbox")
                content_source = shape_def.get("content_source")
                # Use truncated content if overflow occurred
                if role in truncated_contents:
                    content = truncated_contents[role]
                else:
                    content = slide_spec.get(content_source, "") if content_source else shape_def.get("content", "")
                if shape_type == "textbox":
                    shape = render_text_shape(slide, region, content, style, adaptive=None)
                elif shape_type == "shape":
                    shape = render_shape_shape(slide, region, shape_def.get("shape_type", "rectangle"), style, content)
                elif shape_type == "image":
                    shape = render_image_shape(slide, region, content)
                else:
                    shape = render_text_shape(slide, region, content, style, adaptive=None)
                rendered.append(shape)
                if role: shape_registry[role] = shape
                actual_size = _get_shape_font_size(shape)
                if actual_size is not None:
                    font_size_registry[role] = actual_size
            return rendered
        except Exception as e:
            print(f"  Warning: constraint layout failed ({e}), falling back to fixed layout")

    # ── Fixed layout mode (default) ──
    # Separate into independent (no font_size_ref) and dependent shapes
    independent = []
    dependent = []
    for shape_def in shapes_defs:
        style = resolve_style(shape_def.get("style", "body"), theme)
        if style.get("font_size_ref"):
            dependent.append(shape_def)
        else:
            independent.append(shape_def)

    # Phase 1: render independent shapes
    for shape_def in independent:
        role = shape_def.get("role", "unknown")
        region_spec = shape_def.get("region", "full")
        attach = shape_def.get("attach")
        if attach and attach.get("to") in shape_registry:
            base_shape = shape_registry[attach["to"]]
            base_region = {"left": base_shape.left, "top": base_shape.top, "width": base_shape.width, "height": base_shape.height}
            region = resolve_region(region_spec, sw, sh)
            edge = attach.get("edge", "bottom")
            offset = _parse_relative(attach.get("offset", "0%"), sh if edge in ("top", "bottom") else sw)
            if edge == "bottom": region["top"] = base_region["top"] + base_region["height"] + offset
            elif edge == "top": region["top"] = base_region["top"] - region["height"] - offset
            elif edge == "right": region["left"] = base_region["left"] + base_region["width"] + offset
            elif edge == "left": region["left"] = base_region["left"] - region["width"] - offset
        else:
            region = resolve_region(region_spec, sw, sh)
        shape_type = shape_def.get("type", "textbox")
        style_name = shape_def.get("style", "body")
        style = resolve_style(style_name, theme)
        adaptive = shape_def.get("adaptive")
        content_source = shape_def.get("content_source")
        content = slide_spec.get(content_source, "") if content_source else shape_def.get("content", "")
        if shape_type == "textbox": shape = render_text_shape(slide, region, content, style, adaptive)
        elif shape_type == "shape": shape = render_shape_shape(slide, region, shape_def.get("shape_type", "rectangle"), style, content)
        elif shape_type == "image": shape = render_image_shape(slide, region, content)
        else: shape = render_text_shape(slide, region, content, style, adaptive)
        rendered.append(shape)
        if role: shape_registry[role] = shape
        actual_size = _get_shape_font_size(shape)
        if actual_size is not None:
            font_size_registry[role] = actual_size

    # Phase 2: render dependent shapes with relative font sizes
    for shape_def in dependent:
        role = shape_def.get("role", "unknown")
        style_name = shape_def.get("style", "body")
        style = resolve_style(style_name, theme)
        ref_role = style.get("font_size_ref")
        ratio = style.get("font_size_ratio", 1.0)
        base_size = font_size_registry.get(ref_role)
        if base_size is not None:
            style = copy.deepcopy(style)
            style["font_size"] = max(int(base_size * ratio), 8)
            style.pop("font_size_ref", None)
            style.pop("font_size_ratio", None)

        region_spec = shape_def.get("region", "full")
        attach = shape_def.get("attach")
        if attach and attach.get("to") in shape_registry:
            base_shape = shape_registry[attach["to"]]
            base_region = {"left": base_shape.left, "top": base_shape.top, "width": base_shape.width, "height": base_shape.height}
            region = resolve_region(region_spec, sw, sh)
            edge = attach.get("edge", "bottom")
            offset = _parse_relative(attach.get("offset", "0%"), sh if edge in ("top", "bottom") else sw)
            if edge == "bottom": region["top"] = base_region["top"] + base_region["height"] + offset
            elif edge == "top": region["top"] = base_region["top"] - region["height"] - offset
            elif edge == "right": region["left"] = base_region["left"] + base_region["width"] + offset
            elif edge == "left": region["left"] = base_region["left"] - region["width"] - offset
        else:
            region = resolve_region(region_spec, sw, sh)
        shape_type = shape_def.get("type", "textbox")
        adaptive = shape_def.get("adaptive")
        content_source = shape_def.get("content_source")
        content = slide_spec.get(content_source, "") if content_source else shape_def.get("content", "")
        if shape_type == "textbox": shape = render_text_shape(slide, region, content, style, adaptive)
        elif shape_type == "shape": shape = render_shape_shape(slide, region, shape_def.get("shape_type", "rectangle"), style, content)
        elif shape_type == "image": shape = render_image_shape(slide, region, content)
        else: shape = render_text_shape(slide, region, content, style, adaptive)
        rendered.append(shape)
        if role: shape_registry[role] = shape
        actual_size = _get_shape_font_size(shape)
        if actual_size is not None:
            font_size_registry[role] = actual_size

    # Phase 4: uniform font size groups
    uniform_groups = {}
    for shape_def in shapes_defs:
        group = shape_def.get("uniform_group")
        role = shape_def.get("role")
        if group and role:
            uniform_groups.setdefault(group, []).append(role)
    for group_name, roles in uniform_groups.items():
        sizes = [font_size_registry.get(r) for r in roles]
        sizes = [s for s in sizes if s is not None]
        if not sizes:
            continue
        unified_size = min(sizes)
        for r in roles:
            shape = shape_registry.get(r)
            if not shape or not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size is not None:
                        run.font.size = Pt(unified_size)
            font_size_registry[r] = unified_size

    return rendered