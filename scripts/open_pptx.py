#!/usr/bin/env python3
"""
Open/read an existing PPTX and output its structure as JSON.
Usage:
  python open_pptx.py input.pptx [--output structure.json]
"""
import sys
import json
import argparse
import os

from pptx import Presentation
from pptx.util import Inches, Pt

def extract_structure(pptx_path):
    prs = Presentation(pptx_path)
    result = {
        "file": pptx_path,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slide_count": len(prs.slides),
        "slides": []
    }
    for idx, slide in enumerate(prs.slides):
        slide_info = {
            "index": idx + 1,
            "name": slide.name,
            "layout_name": slide.slide_layout.name if slide.slide_layout else None,
            "shapes": []
        }
        for shape in slide.shapes:
            shape_info = {
                "type": str(shape.shape_type),
                "has_text_frame": shape.has_text_frame,
                "text": shape.text_frame.text.strip() if shape.has_text_frame else "",
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
            if shape.has_chart:
                chart = shape.chart
                shape_info["chart_type"] = str(chart.chart_type)
                if hasattr(chart, 'plots') and chart.plots:
                    plot = chart.plots[0]
                    cats = []
                    if hasattr(plot, 'categories') and plot.categories:
                        cats = [str(c.label) for c in plot.categories]
                    vals = []
                    if hasattr(plot, 'series') and plot.series:
                        for s in plot.series:
                            vals.append({"name": s.name, "values": list(s.values) if hasattr(s, 'values') else []})
                    shape_info["chart_categories"] = cats
                    shape_info["chart_series"] = vals
            if shape.has_table:
                table = shape.table
                rows = []
                for r in table.rows:
                    row_text = [cell.text.strip() for cell in r.cells]
                    rows.append(row_text)
                shape_info["table_rows"] = rows
            slide_info["shapes"].append(shape_info)
        result["slides"].append(slide_info)
    return result

def main():
    parser = argparse.ArgumentParser(description="Open and analyze a PPTX file")
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument("--output", "-o", help="Output JSON file (default: print to stdout)")
    args = parser.parse_args()

    result = extract_structure(args.pptx)
    json_str = json.dumps(result, ensure_ascii=False, indent=2)

    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(json_str)
        print(f"Structure saved to: {args.output}")
    else:
        print(json_str)

if __name__ == "__main__":
    main()
