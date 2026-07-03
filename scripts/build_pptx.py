#!/usr/bin/env python3
"""Build a PPTX from a JSON deck spec using a template."""
import sys, json, argparse, copy, os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

sys.path.insert(0, SCRIPT_DIR)
try:
    from animator import animate_presentation
    HAS_ANIMATOR = True
except Exception:
    HAS_ANIMATOR = False

# Import advanced charts
try:
    from advanced_charts import add_waterfall_chart, add_funnel_chart, add_gantt_chart, generate_wordcloud
    HAS_ADVANCED = True
except Exception:
    HAS_ADVANCED = False

# Import layout schema engine
try:
    import layout_schema
    HAS_SCHEMA = True
except Exception as e:
    HAS_SCHEMA = False
    print(f"Warning: layout_schema not available: {e}")

# Import intent router
try:
    import intent_router
    HAS_ROUTER = True
except Exception as e:
    HAS_ROUTER = False
    print(f"Warning: intent_router not available: {e}")

# Import theme engine
try:
    from theme_engine import apply_theme
    HAS_THEME = True
except Exception:
    HAS_THEME = False

# Import color guard for contrast audit
try:
    from color_guard import analyze_presentation
    HAS_COLOR_GUARD = True
except Exception:
    HAS_COLOR_GUARD = False

# Import text fitter for auto-sizing fonts
try:
    from text_fitter import fit_table, fit_shape_text, best_fit_font_size
    HAS_FITTER = True
except Exception as e:
    HAS_FITTER = False
    print(f"Warning: text_fitter not available: {e}")

# Import adaptive layout strategies
try:
    import adaptive
    HAS_ADAPTIVE = True
except Exception as e:
    HAS_ADAPTIVE = False

# Import plugin engine
try:
    from plugin_engine import load_plugins, run_hooks
    HAS_PLUGINS = True
except Exception as e:
    HAS_PLUGINS = False

# Import deck logic for conditions and loops
try:
    from deck_logic import evaluate_deck
    HAS_DECK_LOGIC = True
except Exception as e:
    HAS_DECK_LOGIC = False

PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
DEFAULT_TEMPLATE = os.path.join(PROJECT_DIR, "templates", "base_template.pptx")

DARK_NAVY = RGBColor(0x25, 0x63, 0xEB)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
CHART_COLORS = [
    ACCENT_BLUE, RGBColor(0x10, 0xB9, 0x81), RGBColor(0xD9, 0x77, 0x06),
    RGBColor(0xB9, 0x1C, 0x1C), RGBColor(0x7C, 0x3A, 0xED), RGBColor(0xDB, 0x27, 0x77),
]

# Tree node colors: darker shades that all work with white text (contrast >= 4.5)
TREE_COLORS = [
    RGBColor(0x25, 0x63, 0xEB),  # blue
    RGBColor(0x04, 0x78, 0x57),  # emerald
    RGBColor(0xB4, 0x53, 0x09),  # orange
    RGBColor(0xDC, 0x26, 0x26),  # red
    RGBColor(0x7C, 0x3A, 0xED),  # purple
    RGBColor(0xBE, 0x12, 0x3C),  # rose
]

def find_template_slide(prs, name):
    for idx, slide in enumerate(prs.slides):
        if slide.name == name:
            return idx, slide
    return None, None

def clone_slide(prs, template_slide):
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    new_slide = prs.slides.add_slide(blank_layout)
    new_slide.name = template_slide.name
    for shape in template_slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide

def replace_placeholder_text(slide, old_text, new_text, color=None, size_pt=None,
                              auto_fit=False, auto_fit_max=None, auto_fit_min=14):
    """Replace placeholder text in a slide.

    *new_text* can be a plain string or a list of run dicts for rich text.

    If auto_fit is True and text_fitter is available, the font size will be
    adjusted so the text fits within its shape bounds.  auto_fit_max defaults
    to size_pt when not provided.

    For titles (size_pt >= 28), we first try single-line mode (word_wrap=False)
    to keep titles on one line.  If that shrinks the font below a comfortable
    threshold (28 pt), we fall back to multi-line but extend the shape so it
    doesn't collide with shapes below it.
    """
    target_shape = None
    is_rich = isinstance(new_text, list)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    if is_rich:
                        run.text = run.text.replace(old_text, "")
                        if not run.text:
                            paragraph._p.remove(run._r)
                        add_rich_text(shape, new_text, default_size=size_pt or 16,
                                      default_color=color or DARK_TEXT)
                    else:
                        run.text = run.text.replace(old_text, new_text)
                        if color:
                            run.font.color.rgb = color
                        if size_pt:
                            run.font.size = Pt(size_pt)
                        if not run.font.name:
                            run.font.name = "Microsoft YaHei"
                    target_shape = shape
            if old_text in paragraph.text:
                if is_rich:
                    paragraph.clear()
                    add_rich_text(shape, new_text, default_size=size_pt or 16,
                                  default_color=color or DARK_TEXT)
                    target_shape = shape
                elif len(paragraph.runs) == 0:
                    paragraph.text = paragraph.text.replace(old_text, new_text)
                    target_shape = shape
                else:
                    full_text = paragraph.text
                    paragraph.clear()
                    run = paragraph.add_run()
                    run.text = full_text.replace(old_text, new_text)
                    run.font.size = Pt(size_pt) if size_pt else Pt(16)
                    run.font.color.rgb = color if color else DARK_TEXT
                    run.font.name = "Microsoft YaHei"
                    target_shape = shape
    # Auto-fit after replacement so we measure the actual text
    if auto_fit and HAS_FITTER and target_shape is not None:
        max_sz = auto_fit_max if auto_fit_max is not None else (size_pt or 48)
        is_title = (size_pt or 0) >= 28

        if is_title:
            # Try single-line first for titles
            single_size = fit_shape_text(target_shape, max_size=max_sz, min_size=auto_fit_min,
                                         word_wrap=False)
            if single_size is not None:
                # Title fits on single line at some size — good
                pass
            else:
                # Even min_size doesn't fit on one line.
                # Force min_size single-line; _adjust_slide_for_overflow will
                # extend the shape if text actually wraps and pushes into others.
                tf = target_shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(auto_fit_min)
        else:
            # Non-title: normal word-wrap fit
            fit_shape_text(target_shape, max_size=max_sz, min_size=auto_fit_min,
                           word_wrap=True)

def _parse_color(color_val):
    """Parse a color value into RGBColor.

    Supports:
      - "#RRGGBB" hex string
      - "RRGGBB" hex string (no hash)
      - dict {"r": R, "g": G, "b": B}
      - RGBColor instance (pass through)
    """
    if color_val is None:
        return None
    if isinstance(color_val, RGBColor):
        return color_val
    if isinstance(color_val, dict):
        return RGBColor(color_val["r"], color_val["g"], color_val["b"])
    s = str(color_val).strip()
    if s.startswith("#"):
        s = s[1:]
    if len(s) == 6:
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    return None


def _url_pattern():
    import re
    return re.compile(r'https?://[^\s<>"\')\]]+')


def _add_hyperlink_runs(paragraph, text, font_size, font_color, font_name, default_size, default_color, default_font):
    """Add text to paragraph, converting URLs to hyperlinks."""
    url_re = _url_pattern()
    pos = 0
    for m in url_re.finditer(text):
        # Text before URL
        if m.start() > pos:
            run = paragraph.add_run()
            run.text = text[pos:m.start()]
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color
            run.font.name = font_name
        # URL as hyperlink
        run = paragraph.add_run()
        run.text = m.group(0)
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)  # blue
        run.font.underline = True
        run.font.name = font_name
        run.hyperlink.address = m.group(0)
        pos = m.end()
    # Remaining text
    if pos < len(text):
        run = paragraph.add_run()
        run.text = text[pos:]
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.name = font_name


def add_rich_text(shape, content, default_size=16, default_color=DARK_TEXT,
                  default_font="Microsoft YaHei", auto_fit=False):
    """Populate *shape* with rich text.

    *content* can be:
      - a plain string (single run, default style)
      - a list of run dicts: [{"text": "...", "bold": true, "color": "#FF0000", "size": 18}, ...]
      - a list of strings (bullet points, one per paragraph)
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()

    # Bullet list
    if isinstance(content, list) and content and isinstance(content[0], str):
        in_code_block = False
        code_lang = ""
        for i, item in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            clean = _clean_bullet_prefix(item)

            # Check for code block markers
            if clean.startswith('[') and clean.endswith(']') and not clean.startswith('[indent:'):
                lang = clean[1:-1]
                if lang in ('code', 'python', 'js', 'java', 'cpp', 'json', 'yaml', 'xml', 'sql', 'bash', 'sh'):
                    in_code_block = True
                    code_lang = lang
                    p.text = f"• [{code_lang}]"
                    p.font.size = Pt(default_size)
                    p.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)  # gray
                    p.font.name = "Courier New"
                    p.space_after = Pt(4)
                    continue
                elif lang == '/code':
                    in_code_block = False
                    code_lang = ""
                    p.text = "• [/code]"
                    p.font.size = Pt(default_size)
                    p.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
                    p.font.name = "Courier New"
                    p.space_after = Pt(8)
                    continue

            # Check for indent prefix: indent:N|text
            indent_level = 0
            indent_match = re.match(r'^indent:(\d+)\|(.*)$', clean)
            if indent_match:
                indent_level = int(indent_match.group(1))
                clean = indent_match.group(2)
                p.level = min(indent_level, 8)

            bullet_char = "•" if not in_code_block else " "
            display_text = f"{bullet_char} {clean}"

            if in_code_block:
                # Code line: monospace, slightly smaller, gray text
                run = p.add_run()
                run.text = display_text
                run.font.size = Pt(max(default_size - 2, 10))
                run.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
                run.font.name = "Courier New"
                p.space_after = Pt(2)
            elif _url_pattern().search(display_text):
                # Text with URLs — add hyperlink runs
                p.text = ""
                _add_hyperlink_runs(p, display_text, default_size, default_color, default_font,
                                    default_size, default_color, default_font)
                p.space_after = Pt(8)
            else:
                p.text = display_text
                p.font.size = Pt(default_size)
                p.font.color.rgb = default_color
                p.font.name = default_font
                p.space_after = Pt(8)
        if auto_fit and HAS_FITTER:
            fit_shape_text(shape, max_size=default_size, min_size=12, word_wrap=True)
        return

    # Single plain string
    if isinstance(content, str):
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(default_size)
        p.font.color.rgb = default_color
        p.font.name = default_font
        if auto_fit and HAS_FITTER:
            fit_shape_text(shape, max_size=default_size, min_size=12, word_wrap=True)
        return

    # Rich runs
    if isinstance(content, list):
        p = tf.paragraphs[0]
        first_size = None
        first_color = None
        first_font = None
        for run_spec in content:
            if isinstance(run_spec, str):
                run_spec = {"text": run_spec}
            text = run_spec.get("text", "")
            # Split on newlines to create multiple paragraphs
            parts = text.split("\n")
            for idx, part in enumerate(parts):
                if idx > 0:
                    # Start a new paragraph for text after \n
                    p = tf.add_paragraph()
                    first_size = None
                    first_color = None
                    first_font = None
                if not part:
                    continue
                run = p.add_run()
                run.text = part
                if run_spec.get("bold"):
                    run.font.bold = True
                if run_spec.get("italic"):
                    run.font.italic = True
                if run_spec.get("underline"):
                    run.font.underline = True
                color = _parse_color(run_spec.get("color"))
                if color:
                    run.font.color.rgb = color
                    if first_color is None:
                        first_color = color
                size = run_spec.get("size")
                if size:
                    run.font.size = Pt(size)
                    if first_size is None:
                        first_size = Pt(size)
                font_name = run_spec.get("font", default_font)
                if font_name:
                    run.font.name = font_name
                    if first_font is None:
                        first_font = font_name
        # Set paragraph-level defaults from first explicitly-styled run
        if p.runs:
            p.font.size = first_size or Pt(default_size)
            p.font.color.rgb = first_color or default_color
            p.font.name = first_font or default_font
        if auto_fit and HAS_FITTER:
            fit_shape_text(shape, max_size=first_size.pt if first_size else default_size,
                           min_size=12, word_wrap=True)
        return


def add_bullet_points(shape, items, font_size=30, auto_fit=False):
    tf = shape.text_frame
    tf.clear()
    for i, item in enumerate(items):
        clean = _clean_bullet_prefix(item)
        if not clean:
            continue
        p = tf.paragraphs[0] if i == 0 and len(tf.paragraphs) > 0 else tf.add_paragraph()
        # Check for indent prefix
        indent_match = re.match(r'^indent:(\d+)\|(.*)$', clean)
        if indent_match:
            p.level = min(int(indent_match.group(1)), 8)
            clean = indent_match.group(2)
        p.text = f"• {clean}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(8)
    if auto_fit and HAS_FITTER:
        fit_shape_text(shape, max_size=font_size, min_size=12, word_wrap=True)

def find_placeholder_shape(slide, keyword):
    for shape in slide.shapes:
        if shape.has_text_frame and keyword in shape.text_frame.text:
            return shape
    return None

def remove_placeholder_shape(slide, keyword):
    shape = find_placeholder_shape(slide, keyword)
    if shape:
        sp = shape.element
        sp.getparent().remove(sp)

def _measure_text_width_emu(text, font_name, font_size_pt):
    """Approximate text width in EMU using Pillow."""
    from PIL import ImageDraw, ImageFont, Image
    import os
    px_size = round(font_size_pt * 96 / 72)
    candidates = [
        os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "fonts", "NotoSansCJKsc-Regular.otf"),
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    font = None
    for fp in candidates:
        if os.path.exists(fp):
            try:
                font = ImageFont.truetype(fp, px_size)
                break
            except Exception:
                pass
    if font is None:
        return None
    draw = ImageDraw.Draw(Image.new("RGB", (1, 1)))
    try:
        bbox = draw.textbbox((0, 0), text, font=font, anchor='lt')
        width_px = bbox[2] - bbox[0]
        return int(width_px * 914400 / 96)
    except Exception:
        return None


def _shape_actual_text_height(shape, safety_factor=0.95):
    """Calculate the actual rendered text height for *shape* in EMU.

    Uses the empirical PowerPoint width model from text_fitter.
    Handles rich text with mixed font sizes paragraph-by-paragraph.
    Returns 0 if shape has no text.
    """
    if not shape.has_text_frame:
        return 0
    tf = shape.text_frame
    text = tf.text.strip()
    if not text:
        return 0
    if not HAS_FITTER:
        return 0

    from text_fitter import _pptx_wrap_text, _pptx_text_height, _emu_to_px

    width = shape.width
    width -= (tf.margin_left or 0) + (tf.margin_right or 0)
    width = max(width, 100000)
    width = int(width * safety_factor)
    max_w_pt = _emu_to_px(width) * 72 / 96  # px → pt

    total_h_pt = 0.0
    for para in tf.paragraphs:
        para_text = para.text
        if not para_text:
            continue
        # Determine effective font size for this paragraph
        para_size = 16
        for run in para.runs:
            if run.font.size:
                pt = run.font.size.pt if hasattr(run.font.size, 'pt') else int(run.font.size) // 12700
                para_size = pt
                break
        lines = _pptx_wrap_text(para_text, para_size, max_w_pt)
        if lines:
            total_h_pt += _pptx_text_height(lines, para_size)

    if total_h_pt <= 0:
        return 0
    # pt → EMU  (1 pt = 12700 EMU)
    return int(total_h_pt * 12700)


def _adjust_slide_for_overflow(slide):
    """Post-process a slide: detect text overflow and apply adaptive strategies.

    Uses the adaptive module for enhanced overflow handling:
    - shrink font size to fit
    - multi-column for long bullet lists
    - expand shape height and push down overlapping shapes
    - truncate with ... as last resort
    """
    if HAS_ADAPTIVE:
        adaptive.adjust_slide_for_overflow(slide, strategy="auto")
        return

    # Fallback to legacy behavior if adaptive module unavailable
    text_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            if shape._element.getparent().tag.endswith('}tc'):
                continue
        except Exception:
            pass
        text_shapes.append(shape)

    if not text_shapes:
        return

    text_shapes.sort(key=lambda s: s.top)

    if HAS_FITTER:
        for shape in text_shapes:
            tf = shape.text_frame
            text = tf.text.strip()
            if not text:
                continue
            actual_h = _shape_actual_text_height(shape)
            if actual_h <= 0 or actual_h <= shape.height:
                continue
            current_size = 16
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pt = run.font.size.pt if hasattr(run.font.size, 'pt') else int(run.font.size) // 12700
                        current_size = pt
                        break
                break
            fit_shape_text(shape, max_size=current_size, min_size=10, word_wrap=True)


def _clean_bullet_prefix(item):
    """Remove leading bullet markers from an item to avoid double bullets."""
    s = str(item).strip()
    # Remove common bullet prefixes
    for prefix in ('• ', '· ', '- ', '* ', '•', '·', '-', '*'):
        if s.startswith(prefix):
            s = s[len(prefix):].strip()
            break
    return s

def sync_accent_line(slide, slide_spec=None, title_text=""):
    """Sync accent line: match title text width, align left with title,
    and position just below the title to avoid overlap.

    If slide_spec is provided and ``accent_line`` is False, the accent line
    is removed instead of being synced.
    """
    if slide_spec is not None and slide_spec.get("accent_line", True) is False:
        # Remove accent line entirely
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                text = shape.text_frame.text.strip() if shape.has_text_frame else ""
                if text:
                    continue
                try:
                    if shape.fill.type == 1 and shape.fill.fore_color.rgb == ACCENT_BLUE:
                        if shape.width > shape.height * 3:
                            sp = shape.element
                            sp.getparent().remove(sp)
                            return
                except Exception:
                    pass
        return

    title_shape = None
    title_pt = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and (not title_text or text == title_text.strip()):
                max_size = 0
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            pt = run.font.size.pt if hasattr(run.font.size, 'pt') else int(run.font.size) // 12700
                            max_size = max(max_size, pt)
                if max_size >= 28:
                    title_shape = shape
                    title_pt = max_size
                    break
    if not title_shape:
        return

    # Measure actual text width
    text_width = _measure_text_width_emu(
        title_shape.text_frame.text.strip(),
        title_shape.text_frame.paragraphs[0].runs[0].font.name or "Microsoft YaHei",
        title_pt,
    )
    if text_width is None:
        text_width = title_shape.width
    else:
        # Compensate for font differences: Pillow uses Noto Sans CJK / DejaVu,
        # but PowerPoint renders with Microsoft YaHei which is slightly wider for Latin chars.
        text_width = int(text_width * 1.04)

    title_bottom = title_shape.top + title_shape.height
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            text = shape.text_frame.text.strip() if shape.has_text_frame else ""
            if text:
                continue
            try:
                if shape.fill.type == 1 and shape.fill.fore_color.rgb == ACCENT_BLUE:
                    # Only adjust horizontal accent lines (wide & short), not vertical ones
                    if shape.width > shape.height * 3:
                        # Ensure the line is directly below the title (not timeline bars etc.)
                        if abs(shape.top - title_bottom) < 400000 and abs(shape.left - title_shape.left) < 200000:
                            shape.width = text_width
                            shape.left = title_shape.left + title_shape.text_frame.margin_left
                            # Only move accent line down if it overlaps the title text.
                            # If the template placed it inside the title area (common
                            # design), preserve that original relative position.
                            if shape.top > title_bottom - 50000:
                                shape.top = title_bottom + 80000
                            return
            except Exception:
                pass


def add_native_chart(slide, chart_data, placeholder_shape):
    if not chart_data or "values" not in chart_data:
        return
    categories = chart_data.get("categories", [])
    values = chart_data["values"]
    chart_type = chart_data.get("type", "column")
    series_name = chart_data.get("series_name", "数据")

    chart_map = {
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "area": XL_CHART_TYPE.AREA,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "radar": XL_CHART_TYPE.RADAR_MARKERS,
        "bubble": XL_CHART_TYPE.BUBBLE,
    }
    xl_type = chart_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Handle scatter and bubble with special data formats
    if chart_type == "scatter":
        from pptx.chart.data import XyChartData
        cd = XyChartData()
        series = cd.add_series(series_name)
        for i, v in enumerate(values):
            x = i + 1
            series.add_data_point(x, v)
    elif chart_type == "bubble":
        from pptx.chart.data import BubbleChartData
        cd = BubbleChartData()
        series = cd.add_series(series_name)
        for i, v in enumerate(values):
            x = i + 1
            size = max(v * 10, 1)
            series.add_data_point(x, v, size)
    else:
        cd = ChartData()
        if categories:
            cd.categories = categories
        cd.add_series(series_name, values)

    x, y, cx, cy = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    chart = slide.shapes.add_chart(xl_type, x, y, cx, cy, cd).chart
    chart.has_legend = True

    # For pie/doughnut charts, place legend on the right side to make best
    # use of the landscape slide aspect ratio.  Other charts keep bottom legend.
    if chart_type in ("pie", "doughnut"):
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
    else:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    if hasattr(plot, 'series'):
        for i, series in enumerate(plot.series):
            if chart_type in ("pie", "doughnut"):
                # WPS compatibility: set colors per-point (c:dPt) rather than
                # per-series, so each slice gets its own color.
                for pt_idx, point in enumerate(series.points):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = CHART_COLORS[pt_idx % len(CHART_COLORS)]
            elif chart_type == "radar":
                # Radar: aggressive no-fill for cross-app consistency
                # Some viewers (QQ) ignore series-level noFill, so we also
                # set point-level noFill.
                # NOTE: do NOT use series.format API here because it can
                # overwrite the XML we carefully construct.
                from lxml import etree
                ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
                ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                spPr = series._element.find(f'{{{ns_c}}}spPr')
                if spPr is None:
                    spPr = etree.SubElement(series._element, f'{{{ns_c}}}spPr')
                # Remove any existing fill children
                for child in list(spPr):
                    if child.tag.endswith('}noFill') or child.tag.endswith('}solidFill') or child.tag.endswith('}gradFill') or child.tag.endswith('}pattFill'):
                        spPr.remove(child)
                # Add transparent solidFill
                solidFill = etree.SubElement(spPr, f'{{{ns_a}}}solidFill')
                srgbClr = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
                srgbClr.set('val', 'FFFFFF')
                alpha = etree.SubElement(srgbClr, f'{{{ns_a}}}alpha')
                alpha.set('val', '0')
                # Ensure line element exists with correct color
                ln = spPr.find(f'{{{ns_a}}}ln')
                if ln is None:
                    ln = etree.SubElement(spPr, f'{{{ns_a}}}ln')
                    ln.set('w', '28575')  # 2.25 pt in EMUs
                else:
                    # Remove existing line fill
                    for child in list(ln):
                        if child.tag.endswith('}solidFill') or child.tag.endswith('}noFill'):
                            ln.remove(child)
                lnSolid = etree.SubElement(ln, f'{{{ns_a}}}solidFill')
                lnSrgb = etree.SubElement(lnSolid, f'{{{ns_a}}}srgbClr')
                lnSrgb.set('val', str(CHART_COLORS[i % len(CHART_COLORS)]))
                # Point level: set each point to noFill
                for point in series.points:
                    point.format.fill.background()
            else:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = CHART_COLORS[i % len(CHART_COLORS)]

def add_native_table(slide, table_data, placeholder_shape):
    if not table_data or "headers" not in table_data or "rows" not in table_data:
        return
    headers = table_data["headers"]
    rows = table_data["rows"]
    num_rows = 1 + len(rows)
    num_cols = len(headers)

    x, y, cx, cy = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    table = slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy).table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(24)
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = "Microsoft YaHei"

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = DARK_TEXT
            p.font.size = Pt(22)
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = "Microsoft YaHei"
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    # Uniform column widths + fixed font sizes
    col_w = cx // num_cols
    for c in table.columns:
        c.width = col_w
    for r in table.rows:
        r.height = cy // num_rows

    for cell in table.rows[0].cells:
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(24)
            para.font.size = Pt(24)
    for r_idx in range(1, len(table.rows)):
        for cell in table.rows[r_idx].cells:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(22)
                para.font.size = Pt(22)

def add_timeline(slide, timeline_data, placeholder_shape):
    if not timeline_data or "milestones" not in timeline_data:
        return
    milestones = timeline_data["milestones"]
    if not milestones:
        return
    x, y, w, h = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    count = len(milestones)

    # Remove template auto-shapes (accent line, pre-drawn bar) so we draw fresh
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            text = shape.text_frame.text.strip() if shape.has_text_frame else ""
            if not text:
                sp = shape.element
                sp.getparent().remove(sp)

    TIMELINE_BAR_BLUE = RGBColor(0x93, 0xC5, 0xFD)
    NODE_BLUE = RGBColor(0x25, 0x63, 0xEB)

    # Bar spans the full placeholder width for a long, confident look
    bar_top = y + int(h * 0.48)
    bar_h = int(h * 0.015)
    bar_start = x
    bar_end = x + w
    bar_length = bar_end - bar_start

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_start, bar_top, bar_length, bar_h)
    line.fill.solid()
    line.fill.fore_color.rgb = TIMELINE_BAR_BLUE
    line.line.fill.background()

    # Arrow at the end of the bar
    arrow_size = int(bar_h * 4)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        bar_end - arrow_size, bar_top - (arrow_size - bar_h) // 2,
        arrow_size, arrow_size,
    )
    arrow.rotation = 90
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = NODE_BLUE
    arrow.line.fill.background()

    # Node radius 3.5% of placeholder height (diameter 7%)
    node_r = int(h * 0.035)

    # spacing: divide placeholder width into  equal segments so the
    # last milestone has the same right-hand room as every other milestone.
    spacing = w // count

    for i, ms in enumerate(milestones):
        mx = bar_start + spacing * i + spacing // 2
        my = bar_top - node_r

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, mx - node_r, my, node_r * 2, node_r * 2)
        circle.fill.solid()
        circle.fill.fore_color.rgb = WHITE
        circle.line.color.rgb = NODE_BLUE
        circle.line.width = Pt(2.5)

        is_top = (i % 2 == 0)

        # Text box spans one full spacing interval, centred on the node
        box_left = bar_start + spacing * i
        box_width = spacing

        if is_top:
            lbl_top = y + int(h * 0.04)
            title_top = y + int(h * 0.14)
            desc_top = y + int(h * 0.26)
        else:
            lbl_top = bar_top + bar_h + int(h * 0.04)
            title_top = bar_top + bar_h + int(h * 0.14)
            desc_top = bar_top + bar_h + int(h * 0.26)

        # Decorative lead line from node to text area
        lead_w = int(Pt(1.5))
        lead_x = mx - lead_w // 2
        lead_gap = int(h * 0.03)
        if is_top:
            lead_y = my - lead_gap
            lead_h = lead_gap
        else:
            lead_y = my + node_r * 2
            lead_h = lead_gap
        lead = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lead_x, lead_y, lead_w, lead_h)
        lead.fill.solid()
        lead.fill.fore_color.rgb = NODE_BLUE
        lead.line.fill.background()

        lbl = slide.shapes.add_textbox(box_left, lbl_top, box_width, int(h * 0.10))
        p = lbl.text_frame.paragraphs[0]
        p.text = ms.get("date", "")
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = NODE_BLUE
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER

        title = slide.shapes.add_textbox(box_left, title_top, box_width, int(h * 0.10))
        tp = title.text_frame.paragraphs[0]
        tp.text = ms.get("title", "")
        tp.font.bold = True
        tp.font.size = Pt(18)
        tp.font.color.rgb = DARK_TEXT
        tp.font.name = "Microsoft YaHei"
        tp.alignment = PP_ALIGN.CENTER

        desc = slide.shapes.add_textbox(box_left, desc_top, box_width, int(h * 0.22))
        p2 = desc.text_frame.paragraphs[0]
        p2.text = ms.get("description", ms.get("desc", ""))
        p2.font.size = Pt(16)
        p2.font.color.rgb = DARK_TEXT
        p2.font.name = "Microsoft YaHei"
        p2.alignment = PP_ALIGN.CENTER
        desc.text_frame.word_wrap = True
def add_process(slide, process_data, placeholder_shape):
    if not process_data or "steps" not in process_data:
        return
    steps = process_data["steps"]
    if not steps:
        return
    x, y, w, h = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    count = len(steps)
    gap = int(w * 0.02)
    step_w = (w - gap * (count - 1)) // count
    step_h = int(h * 0.6)
    step_y = y + int(h * 0.15)

    # Dark saturated colors for strong white-text contrast (all >= 3:1)
    PROCESS_COLORS = [
        RGBColor(0x06, 0x5F, 0x46),  # emerald-800
        RGBColor(0x92, 0x40, 0x0E),  # amber-800
        RGBColor(0xB9, 0x1C, 0x1C),  # red-700
        RGBColor(0x5B, 0x21, 0xB6),  # violet-700
    ]
    for i, step in enumerate(steps):
        sx = x + i * (step_w + gap)
        color = PROCESS_COLORS[i % len(PROCESS_COLORS)]
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, step_y, step_w, step_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        num_shape = slide.shapes.add_textbox(sx, step_y + int(step_h * 0.1), step_w, int(step_h * 0.3))
        np = num_shape.text_frame.paragraphs[0]
        np.text = str(i + 1)
        np.font.size = Pt(36)
        np.font.bold = True
        np.font.color.rgb = WHITE
        np.font.name = "Microsoft YaHei"
        np.alignment = PP_ALIGN.CENTER

        title_shape = slide.shapes.add_textbox(sx, step_y + int(step_h * 0.45), step_w, int(step_h * 0.5))
        tp = title_shape.text_frame
        tp.word_wrap = True
        tp.paragraphs[0].text = step.get("title", "")
        tp.paragraphs[0].font.size = Pt(22)
        tp.paragraphs[0].font.bold = True
        tp.paragraphs[0].font.color.rgb = WHITE
        tp.paragraphs[0].font.name = "Microsoft YaHei"
        tp.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_team_cards(slide, team_data, placeholder_shape):
    if not team_data or "members" not in team_data:
        return
    members = team_data["members"]
    if not members:
        return
    x, y, w, h = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    count = len(members)
    gap = int(w * 0.03)
    card_w = (w - gap * (count - 1)) // count
    card_h = int(h * 0.85)
    card_y = y + int(h * 0.05)

    for i, member in enumerate(members):
        cx = x + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()

        avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx + int(card_w * 0.35), card_y + int(card_h * 0.12),
                                         int(card_w * 0.3), int(card_w * 0.3))
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = ACCENT_BLUE
        avatar.line.fill.background()

        name_shape = slide.shapes.add_textbox(cx, card_y + int(card_h * 0.48), card_w, int(card_h * 0.14))
        np = name_shape.text_frame.paragraphs[0]
        np.text = member.get("name", "")
        np.font.size = Pt(24)
        np.font.bold = True
        np.font.color.rgb = DARK_TEXT
        np.font.name = "Microsoft YaHei"
        np.alignment = PP_ALIGN.CENTER

        role_shape = slide.shapes.add_textbox(cx, card_y + int(card_h * 0.63), card_w, int(card_h * 0.12))
        rp = role_shape.text_frame.paragraphs[0]
        rp.text = member.get("role", "")
        rp.font.size = Pt(18)
        rp.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
        rp.font.name = "Microsoft YaHei"
        rp.alignment = PP_ALIGN.CENTER

        desc_shape = slide.shapes.add_textbox(cx + int(card_w * 0.08), card_y + int(card_h * 0.76),
                                              int(card_w * 0.84), int(card_h * 0.20))
        dp = desc_shape.text_frame
        dp.word_wrap = True
        dp.paragraphs[0].text = member.get("desc", "")
        dp.paragraphs[0].font.size = Pt(16)
        dp.paragraphs[0].font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
        dp.paragraphs[0].font.name = "Microsoft YaHei"
        dp.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_tree_diagram(slide, tree_data, placeholder_shape):
    """Draw a 2-level tree diagram: root node on top, children in a row below."""
    if not tree_data or "root" not in tree_data:
        return
    root_label = tree_data["root"]
    children = tree_data.get("children", [])
    if not children:
        return
    x, y, w, h = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    count = len(children)

    # Node dimensions
    node_w = min(int(w * 0.22), int(w // count * 0.85))
    node_h = int(h * 0.22)
    root_w = min(int(w * 0.35), node_w * 1.5)
    root_h = node_h

    # Positions
    root_x = x + (w - root_w) // 2
    root_y = y + int(h * 0.05)
    child_y = y + int(h * 0.55)
    total_children_w = count * node_w + (count - 1) * int(w * 0.03)
    child_start_x = x + (w - total_children_w) // 2

    # Tree nodes use a consistent white text color across all nodes.
    # TREE_COLORS are chosen so every color has white contrast >= 4.5:1.
    tree_text_color = WHITE

    # Use a darker blue that won't be affected by any theme remap issues
    root_color = RGBColor(0x1E, 0x40, 0xAF)
    root_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, root_x, root_y, root_w, root_h)
    root_box.fill.solid()
    root_box.fill.fore_color.rgb = root_color
    root_tf = root_box.text_frame
    root_tf.word_wrap = True
    root_p = root_tf.paragraphs[0]
    root_p.text = root_label
    root_p.font.size = Pt(18)
    root_p.font.bold = True
    root_p.font.name = "Microsoft YaHei"
    root_p.alignment = PP_ALIGN.CENTER
    for run in root_p.runs:
        run.font.color.rgb = tree_text_color

    # Children
    for i, child in enumerate(children):
        cx = child_start_x + i * (node_w + int(w * 0.03))
        child_color = TREE_COLORS[(i + 1) % len(TREE_COLORS)]
        child_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, child_y, node_w, node_h)
        child_box.fill.solid()
        child_box.fill.fore_color.rgb = child_color
        child_box.line.fill.background()
        ctf = child_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        if isinstance(child, dict):
            cp.text = child.get("label", child.get("name", ""))
        else:
            cp.text = str(child)
        cp.font.size = Pt(14)
        cp.font.bold = True
        cp.font.name = "Microsoft YaHei"
        cp.alignment = PP_ALIGN.CENTER
        for run in cp.runs:
            run.font.color.rgb = tree_text_color

        # Connector line: from bottom-center of root to top-center of child
        root_bottom_x = root_x + root_w // 2
        root_bottom_y = root_y + root_h
        child_top_x = cx + node_w // 2
        child_top_y = child_y
        # Draw elbow connector (vertical then horizontal then vertical)
        mid_y = root_bottom_y + (child_top_y - root_bottom_y) // 2
        # Vertical from root
        v1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     root_bottom_x - int(Pt(1.5)) // 2, root_bottom_y,
                                     int(Pt(1.5)), mid_y - root_bottom_y)
        v1.fill.solid()
        v1.fill.fore_color.rgb = MID_GRAY
        v1.line.fill.background()
        # Horizontal across
        h_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         min(root_bottom_x, child_top_x), mid_y - int(Pt(1.5)) // 2,
                                         abs(child_top_x - root_bottom_x), int(Pt(1.5)))
        h_line.fill.solid()
        h_line.fill.fore_color.rgb = MID_GRAY
        h_line.line.fill.background()
        # Vertical to child
        v2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     child_top_x - int(Pt(1.5)) // 2, mid_y,
                                     int(Pt(1.5)), child_top_y - mid_y)
        v2.fill.solid()
        v2.fill.fore_color.rgb = MID_GRAY
        v2.line.fill.background()

def add_image_placeholder(slide, image_path, placeholder_shape, mode="fit"):
    """Add an image to the slide, replacing a placeholder shape.

    mode:
      "fit"    - scale to fit within placeholder, preserving aspect ratio, centered.
      "stretch"- stretch to exactly match placeholder dimensions.
    """
    if not image_path or not os.path.exists(image_path):
        return
    x, y, cx, cy = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    if mode == "stretch":
        slide.shapes.add_picture(image_path, x, y, cx, cy)
    else:
        from PIL import Image as PILImage
        with PILImage.open(image_path) as img:
            img_w, img_h = img.size
        if img_w <= 0 or img_h <= 0:
            slide.shapes.add_picture(image_path, x, y, cx, cy)
            return
        ratio_img = img_w / img_h
        ratio_ph = cx / cy
        if ratio_img > ratio_ph:
            new_cx = cx
            new_cy = int(cx / ratio_img)
        else:
            new_cy = cy
            new_cx = int(cy * ratio_img)
        new_x = x + (cx - new_cx) // 2
        new_y = y + (cy - new_cy) // 2
        slide.shapes.add_picture(image_path, new_x, new_y, new_cx, new_cy)


def add_hyperlink(shape, url):
    """Add a hyperlink to a shape or its text."""
    if not url:
        return
    if shape.has_text_frame and shape.text_frame.text:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.hyperlink.address = url
    else:
        try:
            shape.click_action.hyperlink.address = url
        except Exception:
            pass

def add_video(slide, video_path, placeholder_shape):
    """Add a video to the slide."""
    if not video_path or not os.path.exists(video_path):
        return
    x, y, cx, cy = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    try:
        slide.shapes.add_movie(video_path, x, y, cx, cy)
    except Exception as e:
        print(f"Warning: could not add video {video_path}: {e}")

def build_deck(deck_spec, template_path, output_path):
    # Load plugins
    plugins = load_plugins() if HAS_PLUGINS else []
    context = {"deck_spec": deck_spec, "template_path": template_path}

    # pre_build hook
    if plugins:
        deck_spec = run_hooks(plugins, 'pre_build', deck_spec, context)

    template_prs = Presentation(template_path)
    out_prs = Presentation()
    out_prs.slide_width = template_prs.slide_width
    out_prs.slide_height = template_prs.slide_height
    context["prs"] = out_prs

    total_slides = len(deck_spec.get("slides", []))
    for idx, slide_spec in enumerate(deck_spec.get("slides", [])):
        slide_context = dict(context, slide_number=idx + 1, total_slides=total_slides, slide_index=idx)
        layout_type = slide_spec.get("layout", "title_content")
        template_name = slide_spec.get("template_slide_name", layout_type)
        # Fallback for layouts not present in base template
        if template_name not in ("cover", "toc", "title_content", "two_column", "chart", "closing",
                                  "table", "timeline", "image_content", "quote", "team",
                                  "data_highlight", "process", "hero_top", "three_column",
                                  "asymmetric", "single_focus", "l_shape", "t_shape",
                                  "symmetric", "section", "comparison", "mixed_grid",
                                  "primary_secondary", "waterfall_layout", "image_hero",
                                  "people", "timeline_h",
                                  "title_content_constraint", "two_column_constraint",
                                  "cover_constraint", "quote_constraint",
                                  "data_highlight_constraint", "section_constraint",
                                  "three_column_constraint", "hero_top_constraint",
                                  "image_hero_constraint", "side_by_side_constraint",
                                  "chart_constraint", "table_constraint",
                                  "timeline_constraint", "team_constraint",
                                  "process_constraint", "tree_constraint",
                                  "image_text_split_constraint"):
            template_name = "title_content"

        # Try schema-based rendering first
        schema_rendered = False
        if HAS_SCHEMA:
            schema = layout_schema.load_schema(layout_type)
            if schema:
                blank_layout = out_prs.slide_layouts[6] if len(out_prs.slide_layouts) > 6 else out_prs.slide_layouts[-1]
                new_slide = out_prs.slides.add_slide(blank_layout)
                new_slide.name = layout_type
                layout_schema.render_slide_from_schema(new_slide, schema, slide_spec, out_prs)
                schema_rendered = True
                # Still process free-form shapes and hyperlinks below

        if not schema_rendered:
            idx, template_slide = find_template_slide(template_prs, template_name)
            if template_slide is None:
                print(f"Warning: template slide '{template_name}' not found, skipping")
                continue

            new_slide = clone_slide(out_prs, template_slide)

            # === cover ===
        if layout_type == "cover":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=WHITE, size_pt=54, auto_fit=True)
            replace_placeholder_text(new_slide, "__SUBTITLE__", slide_spec.get("subtitle", ""), color=WHITE, size_pt=28, auto_fit=True, auto_fit_min=12)
            replace_placeholder_text(new_slide, "__DATE__", slide_spec.get("date", ""), color=WHITE, size_pt=20)
            # Center cover content vertically: shrink title height, then
            # shift title/subtitle/date/accent-bar as a group so the block
            # sits in the vertical middle of the slide.
            text_shapes = [s for s in new_slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
            if len(text_shapes) >= 2:
                text_shapes.sort(key=lambda s: s.top)
                title_shape = text_shapes[0]
                subtitle_shape = text_shapes[1] if len(text_shapes) > 1 else None
                date_shape = text_shapes[2] if len(text_shapes) > 2 else None
                # Shrink title height to actual rendered text + padding
                actual_h = _shape_actual_text_height(title_shape, safety_factor=1.0)
                if actual_h > 0 and actual_h < title_shape.height:
                    title_shape.height = actual_h + 120000
                # Compute content block top/bottom
                content_top = title_shape.top
                content_bottom = date_shape.top + date_shape.height if date_shape else subtitle_shape.top + subtitle_shape.height
                slide_h = out_prs.slide_height
                target_center = slide_h // 2
                current_center = (content_top + content_bottom) // 2
                shift = target_center - current_center
                # Move all cover text shapes first, then sync the vertical
                # accent bar so it aligns exactly with the title top.
                bar_shape = None
                for shape in new_slide.shapes:
                    is_text = shape.has_text_frame and shape.text_frame.text.strip()
                    is_bar = False
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        try:
                            if shape.fill.type == 1 and shape.fill.fore_color.rgb == ACCENT_BLUE:
                                if shape.height > shape.width * 3:
                                    is_bar = True
                        except Exception:
                            pass
                    if is_text:
                        shape.top += shift
                    elif is_bar:
                        bar_shape = shape
                # Sync bar after all text shapes have moved
                if bar_shape and title_shape and subtitle_shape:
                    bar_top = title_shape.top
                    bar_bottom = date_shape.top + date_shape.height if date_shape else subtitle_shape.top + subtitle_shape.height
                    bar_shape.top = bar_top
                    bar_shape.height = bar_bottom - bar_top

        # === toc ===

        elif layout_type == "toc":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", "目录"), color=DARK_TEXT, size_pt=48, auto_fit=True)
            items = slide_spec.get("items", [])
            for shape in new_slide.shapes:
                if shape.has_text_frame and "__ITEMS__" in shape.text_frame.text:
                    add_bullet_points(shape, items, font_size=30)
                    break
            sync_accent_line(new_slide, slide_spec)

        # === title_content ===

        elif layout_type == "title_content":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            content = slide_spec.get("content", "")
            for shape in new_slide.shapes:
                if shape.has_text_frame and "__CONTENT__" in shape.text_frame.text:
                    if isinstance(content, list) and content and isinstance(content[0], str):
                        add_bullet_points(shape, content, font_size=30, auto_fit=True)
                    else:
                        add_rich_text(shape, content, default_size=24, default_color=DARK_TEXT, auto_fit=True)
                    break
            sync_accent_line(new_slide, slide_spec)

        # === two_column ===

        elif layout_type == "two_column":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            left = slide_spec.get("left_content", "")
            right = slide_spec.get("right_content", "")
            for shape in new_slide.shapes:
                if shape.has_text_frame and "__LEFT__" in shape.text_frame.text:
                    if isinstance(left, list) and left and isinstance(left[0], str):
                        add_bullet_points(shape, left, font_size=28, auto_fit=True)
                    else:
                        add_rich_text(shape, left, default_size=28, default_color=DARK_TEXT, auto_fit=True)
                if shape.has_text_frame and "__RIGHT__" in shape.text_frame.text:
                    if isinstance(right, list) and right and isinstance(right[0], str):
                        add_bullet_points(shape, right, font_size=28, auto_fit=True)
                    else:
                        add_rich_text(shape, right, default_size=28, default_color=DARK_TEXT, auto_fit=True)
            sync_accent_line(new_slide, slide_spec)

        # === chart ===

        elif layout_type == "chart":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            replace_placeholder_text(new_slide, "__DESC__", slide_spec.get("description", ""), color=DARK_TEXT, size_pt=28, auto_fit=True, auto_fit_max=32, auto_fit_min=14)
            ph = find_placeholder_shape(new_slide, "__CHART__")
            if ph:
                remove_placeholder_shape(new_slide, "__CHART__")
                add_native_chart(new_slide, slide_spec.get("chart_data"), ph)
            sync_accent_line(new_slide, slide_spec)

        # === closing ===
        elif layout_type == "closing":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", "Thank You"), color=WHITE, size_pt=64, auto_fit=True)
            replace_placeholder_text(new_slide, "__SUBTITLE__", slide_spec.get("subtitle", ""), color=WHITE, size_pt=28, auto_fit=True, auto_fit_min=12)

        # === table ===

        elif layout_type == "table":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            replace_placeholder_text(new_slide, "__DESC__", slide_spec.get("description", ""), color=DARK_TEXT, size_pt=28, auto_fit=True, auto_fit_max=32, auto_fit_min=14)
            ph = find_placeholder_shape(new_slide, "__TABLE__")
            if ph:
                remove_placeholder_shape(new_slide, "__TABLE__")
                add_native_table(new_slide, slide_spec.get("table_data"), ph)
            sync_accent_line(new_slide, slide_spec)

        # === timeline ===

        elif layout_type == "timeline":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            ph = find_placeholder_shape(new_slide, "__TIMELINE__")
            if ph:
                remove_placeholder_shape(new_slide, "__TIMELINE__")
                add_timeline(new_slide, slide_spec.get("timeline_data"), ph)
            # Do NOT sync accent line on timeline slides — the timeline bar
            # itself is blue and would be mistaken for the accent decoration.

        # === image_content ===

        elif layout_type == "image_content":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            img_ph = find_placeholder_shape(new_slide, "__IMAGE__")
            if img_ph:
                # Also remove the gray background rectangle that sits under the
                # image placeholder (same position, solid light-gray fill).
                for shape in list(new_slide.shapes):
                    if shape == img_ph:
                        continue
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        if (shape.left == img_ph.left and shape.top == img_ph.top and
                            shape.width == img_ph.width and shape.height == img_ph.height):
                            try:
                                if shape.fill.type is not None and shape.fill.fore_color and shape.fill.fore_color.rgb:
                                    sp = shape.element
                                    sp.getparent().remove(sp)
                            except Exception:
                                pass
                remove_placeholder_shape(new_slide, "__IMAGE__")
                image_path = slide_spec.get("image_path")
                if image_path and os.path.exists(image_path):
                    add_image_placeholder(new_slide, image_path, img_ph)
                else:
                    # No image provided — draw a placeholder gray box with hint text
                    x, y, cx, cy = img_ph.left, img_ph.top, img_ph.width, img_ph.height
                    box = new_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
                    box.fill.solid()
                    box.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
                    box.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
                    box.line.width = Pt(1)
                    tf = box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = "[Image placeholder]"
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
                    p.font.name = "Microsoft YaHei"
                    p.alignment = PP_ALIGN.CENTER
            content = slide_spec.get("content", "")
            for shape in new_slide.shapes:
                if shape.has_text_frame and "__CONTENT__" in shape.text_frame.text:
                    if isinstance(content, list) and content and isinstance(content[0], str):
                        add_bullet_points(shape, content, font_size=28, auto_fit=True)
                    else:
                        add_rich_text(shape, content, default_size=22, default_color=DARK_TEXT, auto_fit=True)
                    break
            sync_accent_line(new_slide, slide_spec)

        # === quote ===
        elif layout_type == "quote":
            replace_placeholder_text(new_slide, "__QUOTE__", slide_spec.get("quote", ""), color=WHITE, size_pt=30)
            replace_placeholder_text(new_slide, "__AUTHOR__", slide_spec.get("author", ""), color=WHITE, size_pt=20)

        # === team ===

        elif layout_type == "team":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            ph = find_placeholder_shape(new_slide, "__TEAM__")
            if ph:
                remove_placeholder_shape(new_slide, "__TEAM__")
                add_team_cards(new_slide, slide_spec.get("team_data"), ph)
            sync_accent_line(new_slide, slide_spec)

        # === data_highlight ===

        elif layout_type == "data_highlight":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            replace_placeholder_text(new_slide, "__BIG_NUMBER__", slide_spec.get("big_number", ""), color=DARK_TEXT, size_pt=64, auto_fit=True, auto_fit_min=24)
            replace_placeholder_text(new_slide, "__LABEL__", slide_spec.get("label", ""), color=DARK_TEXT, size_pt=22, auto_fit=True, auto_fit_min=10)
            sync_accent_line(new_slide, slide_spec)

        # === process ===

        elif layout_type == "process":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            ph = find_placeholder_shape(new_slide, "__PROCESS__")
            if ph:
                remove_placeholder_shape(new_slide, "__PROCESS__")
                add_process(new_slide, slide_spec.get("process_data"), ph)
            sync_accent_line(new_slide, slide_spec)

        # === tree ===

        elif layout_type == "tree":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            ph = find_placeholder_shape(new_slide, "__TREE__")
            if not ph:
                ph = find_placeholder_shape(new_slide, "__CHART__")
            if not ph:
                ph = find_placeholder_shape(new_slide, "__CONTENT__")
            if ph:
                keyword = ph.text_frame.text.strip() if ph.has_text_frame else ""
                if keyword:
                    remove_placeholder_shape(new_slide, keyword)
                add_tree_diagram(new_slide, slide_spec.get("tree_data"), ph)
            sync_accent_line(new_slide, slide_spec)

        elif layout_type == "waterfall":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            replace_placeholder_text(new_slide, "__DESC__", slide_spec.get("description", ""), color=DARK_TEXT, size_pt=22, auto_fit=True, auto_fit_min=12)
            ph = find_placeholder_shape(new_slide, "__CHART__")
            if ph and HAS_ADVANCED:
                remove_placeholder_shape(new_slide, "__CHART__")
                add_waterfall_chart(new_slide, slide_spec.get("chart_data"), ph)
            sync_accent_line(new_slide, slide_spec)

        elif layout_type == "funnel":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            replace_placeholder_text(new_slide, "__DESC__", slide_spec.get("description", ""), color=DARK_TEXT, size_pt=22, auto_fit=True, auto_fit_min=12)
            ph = find_placeholder_shape(new_slide, "__CHART__")
            if ph and HAS_ADVANCED:
                remove_placeholder_shape(new_slide, "__CHART__")
                add_funnel_chart(new_slide, slide_spec.get("chart_data"), ph)
            sync_accent_line(new_slide, slide_spec)

        elif layout_type == "gantt":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            replace_placeholder_text(new_slide, "__DESC__", slide_spec.get("description", ""), color=DARK_TEXT, size_pt=22, auto_fit=True, auto_fit_min=12)
            ph = find_placeholder_shape(new_slide, "__CHART__")
            if ph and HAS_ADVANCED:
                remove_placeholder_shape(new_slide, "__CHART__")
                add_gantt_chart(new_slide, slide_spec.get("chart_data"), ph)
            sync_accent_line(new_slide, slide_spec)

        elif layout_type == "wordcloud":
            replace_placeholder_text(new_slide, "__TITLE__", slide_spec.get("title", ""), color=DARK_TEXT, size_pt=48, auto_fit=True)
            replace_placeholder_text(new_slide, "__DESC__", slide_spec.get("description", ""), color=DARK_TEXT, size_pt=22, auto_fit=True, auto_fit_min=12)

            # Remove content placeholder from fallback layouts (e.g. title_content -> __CONTENT__)
            remove_placeholder_shape(new_slide, "__CONTENT__")
            # Also remove any default text placeholders
            for shape in list(new_slide.shapes):
                if shape.has_text_frame and "Click to edit Master text" in shape.text_frame.text:
                    sp = shape._element
                    sp.getparent().remove(sp)

            # Try to place wordcloud image
            ph = find_placeholder_shape(new_slide, "__CHART__")
            if ph and HAS_ADVANCED:
                words = slide_spec.get("words", [])
                wc_path = os.path.join(os.path.dirname(output_path), "_wordcloud.png")
                result = generate_wordcloud(words, wc_path)
                if result:
                    remove_placeholder_shape(new_slide, "__CHART__")
                    add_image_placeholder(new_slide, wc_path, ph)
            elif HAS_ADVANCED:
                # No __CHART__ placeholder - add image directly centered
                words = slide_spec.get("words", [])
                wc_path = os.path.join(os.path.dirname(output_path), "_wordcloud.png")
                result = generate_wordcloud(words, wc_path)
                if result and os.path.exists(wc_path):
                    from pptx.util import Inches
                    slide_w = new_slide.shapes._spTree.getparent().getparent().attrib.get('cx', 12192000)
                    if isinstance(slide_w, str):
                        slide_w = int(slide_w)
                    slide_h = new_slide.shapes._spTree.getparent().getparent().attrib.get('cy', 6858000)
                    if isinstance(slide_h, str):
                        slide_h = int(slide_h)
                    img_w = int(slide_w * 0.7)
                    img_h = int(slide_h * 0.6)
                    left = (slide_w - img_w) // 2
                    top = int(slide_h * 0.25)
                    new_slide.shapes.add_picture(wc_path, left, top, img_w, img_h)
            sync_accent_line(new_slide, slide_spec)

        # Apply hyperlink if specified
        hyperlink = slide_spec.get("hyperlink")
        if hyperlink:
            for shape in new_slide.shapes:
                if shape.has_text_frame:
                    add_hyperlink(shape, hyperlink)

        # Apply video if specified (replaces image placeholder)
        video_path = slide_spec.get("video_path")
        if video_path:
            for shape in new_slide.shapes:
                if shape.has_text_frame and "__IMAGE__" in shape.text_frame.text:
                    remove_placeholder_shape(new_slide, "__IMAGE__")
                    add_video(new_slide, video_path, shape)
                    break

        # === Free-form shapes (user-defined absolute positioning) ===
        for shape_spec in slide_spec.get("shapes", []):
            _add_free_shape(new_slide, shape_spec, output_path)

        # Post-process: detect text overflow and adjust shape positions
        # Skip for schema-rendered (especially constraint mode) slides
        if not schema_rendered:
            _adjust_slide_for_overflow(new_slide)

        # Post-process: render chart placeholders for constraint-mode slides
        chart_ph = find_placeholder_shape(new_slide, "__CHART__")
        if chart_ph and slide_spec.get("chart_data"):
            remove_placeholder_shape(new_slide, "__CHART__")
            add_native_chart(new_slide, slide_spec["chart_data"], chart_ph)

        # Post-process: render table placeholders for constraint-mode slides
        table_ph = find_placeholder_shape(new_slide, "__TABLE__")
        if table_ph and slide_spec.get("table_data"):
            remove_placeholder_shape(new_slide, "__TABLE__")
            add_native_table(new_slide, slide_spec["table_data"], table_ph)

        # Post-process: render timeline placeholders for constraint-mode slides
        timeline_ph = find_placeholder_shape(new_slide, "__TIMELINE__")
        print(f"  [debug] timeline_ph={timeline_ph}, timeline_data={slide_spec.get('timeline_data') is not None}")
        if timeline_ph and slide_spec.get("timeline_data"):
            remove_placeholder_shape(new_slide, "__TIMELINE__")
            add_timeline(new_slide, slide_spec["timeline_data"], timeline_ph)

        # Post-process: render team placeholders for constraint-mode slides
        team_ph = find_placeholder_shape(new_slide, "__TEAM__")
        if team_ph and slide_spec.get("team_data"):
            remove_placeholder_shape(new_slide, "__TEAM__")
            add_team_cards(new_slide, slide_spec["team_data"], team_ph)

        # Post-process: render process placeholders for constraint-mode slides
        process_ph = find_placeholder_shape(new_slide, "__PROCESS__")
        if process_ph and slide_spec.get("process_data"):
            remove_placeholder_shape(new_slide, "__PROCESS__")
            add_process(new_slide, slide_spec["process_data"], process_ph)

        # Post-process: render tree placeholders for constraint-mode slides
        tree_ph = find_placeholder_shape(new_slide, "__TREE__")
        if tree_ph and slide_spec.get("tree_data"):
            remove_placeholder_shape(new_slide, "__TREE__")
            add_tree_diagram(new_slide, slide_spec["tree_data"], tree_ph)

        # Re-sync accent lines after overflow adjustments so they stay
        # attached to their titles even when shapes were extended.
        if layout_type != "timeline":
            sync_accent_line(new_slide, slide_spec)

        # post_slide hook
        if plugins:
            run_hooks(plugins, 'post_slide', slide_spec, new_slide, slide_context)

    # post_build hook
    if plugins:
        out_prs = run_hooks(plugins, 'post_build', out_prs, context)

    out_prs.save(output_path)
    print(f"Saved: {output_path} ({len(out_prs.slides)} slides)")

    # Font size guard check
    try:
        from font_size_guard import check_presentation, print_report
        issues = check_presentation(out_prs)
        print_report(issues)
    except Exception as e:
        print(f"  Font size guard skipped: {e}")

    return output_path


def _add_free_shape(slide, shape_spec, output_path):
    """Add a user-defined free-form shape to a slide.

    shape_spec fields:
      type: "text" | "image" | "shape" | "chart"
      left, top, width, height: integers in EMU
      --- text ---
      content: string | list of run dicts
      font_size, color, alignment: optional defaults
      --- image ---
      path: image file path
      mode: "fit" | "stretch" (default "fit")
      mask: "circle" | "rounded_rectangle" | "oval" | ...
      crop: {"left":0, "right":0, "top":0, "bottom":0}
      --- shape ---
      shape_type: "rectangle" | "rounded_rectangle" | "oval" | ...
      fill: "#RRGGBB" or {"r":R,"g":G,"b":B}
      border_color, border_width: optional
      --- chart ---
      chart_data: same format as layout chart_data
      --- all ---
      z_order: integer (higher = more front)
    """
    stype = shape_spec.get("type", "text")
    # Support relative coordinates (% of slide dimensions)
    # Resolve slide size from the presentation part
    try:
        prs = slide.part.package.presentation_part.presentation
        sw, sh = prs.slide_width, prs.slide_height
    except Exception:
        sw, sh = 12192000, 6858000  # fallback: 13.333x7.5in
    def _parse_coord(v, total):
        if isinstance(v, str) and v.endswith("%"):
            return int(total * float(v[:-1]) / 100.0)
        return int(v)
    x = _parse_coord(shape_spec.get("left", 0), sw)
    y = _parse_coord(shape_spec.get("top", 0), sh)
    w = _parse_coord(shape_spec.get("width", 1000000), sw)
    h = _parse_coord(shape_spec.get("height", 500000), sh)
    shape = None

    if stype == "text":
        shape = slide.shapes.add_textbox(x, y, w, h)
        content = shape_spec.get("content", "")
        align_map = {
            "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
        }
        align = align_map.get(shape_spec.get("alignment", "left"), PP_ALIGN.LEFT)
        shape.text_frame.word_wrap = True
        if shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].alignment = align
        add_rich_text(
            shape, content,
            default_size=shape_spec.get("font_size", 16),
            default_color=_parse_color(shape_spec.get("color")) or DARK_TEXT,
        )
        if shape_spec.get("vertical_center"):
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    elif stype == "image":
        img_path = shape_spec.get("path", "")
        mask_type = shape_spec.get("mask")
        crop = shape_spec.get("crop")

        if img_path and os.path.exists(img_path):
            if mask_type:
                shape = _add_masked_image(slide, img_path, mask_type, x, y, w, h)
            else:
                mode = shape_spec.get("mode", "fit")
                if mode == "stretch":
                    shape = slide.shapes.add_picture(img_path, x, y, w, h)
                else:
                    from PIL import Image as PILImage
                    with PILImage.open(img_path) as img:
                        img_w, img_h = img.size
                    if img_w > 0 and img_h > 0:
                        ratio_img = img_w / img_h
                        ratio_ph = w / h
                        if ratio_img > ratio_ph:
                            new_w = w
                            new_h = int(w / ratio_img)
                        else:
                            new_h = h
                            new_w = int(h * ratio_img)
                        new_x = x + (w - new_w) // 2
                        new_y = y + (h - new_h) // 2
                        shape = slide.shapes.add_picture(img_path, new_x, new_y, new_w, new_h)
                    else:
                        shape = slide.shapes.add_picture(img_path, x, y, w, h)

                # Apply rectangular crop if specified
                if crop and hasattr(shape, 'crop_left'):
                    shape.crop_left = crop.get("left", 0)
                    shape.crop_right = crop.get("right", 0)
                    shape.crop_top = crop.get("top", 0)
                    shape.crop_bottom = crop.get("bottom", 0)

    elif stype == "shape":
        shape_type = shape_spec.get("shape_type", "rectangle")
        mso_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "pentagon": MSO_SHAPE.PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "star": MSO_SHAPE.STAR_5_POINT,
            "arrow_right": MSO_SHAPE.RIGHT_ARROW,
            "arrow_left": MSO_SHAPE.LEFT_ARROW,
            "arrow_up": MSO_SHAPE.UP_ARROW,
            "arrow_down": MSO_SHAPE.DOWN_ARROW,
            "chevron": MSO_SHAPE.CHEVRON,
            "parallelogram": MSO_SHAPE.PARALLELOGRAM,
            "trapezoid": MSO_SHAPE.TRAPEZOID,
            "donut": MSO_SHAPE.DONUT,
        }
        mso = mso_map.get(shape_type, MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(mso, x, y, w, h)
        fill = _parse_color(shape_spec.get("fill"))
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        border = _parse_color(shape_spec.get("border_color"))
        if border:
            shape.line.color.rgb = border
        if shape_spec.get("border_width"):
            shape.line.width = Pt(shape_spec["border_width"])

        # Support text content on shapes
        content = shape_spec.get("content")
        if content:
            align_map = {
                "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
            }
            align = align_map.get(shape_spec.get("alignment", "left"), PP_ALIGN.LEFT)
            shape.text_frame.word_wrap = True
            add_rich_text(
                shape, content,
                default_size=shape_spec.get("font_size", 16),
                default_color=_parse_color(shape_spec.get("color")) or DARK_TEXT,
            )
            # Re-apply alignment after add_rich_text (which calls tf.clear())
            tf = shape.text_frame
            for p in tf.paragraphs:
                p.alignment = align
                p.space_before = 0
                p.space_after = 0
                p.line_spacing = 1.0
            if shape_spec.get("vertical_center"):
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_top = 0
                tf.margin_bottom = 0
                tf.margin_left = 0
                tf.margin_right = 0

    elif stype == "chart":
        chart_data = shape_spec.get("chart_data")
        if chart_data:
            class _Geom:
                def __init__(self, left, top, width, height):
                    self.left, self.top, self.width, self.height = left, top, width, height
            chart_area = _Geom(x, y, w, h)
            shape = add_native_chart(slide, chart_data, chart_area)

    # Apply z-order if specified (higher number = more front)
    z = shape_spec.get("z_order")
    if z is not None and shape is not None:
        _set_z_order(shape, z)


def _add_masked_image(slide, image_path, mask_type, x, y, w, h):
    """Add an image masked into a shape (e.g. circle, rounded rectangle).

    Works by creating a temporary picture to establish the image relationship,
    then creating a shape with blipFill pointing to the same image.
    """
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    # Create a dummy picture to get the image relationship ID
    pic = slide.shapes.add_picture(image_path, 0, 0, 1, 1)
    blip_el = pic._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
    rId = blip_el.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
    # Remove the dummy picture
    pic._element.getparent().remove(pic._element)

    # Map mask type to MSO shape
    mask_map = {
        "circle": MSO_SHAPE.OVAL,
        "oval": MSO_SHAPE.OVAL,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rectangle": MSO_SHAPE.RECTANGLE,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "diamond": MSO_SHAPE.DIAMOND,
        "hexagon": MSO_SHAPE.HEXAGON,
        "star": MSO_SHAPE.STAR_5_POINT,
        "pentagon": MSO_SHAPE.PENTAGON,
        "heart": MSO_SHAPE.HEART,
        "cloud": MSO_SHAPE.CLOUD,
        "sun": MSO_SHAPE.SUN,
        "moon": MSO_SHAPE.MOON,
    }
    mso = mask_map.get(mask_type, MSO_SHAPE.OVAL)

    shape = slide.shapes.add_shape(mso, x, y, w, h)
    spPr = shape._element.spPr
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Remove any existing fill
    for tag in ['noFill', 'solidFill', 'gradFill', 'blipFill', 'pattFill', 'grpFill']:
        child = spPr.find(f'{{{ns_a}}}{tag}')
        if child is not None:
            spPr.remove(child)

    # Add blipFill using the image relationship
    blipFill = etree.SubElement(spPr, f'{{{ns_a}}}blipFill')
    blip = etree.SubElement(blipFill, f'{{{ns_a}}}blip')
    blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', rId)
    stretch = etree.SubElement(blipFill, f'{{{ns_a}}}stretch')
    etree.SubElement(stretch, f'{{{ns_a}}}fillRect')

    # Remove outline for a cleaner look
    shape.line.fill.background()

    return shape


def _set_z_order(shape, z_order):
    """Move *shape* to the specified position in the slide's shape tree.

    z_order=0 means back (bottom), higher values mean front (top).
    The shape tree order determines visual stacking.
    """
    spTree = shape._parent._spTree
    sp = shape._element
    # Remove from current position
    spTree.remove(sp)
    # Get all current shape elements
    children = list(spTree)
    # Find where to insert (skip non-shape elements like extLst)
    shape_children = [c for c in children if c.tag.endswith(('}sp', '}pic', '}grpSp', '}cxnSp', '}graphicFrame'))]
    insert_idx = min(z_order, len(shape_children))
    # Insert before the element currently at that position
    if insert_idx < len(shape_children):
        target = shape_children[insert_idx]
        target_idx = children.index(target)
        spTree.insert(target_idx, sp)
    else:
        # Append before extLst if present
        extLst = spTree.find('{http://schemas.openxmlformats.org/presentationml/2006/main}extLst')
        if extLst is not None:
            extIdx = list(spTree).index(extLst)
            spTree.insert(extIdx, sp)
        else:
            spTree.append(sp)

def main():
    parser = argparse.ArgumentParser(description="Build PPTX from JSON spec")
    parser.add_argument("deck_json", help="Path to deck JSON file")
    parser.add_argument("output_pptx", help="Path to output PPTX file")
    parser.add_argument("--template", default=DEFAULT_TEMPLATE, help="Template PPTX path")
    parser.add_argument("--animate", choices=["fade", "fly_in", "wipe", "zoom", "bounce", "appear", "swivel", "float"],
                        help="Add entrance animation to all shapes")
    parser.add_argument("--no-stagger", action="store_true", help="Disable staggered animation delays")
    parser.add_argument("--theme", default="default", help="Apply theme (default/dark/warm/forest/minimal or JSON path)")
    parser.add_argument("--export", choices=["pdf", "long_image", "html"], help="Export to additional format after build")
    parser.add_argument("--export-output", help="Export output file path")
    parser.add_argument("--auto-route", action="store_true", help="Auto-infer layout from slide content fields")
    parser.add_argument("--audit", action="store_true", help="Run color contrast audit after build and print report")
    args = parser.parse_args()

    # Support both JSON and YAML input
    with open(args.deck_json, "r", encoding="utf-8") as f:
        if args.deck_json.lower().endswith((".yaml", ".yml")):
            try:
                import yaml
                deck_spec = yaml.safe_load(f)
            except ImportError:
                print("Error: PyYAML is required for YAML files. Install with: pip install pyyaml")
                sys.exit(1)
        else:
            deck_spec = json.load(f)

    # Evaluate conditions and loops
    if HAS_DECK_LOGIC:
        deck_spec = evaluate_deck(deck_spec)
        if deck_spec.get("vars"):
            print(f"Deck vars: {list(deck_spec['vars'].keys())}")

    # Auto-paginate long content slides
    if HAS_ADAPTIVE:
        original_count = len(deck_spec.get("slides", []))
        deck_spec = adaptive.auto_paginate_deck(deck_spec, max_items_per_slide=10)
        new_count = len(deck_spec.get("slides", []))
        if new_count > original_count:
            print(f"Auto-paginated: {original_count} -> {new_count} slides")

    # Auto-route layouts if requested
    if args.auto_route and HAS_ROUTER:
        deck_spec = intent_router.auto_route(deck_spec)
        for i, s in enumerate(deck_spec.get("slides", [])):
            if s.get("_inferred"):
                print(f"  Slide {i+1}: inferred layout '{s['layout']}' (conf={s['_confidence']})")

    build_deck(deck_spec, args.template, args.output_pptx)

    # Apply theme
    if args.theme and HAS_THEME:
        prs = Presentation(args.output_pptx)
        apply_theme(prs, args.theme)
        prs.save(args.output_pptx)
        print(f"Theme applied: {args.theme}")



    # Apply animation
    if args.animate and HAS_ANIMATOR:
        prs = Presentation(args.output_pptx)
        animate_presentation(prs, effect=args.animate, stagger=not args.no_stagger)
        prs.save(args.output_pptx)
        print(f"Animation applied: {args.animate}")

    # Contrast audit
    if args.audit and HAS_COLOR_GUARD:
        prs = Presentation(args.output_pptx)
        theme = None
        if args.theme and os.path.exists(args.theme):
            with open(args.theme, "r", encoding="utf-8") as f:
                theme = json.load(f)
        report = analyze_presentation(prs, theme=theme)
        print("\n" + report.markdown())

    # Export
    if args.export:
        export_script = os.path.join(SCRIPT_DIR, "export.py")
        if os.path.exists(export_script):
            import subprocess
            cmd = [sys.executable, export_script, args.output_pptx, "--format", args.export, "--output", args.export_output]
            subprocess.run(cmd)
        else:
            print(f"Warning: export.py not found at {export_script}")

if __name__ == "__main__":
    main()
