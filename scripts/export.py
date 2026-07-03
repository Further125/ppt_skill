#!/usr/bin/env python3
"""
Export PPTX to various formats: PDF, long image, HTML slideshow.
Usage:
  python export.py input.pptx --format pdf --output out.pdf
  python export.py input.pptx --format long_image --output out.png
  python export.py input.pptx --format html --output out.html
"""
import sys, os, argparse, subprocess, shutil, base64


from pptx import Presentation
from PIL import Image


def _libreoffice_env(work_dir):
    """Build a writable env for soffice (containers often have read-only HOME).

    In read-only sandboxes only /tmp is writable (tmpfs), so the LO
    user profile goes there instead of work_dir (which may be read-only).
    """
    env = dict(os.environ)
    profile_dir = os.path.join("/tmp", f"lo_profile_{os.getpid()}")
    os.makedirs(profile_dir, exist_ok=True)
    env["HOME"] = "/tmp"  # give LO a writable home (tmpfs)
    return env, profile_dir


def _run_soffice(cmd, work_dir, timeout=120):
    """Run soffice with writable HOME + UserInstallation, capture output, never hang."""
    env, profile_dir = _libreoffice_env(work_dir)
    cmd = list(cmd) + [f"-env:UserInstallation=file://{profile_dir}"]
    try:
        return subprocess.run(cmd, capture_output=True, text=True, env=env, timeout=timeout)
    except subprocess.TimeoutExpired:
        print("LibreOffice timed out, falling back...")
        return None


def export_pdf(pptx_path, output_path):
    """Export PPTX to PDF using LibreOffice (best) or fallback."""
    if shutil.which("soffice"):
        out_dir = os.path.dirname(os.path.abspath(output_path)) or "."
        os.makedirs(out_dir, exist_ok=True)
        base = os.path.basename(pptx_path)
        cmd = ["soffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_path]
        result = _run_soffice(cmd, out_dir, timeout=120)
        if result is not None and result.returncode == 0:
            # LibreOffice uses original filename
            generated = os.path.join(out_dir, base.replace(".pptx", ".pdf"))
            if os.path.exists(generated) and generated != output_path:
                os.rename(generated, output_path)
            print(f"PDF exported: {output_path}")
            return True
        else:
            print(f"LibreOffice error: {result.stderr}")
    # Fallback: use python-pptx to read structure and generate a minimal report
    print("LibreOffice not available, generating text report...")
    prs = Presentation(pptx_path)
    with open(output_path.replace(".pdf", ".txt"), "w", encoding="utf-8") as f:
        f.write(f"PPTX Export Report: {pptx_path}\n")
        f.write(f"Slides: {len(prs.slides)}\n\n")
        for i, slide in enumerate(prs.slides):
            f.write(f"--- Slide {i+1} ---\n")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        f.write(text + "\n")
            f.write("\n")
    print(f"Text report saved: {output_path.replace('.pdf', '.txt')}")
    return False


def export_long_image(pptx_path, output_path, dpi=144):
    """Export all slides as a vertically stitched long image."""
    prs = Presentation(pptx_path)
    slide_w = int(prs.slide_width / 914400 * dpi)
    slide_h = int(prs.slide_height / 914400 * dpi)

    total_h = slide_h * len(prs.slides)
    long_img = Image.new("RGB", (slide_w, total_h), (255, 255, 255))

    # Try to use existing preview images first
    preview_dir = os.path.join(os.path.dirname(pptx_path), "preview")
    if not os.path.exists(preview_dir):
        # Generate soft previews
        script_dir = os.path.dirname(os.path.abspath(__file__))
        render_script = os.path.join(script_dir, "render_slides.py")
        if os.path.exists(render_script):
            preview_dir = os.path.join(os.path.dirname(pptx_path), "_temp_preview")
            try:
                subprocess.run([sys.executable, render_script, pptx_path, preview_dir],
                              capture_output=True, timeout=180)
            except subprocess.TimeoutExpired:
                print("Warning: soft preview timed out during long image build")

    if os.path.exists(preview_dir):
        files = sorted([f for f in os.listdir(preview_dir) if f.endswith(".png")])
        for i, fname in enumerate(files):
            img_path = os.path.join(preview_dir, fname)
            try:
                img = Image.open(img_path)
                img = img.resize((slide_w, slide_h), Image.LANCZOS)
                long_img.paste(img, (0, i * slide_h))
            except Exception as e:
                print(f"Warning: could not load {img_path}: {e}")

    long_img.save(output_path)
    print(f"Long image exported: {output_path}")
    return True


def export_html(pptx_path, output_path):
    """Export PPTX to a simple HTML slideshow using reveal.js style."""
    prs = Presentation(pptx_path)
    slides_html = []
    bg_color = "#0B1F3A"
    text_color = "#FFFFFF"

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
        content = "<br>".join(texts)
        slides_html.append(f'<section class="slide" style="background:{bg_color};color:{text_color}">'
                           f'<div class="slide-content">{content}</div></section>')

    html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>PPTX Export</title>
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{ font-family: "Microsoft YaHei", sans-serif; background:#111; }}
.slideshow {{ width:100vw; height:100vh; overflow:hidden; position:relative; }}
.slide {{ width:100%; height:100%; display:flex; align-items:center; justify-content:center;
          font-size:2rem; text-align:center; padding:4rem; position:absolute; top:0; left:0;
          opacity:0; transition: opacity 0.5s; }}
.slide.active {{ opacity:1; }}
.nav {{ position:fixed; bottom:20px; right:20px; z-index:100; }}
.nav button {{ padding:10px 20px; margin:0 5px; cursor:pointer; background:rgba(255,255,255,0.2);
               color:#fff; border:none; border-radius:4px; }}
.slide-content {{ max-width:900px; line-height:1.6; }}
</style>
</head>
<body>
<div class="slideshow" id="slideshow">
{chr(10).join(slides_html)}
</div>
<div class="nav">
<button onclick="prevSlide()">← Prev</button>
<button onclick="nextSlide()">Next →</button>
<span id="indicator" style="color:#fff;margin-left:10px;">1 / {len(slides_html)}</span>
</div>
<script>
let current = 0;
const slides = document.querySelectorAll('.slide');
function showSlide(n) {{
  slides.forEach(s => s.classList.remove('active'));
  slides[n].classList.add('active');
  document.getElementById('indicator').textContent = (n+1) + ' / ' + slides.length;
}}
function nextSlide() {{ current = (current + 1) % slides.length; showSlide(current); }}
function prevSlide() {{ current = (current - 1 + slides.length) % slides.length; showSlide(current); }}
document.addEventListener('keydown', e => {{ if(e.key==='ArrowRight') nextSlide(); if(e.key==='ArrowLeft') prevSlide(); }});
showSlide(0);
</script>
</body>
</html>"""

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"HTML exported: {output_path}")
    return True


def main():
    parser = argparse.ArgumentParser(description="Export PPTX to various formats")
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument("--format", choices=["pdf", "long_image", "html"], required=True)
    parser.add_argument("--output", "-o", required=True, help="Output file path")
    args = parser.parse_args()

    if args.format == "pdf":
        export_pdf(args.pptx, args.output)
    elif args.format == "long_image":
        export_long_image(args.pptx, args.output)
    elif args.format == "html":
        export_html(args.pptx, args.output)


if __name__ == "__main__":
    main()
