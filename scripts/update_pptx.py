#!/usr/bin/env python3
"""
Apply JSON patch updates to an existing PPTX while preserving formatting.

Usage:
    # Update text on slide 3, shape "TextBox 1"
    python scripts/update_pptx.py input.pptx --slide 3 --shape "TextBox 1" --text "New Text" -o output.pptx

    # Update chart data on slide 4
    python scripts/update_pptx.py input.pptx --slide 4 --chart-patch chart_patch.json -o output.pptx

    # Batch update from JSON patch file
    python scripts/update_pptx.py input.pptx --patch batch_patch.json -o output.pptx
"""

import argparse
import json
import os
import sys
import copy

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def _find_shape_by_name(slide, name):
    """Find a shape by name on a slide."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _find_shape_by_text(slide, text_substring):
    """Find a shape containing specific text."""
    for shape in slide.shapes:
        if shape.has_text_frame and text_substring in shape.text_frame.text:
            return shape
    return None


def _replace_shape_text(shape, new_text):
    """Replace text in a shape while preserving formatting."""
    if not shape.has_text_frame:
        return False

    tf = shape.text_frame
    if not tf.paragraphs:
        return False

    # Strategy: keep the first paragraph's first run formatting as the base style,
    # clear all text, then set new text on the first run.
    first_para = tf.paragraphs[0]

    if first_para.runs:
        # Save formatting from the first run
        first_run = first_para.runs[0]
        saved_font = {
            "name": first_run.font.name,
            "size": first_run.font.size,
            "bold": first_run.font.bold,
            "italic": first_run.font.italic,
            "color": None,
        }
        try:
            saved_font["color"] = first_run.font.color.rgb
        except Exception:
            pass

        # Clear all runs except the first
        for run in first_para.runs[1:]:
            run._r.getparent().remove(run._r)

        # Set new text on first run
        first_run.text = new_text

        # Re-apply formatting (some may be lost during text assignment)
        if saved_font["name"]:
            first_run.font.name = saved_font["name"]
        if saved_font["size"]:
            first_run.font.size = saved_font["size"]
        if saved_font["bold"] is not None:
            first_run.font.bold = saved_font["bold"]
        if saved_font["italic"] is not None:
            first_run.font.italic = saved_font["italic"]
        if saved_font["color"]:
            try:
                first_run.font.color.rgb = saved_font["color"]
            except Exception:
                pass
    else:
        # No runs exist, just set text on the paragraph
        first_para.text = new_text

    # Clear other paragraphs if new_text is single-line
    if "\n" not in new_text and len(tf.paragraphs) > 1:
        for para in tf.paragraphs[1:]:
            para.text = ""

    return True


def _update_chart_data(shape, chart_patch):
    """Update chart data while preserving all formatting."""
    if shape.shape_type != MSO_SHAPE_TYPE.CHART:
        return False, "Shape is not a chart"

    try:
        chart = shape.chart
        from pptx.chart.data import ChartData

        cd = ChartData()

        # Handle series-based updates
        if "categories" in chart_patch:
            cd.categories = chart_patch["categories"]
        else:
            # Try to preserve existing categories
            try:
                plot = chart.plots[0]
                if hasattr(plot, 'categories') and plot.categories:
                    cd.categories = [cat.label if hasattr(cat, 'label') else str(cat) for cat in plot.categories]
            except Exception:
                pass

        if "series" in chart_patch:
            for s in chart_patch["series"]:
                cd.add_series(s.get("name", "Series"), s.get("values", []))
        elif "values" in chart_patch:
            # Single series update
            series_name = chart_patch.get("series_name", "Data")
            cd.add_series(series_name, chart_patch["values"])
        else:
            return False, "No chart data provided in patch"

        chart.replace_data(cd)
        return True, "Chart data updated"

    except Exception as e:
        return False, f"Chart update failed: {e}"


def _update_table_data(shape, table_patch):
    """Update table cell values."""
    if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        return False, "Shape is not a table"

    try:
        table = shape.table
        rows_patch = table_patch.get("rows", [])

        for row_idx, row_data in enumerate(rows_patch):
            if row_idx >= len(table.rows):
                break
            row = table.rows[row_idx]
            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= len(row.cells):
                    break
                row.cells[col_idx].text = str(cell_text)

        return True, f"Table updated ({len(rows_patch)} rows)"
    except Exception as e:
        return False, f"Table update failed: {e}"


def _update_shape_fill(shape, fill_color):
    """Update shape fill color."""
    try:
        shape.fill.solid()
        hex_str = fill_color.lstrip("#")
        shape.fill.fore_color.rgb = RGBColor(
            int(hex_str[0:2], 16),
            int(hex_str[2:4], 16),
            int(hex_str[4:6], 16)
        )
        return True, "Fill updated"
    except Exception as e:
        return False, f"Fill update failed: {e}"


def apply_patch(prs, patch):
    """Apply a single patch to a presentation."""
    results = []

    slide_number = patch.get("slide")
    if slide_number is None:
        return [{"success": False, "error": "slide number is required"}]

    if slide_number < 1 or slide_number > len(prs.slides):
        return [{"success": False, "error": f"Slide {slide_number} not found (total: {len(prs.slides)})"}]

    slide = prs.slides[slide_number - 1]

    # Find target shape
    shape = None
    shape_name = patch.get("shape")
    shape_text = patch.get("shape_text")

    if shape_name:
        shape = _find_shape_by_name(slide, shape_name)
        if shape is None:
            return [{"success": False, "error": f"Shape '{shape_name}' not found on slide {slide_number}"}]
    elif shape_text:
        shape = _find_shape_by_text(slide, shape_text)
        if shape is None:
            return [{"success": False, "error": f"No shape containing '{shape_text}' on slide {slide_number}"}]
    else:
        # Apply to all text shapes on the slide
        shapes = [s for s in slide.shapes if s.has_text_frame]
        if not shapes:
            return [{"success": False, "error": f"No text shapes on slide {slide_number}"}]
        # Apply to all matching shapes
        for s in shapes:
            result = apply_patch_to_shape(s, patch)
            result["slide"] = slide_number
            results.append(result)
        return results

    result = apply_patch_to_shape(shape, patch)
    result["slide"] = slide_number
    results.append(result)
    return results


def apply_patch_to_shape(shape, patch):
    """Apply a patch to a specific shape."""
    # Text update
    if "text" in patch:
        success = _replace_shape_text(shape, patch["text"])
        if success:
            return {"success": True, "shape": shape.name, "action": "text_updated"}
        else:
            return {"success": False, "shape": shape.name, "error": "Text replacement failed"}

    # Chart data update
    if "chart_data" in patch:
        success, msg = _update_chart_data(shape, patch["chart_data"])
        return {"success": success, "shape": shape.name, "action": "chart_updated", "message": msg}

    # Table data update
    if "table_data" in patch:
        success, msg = _update_table_data(shape, patch["table_data"])
        return {"success": success, "shape": shape.name, "action": "table_updated", "message": msg}

    # Fill color update
    if "fill" in patch:
        success, msg = _update_shape_fill(shape, patch["fill"])
        return {"success": success, "shape": shape.name, "action": "fill_updated", "message": msg}

    return {"success": False, "shape": shape.name, "error": "No valid patch action found"}


def main():
    parser = argparse.ArgumentParser(description="Apply patches to an existing PPTX")
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument("--output", "-o", required=True, help="Output PPTX file")
    parser.add_argument("--slide", type=int, help="Slide number (1-based)")
    parser.add_argument("--shape", help="Shape name to update")
    parser.add_argument("--shape-text", help="Find shape by text content")
    parser.add_argument("--text", help="New text content")
    parser.add_argument("--chart-patch", help="JSON file with chart data patch")
    parser.add_argument("--table-patch", help="JSON file with table data patch")
    parser.add_argument("--fill", help="New fill color (hex #RRGGBB)")
    parser.add_argument("--patch", help="JSON patch file (batch operations)")
    parser.add_argument("--dry-run", action="store_true", help="Preview changes without writing")
    args = parser.parse_args()

    if not HAS_PPTX:
        print("Error: python-pptx is required", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(args.pptx)
    all_results = []

    # Single patch mode
    if args.slide:
        patch = {"slide": args.slide}
        if args.shape:
            patch["shape"] = args.shape
        elif args.shape_text:
            patch["shape_text"] = args.shape_text

        if args.text:
            patch["text"] = args.text
        elif args.chart_patch:
            with open(args.chart_patch, "r", encoding="utf-8") as f:
                patch["chart_data"] = json.load(f)
        elif args.table_patch:
            with open(args.table_patch, "r", encoding="utf-8") as f:
                patch["table_data"] = json.load(f)
        elif args.fill:
            patch["fill"] = args.fill
        else:
            print("Error: No patch action specified (--text, --chart-patch, --table-patch, --fill)", file=sys.stderr)
            sys.exit(1)

        results = apply_patch(prs, patch)
        all_results.extend(results)

    # Batch patch mode
    elif args.patch:
        with open(args.patch, "r", encoding="utf-8") as f:
            patches = json.load(f)

        if not isinstance(patches, list):
            patches = [patches]

        for patch in patches:
            results = apply_patch(prs, patch)
            all_results.extend(results)

    else:
        print("Error: Either --slide or --patch is required", file=sys.stderr)
        sys.exit(1)

    # Report
    success_count = sum(1 for r in all_results if r.get("success"))
    fail_count = len(all_results) - success_count

    print(f"Results: {success_count} succeeded, {fail_count} failed")
    for r in all_results:
        status = "OK" if r.get("success") else "FAIL"
        msg = r.get("message", r.get("error", ""))
        print(f"  [{status}] Slide {r.get('slide', '?')}, Shape '{r.get('shape', '?')}': {msg}")

    if not args.dry_run:
        prs.save(args.output)
        print(f"Saved to {args.output}")
    else:
        print("(dry run, no file written)")


if __name__ == "__main__":
    main()
