#!/usr/bin/env python3
"""
Install a custom PPTX template into the template market.

Usage:
    python scripts/install_template.py my_template.pptx --name consulting --category business
"""

import argparse
import json
import os
import shutil
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

MARKET_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "templates", "market")


def analyze_template(pptx_path):
    """Analyze a template PPTX and extract its structure."""
    if not HAS_PPTX:
        raise ImportError("python-pptx is required")

    prs = Presentation(pptx_path)

    info = {
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "layouts": [],
        "slides": []
    }

    # Extract slide layouts
    for layout in prs.slide_layouts:
        layout_info = {
            "name": layout.name,
            "placeholders": []
        }
        for ph in layout.placeholders:
            layout_info["placeholders"].append({
                "name": ph.name,
                "type": str(ph.placeholder_format.type).split(".")[-1] if ph.placeholder_format else "unknown",
                "left": ph.left,
                "top": ph.top,
                "width": ph.width,
                "height": ph.height
            })
        info["layouts"].append(layout_info)

    # Extract template slides (if any)
    for slide in prs.slides:
        slide_info = {
            "shapes": []
        }
        for shape in slide.shapes:
            shape_info = {"name": shape.name, "type": str(shape.shape_type).split(".")[-1].lower()}
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    shape_info["text"] = text
            try:
                if shape.fill.type is not None and shape.fill.fore_color:
                    rgb = shape.fill.fore_color.rgb
                    shape_info["fill"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
            except Exception:
                pass
            slide_info["shapes"].append(shape_info)
        info["slides"].append(slide_info)

    return info


def install_template(pptx_path, name, category="general", description=""):
    """Install a template into the market."""
    os.makedirs(MARKET_DIR, exist_ok=True)

    template_dir = os.path.join(MARKET_DIR, name)
    if os.path.exists(template_dir):
        print(f"Warning: Template '{name}' already exists. Overwriting.")

    os.makedirs(template_dir, exist_ok=True)

    # Copy PPTX
    dest_pptx = os.path.join(template_dir, "template.pptx")
    shutil.copy2(pptx_path, dest_pptx)

    # Analyze and write config
    info = analyze_template(pptx_path)
    config = {
        "name": name,
        "category": category,
        "description": description,
        "source": os.path.basename(pptx_path),
        "slide_width": info["slide_width"],
        "slide_height": info["slide_height"],
        "layouts": [l["name"] for l in info["layouts"]],
        "placeholders": {}
    }

    # Extract placeholder mappings
    for layout in info["layouts"]:
        for ph in layout["placeholders"]:
            text = ph["name"].lower()
            if any(k in text for k in ("title", "标题")):
                config["placeholders"][ph["name"]] = "title"
            elif any(k in text for k in ("content", "body", "text", "内容")):
                config["placeholders"][ph["name"]] = "content"
            elif any(k in text for k in ("subtitle", "副标题")):
                config["placeholders"][ph["name"]] = "subtitle"
            elif any(k in text for k in ("date", "日期")):
                config["placeholders"][ph["name"]] = "date"

    config_path = os.path.join(template_dir, "config.json")
    with open(config_path, "w", encoding="utf-8") as f:
        json.dump(config, f, ensure_ascii=False, indent=2)

    # Write full analysis
    analysis_path = os.path.join(template_dir, "analysis.json")
    with open(analysis_path, "w", encoding="utf-8") as f:
        json.dump(info, f, ensure_ascii=False, indent=2)

    print(f"Template '{name}' installed to {template_dir}")
    print(f"  Layouts: {len(info['layouts'])}")
    print(f"  Placeholders: {len(config['placeholders'])}")
    print(f"  Config: {config_path}")
    print(f"\nUse with: python scripts/build_pptx.py deck.json out.pptx --template {dest_pptx}")


def list_templates():
    """List all installed templates."""
    if not os.path.exists(MARKET_DIR):
        print("No templates installed.")
        return

    templates = []
    for name in os.listdir(MARKET_DIR):
        config_path = os.path.join(MARKET_DIR, name, "config.json")
        if os.path.exists(config_path):
            with open(config_path, "r", encoding="utf-8") as f:
                config = json.load(f)
            templates.append(config)

    if not templates:
        print("No templates installed.")
        return

    print(f"Installed templates ({len(templates)}):")
    print()
    for t in sorted(templates, key=lambda x: x.get("category", "") + x["name"]):
        cat = t.get("category", "general")
        desc = t.get("description", "")
        layouts = len(t.get("layouts", []))
        print(f"  {t['name']:20}  [{cat:12}]  {layouts} layouts  {desc}")


def main():
    parser = argparse.ArgumentParser(description="Install or list PPTX templates")
    subparsers = parser.add_subparsers(dest="command", help="Command")

    install_parser = subparsers.add_parser("install", help="Install a template")
    install_parser.add_argument("pptx", help="Source PPTX file")
    install_parser.add_argument("--name", "-n", required=True, help="Template name")
    install_parser.add_argument("--category", "-c", default="general", help="Category")
    install_parser.add_argument("--description", "-d", default="", help="Description")

    list_parser = subparsers.add_parser("list", help="List installed templates")

    args = parser.parse_args()

    if args.command == "install":
        install_template(args.pptx, args.name, args.category, args.description)
    elif args.command == "list":
        list_templates()
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
