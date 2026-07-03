#!/usr/bin/env python3
"""color_guard.py - Contrast & legibility guardian for PPTX slides."""
import sys, os, io
from typing import List, Dict, Tuple, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# ---------------------------------------------------------------------------
# Color math
# ---------------------------------------------------------------------------

def _hex_to_rgb(hex_str: str) -> Tuple[int, int, int]:
    h = hex_str.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)

def _rgb_to_hex(rgb) -> str:
    if isinstance(rgb, RGBColor):
        return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

def _relative_luminance(rgb) -> float:
    def ch(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r = rgb[0] if isinstance(rgb, (tuple, list)) else rgb[0]
    g = rgb[1] if isinstance(rgb, (tuple, list)) else rgb[1]
    b = rgb[2] if isinstance(rgb, (tuple, list)) else rgb[2]
    return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(b)

def contrast_ratio(c1, c2) -> float:
    l1, l2 = _relative_luminance(c1), _relative_luminance(c2)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)

def brightness(rgb) -> float:
    r = rgb[0] if isinstance(rgb, (tuple, list)) else rgb[0]
    g = rgb[1] if isinstance(rgb, (tuple, list)) else rgb[1]
    b = rgb[2] if isinstance(rgb, (tuple, list)) else rgb[2]
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255.0

def recommend_text_color(bg_rgb, theme_text_dark=(31, 41, 55), theme_text_light=(248, 250, 252)) -> Tuple[Tuple[int, int, int], float]:
    dark_ratio = contrast_ratio(bg_rgb, theme_text_dark)
    light_ratio = contrast_ratio(bg_rgb, theme_text_light)
    if light_ratio > dark_ratio:
        return theme_text_light, light_ratio
    return theme_text_dark, dark_ratio

# High-contrast cap for dark themes: when background is near-black,
# pure white text can cause halation / eye strain.
MAX_CONTRAST_DARK = 16.0  # WCAG maximum is 21:1; cap a bit below that.

# Background-background glare threshold for dark slides.
# When slide bg is very dark and a shape (chart/table) is light,
# the high contrast can cause eye strain.  We flag it as INFO.
BG_BG_DARK_THRESHOLD = 55   # max(R,G,B) < 55 => very dark (catches #0F172A etc.)
BG_BG_LIGHT_THRESHOLD = 200 # min(R,G,B) > 200 => very light
BG_BG_CHANNEL_DIFF = 150    # at least 2 channels differ by > 150

def _is_very_dark(rgb):
    return max(rgb) < BG_BG_DARK_THRESHOLD

def _is_very_light(rgb):
    return min(rgb) > BG_BG_LIGHT_THRESHOLD

def _check_bg_bg_contrast(slide_bg_rgb, shape_bg_rgb, slide_idx, shape_name, warnings):
    """Warn when a light shape sits on a very dark slide (glare)."""
    if not _is_very_dark(slide_bg_rgb):
        return
    if not _is_very_light(shape_bg_rgb):
        return
    # Count channels with large difference
    diffs = [abs(shape_bg_rgb[i] - slide_bg_rgb[i]) for i in range(3)]
    if sum(1 for d in diffs if d > BG_BG_CHANNEL_DIFF) < 2:
        return
    warnings.append({
        "slide": slide_idx,
        "element_type": "bg_bg_contrast",
        "shape_text_preview": shape_name,
        "text_color": _rgb_to_hex(shape_bg_rgb),
        "bg_color": _rgb_to_hex(slide_bg_rgb),
        "bg_source": "slide_vs_shape",
        "contrast_ratio": None,
        "severity": "INFO",
        "note": f"Light {shape_name} on very dark slide may cause glare. Consider darkening the shape background.",
        "recommendation": {}
    })

def _check_high_contrast_dark(text_rgb, bg_rgb, slide_idx, shape_text, element_type, warnings):
    """Warn if contrast is excessive on very dark backgrounds (dark-theme halation)."""
    bg_lum = _relative_luminance(bg_rgb)
    text_lum = _relative_luminance(text_rgb)
    if bg_lum > 0.15 or text_lum < 0.7:
        return  # Only check near-black bg + bright text
    ratio = contrast_ratio(text_rgb, bg_rgb)
    if ratio > MAX_CONTRAST_DARK:
        # Suggest dimmed white to bring contrast into comfortable range
        dimmed = (226, 232, 240)  # #E2E8F0
        new_ratio = contrast_ratio(dimmed, bg_rgb)
        warnings.append({
            "slide": slide_idx,
            "element_type": element_type,
            "shape_text_preview": shape_text[:60],
            "text_color": _rgb_to_hex(text_rgb),
            "bg_color": _rgb_to_hex(bg_rgb),
            "bg_source": "slide_bg",
            "contrast_ratio": round(ratio, 2),
            "severity": "INFO",
            "note": f"Excessive contrast on dark bg ({round(ratio,1)}:1). Consider dimming text to {_rgb_to_hex(dimmed)} ({round(new_ratio,1)}:1).",
            "recommendation": {
                "suggested_text_color": _rgb_to_hex(dimmed),
                "projected_contrast": round(new_ratio, 2)
            }
        })

# ---------------------------------------------------------------------------
# Background sampling
# ---------------------------------------------------------------------------

def _shape_fill_color(shape) -> Optional[Tuple[int, int, int]]:
    try:
        if shape.fill.type is not None and shape.fill.fore_color and shape.fill.fore_color.rgb:
            c = shape.fill.fore_color.rgb
            return (c[0], c[1], c[2])
    except Exception:
        pass
    return None

def _dominant_color_from_image_bytes(img_bytes: bytes, crop_box=None):
    if not HAS_PIL:
        return None
    try:
        img = Image.open(io.BytesIO(img_bytes))
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        if crop_box:
            w, h = img.size
            left = int(w * crop_box[0])
            top = int(h * crop_box[1])
            right = int(w * crop_box[2])
            bottom = int(h * crop_box[3])
            img = img.crop((left, top, right, bottom))
        img_small = img.resize((1, 1), Image.Resampling.LANCZOS)
        r, g, b = img_small.getpixel((0, 0))
        return r, g, b
    except Exception:
        return None

def _sample_picture_region(slide, picture_shape, text_shape):
    if not HAS_PIL:
        return None
    try:
        img_bytes = picture_shape.image.blob
    except Exception:
        return None
    px, py, pw, ph = picture_shape.left, picture_shape.top, picture_shape.width, picture_shape.height
    tx, ty, tw, th = text_shape.left, text_shape.top, text_shape.width, text_shape.height
    left = max(0.0, (tx - px) / pw) if pw else 0.0
    top = max(0.0, (ty - py) / ph) if ph else 0.0
    right = min(1.0, (tx + tw - px) / pw) if pw else 1.0
    bottom = min(1.0, (ty + th - py) / ph) if ph else 1.0
    if right <= left or bottom <= top:
        return None
    return _dominant_color_from_image_bytes(img_bytes, (left, top, right, bottom))

def sample_local_background(slide, text_shape, fallback_rgb=(255, 255, 255)):
    tx, ty, tw, th = text_shape.left, text_shape.top, text_shape.width, text_shape.height
    cx, cy = tx + tw // 2, ty + th // 2
    layer_stack = []
    candidates = []
    for shape in slide.shapes:
        if shape == text_shape:
            continue
        sx, sy, sw, sh = shape.left, shape.top, shape.width, shape.height
        if sx <= cx <= sx + sw and sy <= cy <= sy + sh:
            candidates.append(shape)
    for shape in reversed(candidates):
        if shape.shape_type == 13:  # PICTURE
            color = _sample_picture_region(slide, shape, text_shape)
            if color:
                layer_stack.append({"type": "picture", "name": getattr(shape, 'name', ''), "color": _rgb_to_hex(color)})
                return color, "picture_crop", layer_stack
            layer_stack.append({"type": "picture", "name": getattr(shape, 'name', ''), "color": None})
            continue
        fill_color = _shape_fill_color(shape)
        if fill_color:
            layer_stack.append({"type": "solid", "name": getattr(shape, 'name', ''), "color": _rgb_to_hex(fill_color)})
            return fill_color, "solid_shape", layer_stack
        layer_stack.append({"type": str(shape.shape_type), "name": getattr(shape, 'name', '')})
    return fallback_rgb, "slide_bg", layer_stack

# ---------------------------------------------------------------------------
# Table audit
# ---------------------------------------------------------------------------

def _audit_table(shape, slide_idx, warnings, slide_bg=(255,255,255)):
    try:
        table = shape.table
    except Exception:
        return
    # Check overall table background vs slide background
    table_bg = (255, 255, 255)
    try:
        if shape.fill.type is not None and shape.fill.fore_color and shape.fill.fore_color.rgb:
            c = shape.fill.fore_color.rgb
            table_bg = (c[0], c[1], c[2])
    except Exception:
        pass
    _check_bg_bg_contrast(slide_bg, table_bg, slide_idx + 1, "table", warnings)
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            cell_text = cell.text.strip()
            if not cell_text:
                continue
            cell_bg = (255, 255, 255)
            try:
                if cell.fill.type is not None and cell.fill.fore_color and cell.fill.fore_color.rgb:
                    c = cell.fill.fore_color.rgb
                    cell_bg = (c[0], c[1], c[2])
            except Exception:
                pass
            text_rgb = (31, 41, 55)
            try:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            c = run.font.color.rgb
                            text_rgb = (c[0], c[1], c[2])
                        except AttributeError:
                            pass
                        break
                    break
            except Exception:
                pass
            ratio = contrast_ratio(text_rgb, cell_bg)
            if ratio < 4.5:
                rec_color, rec_ratio = recommend_text_color(cell_bg)
                warnings.append({
                    "slide": slide_idx + 1,
                    "element_type": "table_cell",
                    "shape_text_preview": cell_text[:60],
                    "text_color": _rgb_to_hex(text_rgb),
                    "bg_color": _rgb_to_hex(cell_bg),
                    "bg_source": "table_cell_fill",
                    "contrast_ratio": round(ratio, 2),
                    "severity": "CRITICAL" if ratio < 2.0 else "WARNING",
                    "recommendation": {
                        "suggested_text_color": _rgb_to_hex(rec_color),
                        "projected_contrast": round(rec_ratio, 2)
                    }
                })
            _check_high_contrast_dark(text_rgb, cell_bg, slide_idx + 1, cell_text, "table_cell", warnings)

# ---------------------------------------------------------------------------
# Chart audit
# ---------------------------------------------------------------------------

def _audit_chart(shape, slide_idx, warnings, slide_bg=(255, 255, 255)):
    try:
        chart = shape.chart
    except Exception:
        return
    # Read explicit chart background from XML if present; otherwise default white.
    chart_bg = (255, 255, 255)
    try:
        import xml.etree.ElementTree as ET
        c_ns = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        cs_spPr = chart._element.find('{' + c_ns + '}spPr')
        if cs_spPr is not None:
            sf = cs_spPr.find('{' + a_ns + '}solidFill')
            if sf is not None:
                sc = sf.find('{' + a_ns + '}srgbClr')
                if sc is not None:
                    v = sc.get('val')
                    if v:
                        chart_bg = (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
        if chart_bg == (255, 255, 255):
            plot_area = chart._element.find('.//{' + c_ns + '}plotArea')
            if plot_area is not None:
                pa_spPr = plot_area.find('{' + c_ns + '}spPr')
                if pa_spPr is not None:
                    sf = pa_spPr.find('{' + a_ns + '}solidFill')
                    if sf is not None:
                        sc = sf.find('{' + a_ns + '}srgbClr')
                        if sc is not None:
                            v = sc.get('val')
                            if v:
                                chart_bg = (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
    except Exception:
        pass
    _check_bg_bg_contrast(slide_bg, chart_bg, slide_idx + 1, "chart", warnings)
    elements = []
    try:
        if chart.category_axis and chart.category_axis.tick_labels:
            try:
                c = chart.category_axis.tick_labels.font.color.rgb
                elements.append(("cat_axis", (c[0], c[1], c[2])))
            except AttributeError:
                elements.append(("cat_axis", (31, 41, 55)))
    except ValueError:
        pass  # Pie/doughnut has no category axis
    try:
        if chart.value_axis and chart.value_axis.tick_labels:
            try:
                c = chart.value_axis.tick_labels.font.color.rgb
                elements.append(("val_axis", (c[0], c[1], c[2])))
            except AttributeError:
                elements.append(("val_axis", (31, 41, 55)))
    except ValueError:
        pass  # Pie/doughnut has no value axis
    if chart.has_legend:
        try:
            c = chart.legend.font.color.rgb
            elements.append(("legend", (c[0], c[1], c[2])))
        except AttributeError:
            elements.append(("legend", (31, 41, 55)))
    if chart.has_title and chart.chart_title:
        try:
            for para in chart.chart_title.text_frame.paragraphs:
                for run in para.runs:
                    c = run.font.color.rgb
                    elements.append(("title", (c[0], c[1], c[2])))
                    break
                break
        except AttributeError:
            elements.append(("title", (31, 41, 55)))
    for elem_name, text_rgb in elements:
        ratio = contrast_ratio(text_rgb, chart_bg)
        if ratio < 4.5:
            rec_color, rec_ratio = recommend_text_color(chart_bg)
            warnings.append({
                "slide": slide_idx + 1,
                "element_type": f"chart_{elem_name}",
                "shape_text_preview": f"chart {elem_name}",
                "text_color": _rgb_to_hex(text_rgb),
                "bg_color": _rgb_to_hex(chart_bg),
                "bg_source": "chart_plot_area",
                "contrast_ratio": round(ratio, 2),
                "severity": "CRITICAL" if ratio < 2.0 else "WARNING",
                "recommendation": {
                    "suggested_text_color": _rgb_to_hex(rec_color),
                    "projected_contrast": round(rec_ratio, 2)
                }
            })

# ---------------------------------------------------------------------------
# Shape-with-fill audit
# ---------------------------------------------------------------------------

def _audit_shape_fill_text(shape, slide_idx, slide_bg, warnings):
    if not shape.has_text_frame or not shape.text_frame.text.strip():
        return
    fill_color = _shape_fill_color(shape)
    if not fill_color:
        return
    if fill_color == slide_bg:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                c = run.font.color.rgb
                text_rgb = (c[0], c[1], c[2])
            except AttributeError:
                continue
            if not run.text.strip():
                continue
            ratio = contrast_ratio(text_rgb, fill_color)
            if ratio < 4.5:
                rec_color, rec_ratio = recommend_text_color(fill_color)
                warnings.append({
                    "slide": slide_idx + 1,
                    "element_type": "shape_fill_text",
                    "shape_text_preview": run.text.strip()[:60],
                    "text_color": _rgb_to_hex(text_rgb),
                    "bg_color": _rgb_to_hex(fill_color),
                    "bg_source": "shape_solid_fill",
                    "contrast_ratio": round(ratio, 2),
                    "severity": "CRITICAL" if ratio < 2.0 else "WARNING",
                    "recommendation": {
                        "suggested_text_color": _rgb_to_hex(rec_color),
                        "projected_contrast": round(rec_ratio, 2)
                    }
                })
            _check_high_contrast_dark(text_rgb, fill_color, slide_idx + 1, run.text, "shape_fill_text", warnings)

# ---------------------------------------------------------------------------
# Per-slide analysis
# ---------------------------------------------------------------------------

def analyze_slide(prs, slide_idx, theme=None, fallback_bg=(255, 255, 255)):
    slide = prs.slides[slide_idx]
    warnings = []
    slide_area = prs.slide_width * prs.slide_height
    slide_bg = fallback_bg
    for shape in slide.shapes:
        if shape.width * shape.height > slide_area * 0.8:
            try:
                if shape.fill.type is not None and shape.fill.fore_color and shape.fill.fore_color.rgb:
                    c = shape.fill.fore_color.rgb
                    slide_bg = (c[0], c[1], c[2])
            except Exception:
                pass
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            _audit_table(shape, slide_idx, warnings, slide_bg)
            continue
        if shape.shape_type == 3:  # CHART
            _audit_chart(shape, slide_idx, warnings, slide_bg)
            continue
        _audit_shape_fill_text(shape, slide_idx, slide_bg, warnings)
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        text_runs = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    c = run.font.color.rgb
                    text_runs.append((run.text, (c[0], c[1], c[2]), run.font.size.pt if run.font.size else None))
                except AttributeError:
                    fallback_text = _hex_to_rgb(theme["text"]) if theme else (31, 41, 55)
                    text_runs.append((run.text, fallback_text, run.font.size.pt if run.font.size else None))
        if not text_runs:
            continue
        # For shapes with their own solid fill, use the fill as background
        # rather than sampling the slide (avoids false positives on colored
        # shapes like tree nodes or process steps).
        own_fill = _shape_fill_color(shape)
        if own_fill is not None and own_fill != slide_bg:
            local_bg, bg_source, layers = own_fill, "shape_solid_fill", []
        else:
            local_bg, bg_source, layers = sample_local_background(slide, shape, slide_bg)
        for text, text_rgb, size_pt in text_runs:
            if not text.strip():
                continue
            ratio = contrast_ratio(text_rgb, local_bg)
            if ratio < 4.5:
                rec_color, rec_ratio = recommend_text_color(local_bg)
                warnings.append({
                    "slide": slide_idx + 1,
                    "element_type": "textbox",
                    "shape_text_preview": text.strip()[:60],
                    "text_color": _rgb_to_hex(text_rgb),
                    "bg_color": _rgb_to_hex(local_bg),
                    "bg_source": bg_source,
                    "contrast_ratio": round(ratio, 2),
                    "severity": "CRITICAL" if ratio < 2.0 else "WARNING",
                    "font_size_pt": size_pt,
                    "recommendation": {
                        "suggested_text_color": _rgb_to_hex(rec_color),
                        "projected_contrast": round(rec_ratio, 2)
                    },
                    "layer_stack": layers
                })
            _check_high_contrast_dark(text_rgb, local_bg, slide_idx + 1, text, "textbox", warnings)
    return warnings

# ---------------------------------------------------------------------------
# Full presentation analysis
# ---------------------------------------------------------------------------

class ContrastReport:
    def __init__(self, warnings, theme_name=""):
        self.warnings = warnings
        self.theme_name = theme_name
        self.critical_count = sum(1 for w in warnings if w.get("severity") == "CRITICAL")
        self.warning_count = sum(1 for w in warnings if w.get("severity") == "WARNING")

    def __bool__(self):
        return len(self.warnings) == 0

    def json(self) -> str:
        import json
        return json.dumps({
            "theme": self.theme_name,
            "summary": {
                "total_issues": len(self.warnings),
                "critical": self.critical_count,
                "warning": self.warning_count,
                "pass": len(self.warnings) == 0
            },
            "issues": self.warnings
        }, ensure_ascii=False, indent=2)

    def markdown(self) -> str:
        lines = [f"## Contrast Audit: {self.theme_name or 'Untitled'}", ""]
        lines.append(f"**Total issues:** {len(self.warnings)}  ")
        lines.append(f"Critical: {self.critical_count}  Warning: {self.warning_count}  ")
        lines.append("")
        if not self.warnings:
            lines.append("All text shapes meet WCAG AA contrast (>=4.5:1).")
            return "\n".join(lines)
        for w in self.warnings:
            sev = w.get("severity", "WARNING")
            emoji = {"CRITICAL": "CRIT", "WARNING": "WARN", "INFO": "INFO"}.get(sev, "WARN")
            elem = w.get("element_type", "textbox")
            ratio_str = f"ratio {w['contrast_ratio']}:1" if w.get('contrast_ratio') else ""
            lines.append(f"{emoji} Slide {w['slide']} [{elem}] {ratio_str}")
            if elem == "bg_bg_contrast":
                lines.append(f"   {w.get('note', '')}")
            else:
                lines.append(f"   Text: {w['text_color']} -> {w['shape_text_preview']}")
                lines.append(f"   Background ({w['bg_source']}): {w['bg_color']}")
                rec = w.get("recommendation", {})
                lines.append(f"   Fix: change text to {rec.get('suggested_text_color', '?')} "
                            f"(projected {rec.get('projected_contrast', '?')}:1)")
            lines.append("")
        return "\n".join(lines)

def analyze_presentation(prs, theme=None):
    all_warnings = []
    fallback_bg = _hex_to_rgb(theme["background"]) if theme and "background" in theme else (255, 255, 255)
    for i in range(len(prs.slides)):
        all_warnings.extend(analyze_slide(prs, i, theme=theme, fallback_bg=fallback_bg))
    return ContrastReport(all_warnings, theme_name=theme.get("name", "") if theme else "")

# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import argparse, json
    parser = argparse.ArgumentParser(description="Contrast audit for PPTX")
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument("--theme-json", help="Optional theme JSON for background/text fallback colors")
    parser.add_argument("--format", choices=["json", "markdown"], default="markdown", help="Output format")
    parser.add_argument("--output", "-o", help="Write report to file")
    args = parser.parse_args()

    theme = None
    if args.theme_json and os.path.exists(args.theme_json):
        with open(args.theme_json, "r", encoding="utf-8") as f:
            theme = json.load(f)

    prs = Presentation(args.pptx)
    report = analyze_presentation(prs, theme=theme)

    out = report.json() if args.format == "json" else report.markdown()
    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(out)
        print(f"Report saved: {args.output}")
    else:
        print(out)
