#!/usr/bin/env python3
"""Extract visual contract (theme colors, fonts, layout structure) from a PPTX template.

Usage:
  python template_analyzer.py template.pptx --output contract.json
"""
import sys, os, json, argparse, collections

from pptx import Presentation
from pptx.util import Pt


def _rgb_to_hex(rgb_color):
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color[0]:02X}{rgb_color[1]:02X}{rgb_color[2]:02X}"
    except Exception:
        return None


def extract_colors(prs):
    """Extract dominant colors from shapes across all template slides."""
    colors = collections.Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            # Text colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            c = run.font.color.rgb
                            if c:
                                colors[_rgb_to_hex((c[0], c[1], c[2]))] += 1
                        except Exception:
                            pass
            # Fill colors
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                try:
                    c = shape.fill.fore_color.rgb
                    if c:
                        hex_c = _rgb_to_hex((c[0], c[1], c[2]))
                        # Weight background shapes more heavily
                        shape_area = shape.width * shape.height
                        slide_area = prs.slide_width * prs.slide_height
                        weight = 3 if shape_area > slide_area * 0.5 else 1
                        colors[hex_c] += weight
                except Exception:
                    pass
            # Line colors
            if hasattr(shape, 'line') and shape.line.fill.type is not None:
                try:
                    c = shape.line.color.rgb
                    if c:
                        colors[_rgb_to_hex((c[0], c[1], c[2]))] += 1
                except Exception:
                    pass
    return colors


def extract_fonts(prs):
    """Extract font usage statistics."""
    fonts = collections.Counter()
    font_sizes = collections.Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts[run.font.name] += 1
                        if run.font.size:
                            font_sizes[int(run.font.size.pt)] += 1
    return fonts, font_sizes


def extract_layouts(prs):
    """Extract layout structure from template slides."""
    layouts = []
    for slide in prs.slides:
        layout_info = {
            "name": slide.name,
            "shapes": [],
        }
        for shape in slide.shapes:
            sinfo = {
                "type": str(shape.shape_type).split('.')[-1] if hasattr(shape, 'shape_type') else "unknown",
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
            # Relative position
            sinfo["rel_left"] = round(shape.left / prs.slide_width, 3)
            sinfo["rel_top"] = round(shape.top / prs.slide_height, 3)
            sinfo["rel_width"] = round(shape.width / prs.slide_width, 3)
            sinfo["rel_height"] = round(shape.height / prs.slide_height, 3)

            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                sinfo["text_preview"] = text[:50] if text else ""
                # Detect placeholder patterns
                if text.startswith("__") and text.endswith("__"):
                    sinfo["role"] = text.strip("_").lower()
                # Detect font
                try:
                    font = shape.text_frame.paragraphs[0].runs[0].font
                    sinfo["font_name"] = font.name
                    if font.size:
                        sinfo["font_size"] = int(font.size.pt)
                    if font.bold:
                        sinfo["bold"] = True
                except Exception:
                    pass
                # Detect color
                try:
                    c = shape.text_frame.paragraphs[0].runs[0].font.color.rgb
                    if c:
                        sinfo["color"] = _rgb_to_hex((c[0], c[1], c[2]))
                except Exception:
                    pass

            if hasattr(shape, 'fill') and shape.fill.type is not None:
                try:
                    c = shape.fill.fore_color.rgb
                    if c:
                        sinfo["fill"] = _rgb_to_hex((c[0], c[1], c[2]))
                except Exception:
                    pass

            layout_info["shapes"].append(sinfo)
        layouts.append(layout_info)
    return layouts


def analyze_template(pptx_path):
    """Analyze a PPTX template and return visual contract."""
    prs = Presentation(pptx_path)

    colors = extract_colors(prs)
    fonts, font_sizes = extract_fonts(prs)
    layouts = extract_layouts(prs)

    # Identify theme colors
    sorted_colors = colors.most_common()
    primary_candidates = [c for c in sorted_colors if c[1] >= 3]

    # Heuristic: most common non-white/non-black is primary
    primary = None
    secondary = None
    bg = None
    text = None
    for hex_c, count in primary_candidates:
        if hex_c in ("#FFFFFF", "#000000", "#F3F4F6", "#FFFFFF00"):
            continue
        if not primary:
            primary = hex_c
        elif not secondary:
            secondary = hex_c
        break

    # Most common large-area color is likely background
    for slide in prs.slides:
        for shape in slide.shapes:
            shape_area = shape.width * shape.height
            slide_area = prs.slide_width * prs.slide_height
            if shape_area > slide_area * 0.5 and hasattr(shape, 'fill'):
                try:
                    c = shape.fill.fore_color.rgb
                    if c:
                        bg = _rgb_to_hex((c[0], c[1], c[2]))
                        break
                except Exception:
                    pass
        if bg:
            break

    contract = {
        "slide_dimensions": {
            "width_emu": prs.slide_width,
            "height_emu": prs.slide_height,
            "width_in": round(prs.slide_width / 914400, 3),
            "height_in": round(prs.slide_height / 914400, 3),
        },
        "theme": {
            "primary": primary,
            "secondary": secondary,
            "background": bg,
            "text": text,
            "all_colors": dict(sorted_colors[:20]),
        },
        "typography": {
            "fonts": dict(fonts.most_common(10)),
            "font_sizes": dict(font_sizes.most_common(10)),
            "title_font": fonts.most_common(1)[0][0] if fonts else None,
        },
        "layouts": layouts,
    }
    return contract


def main():
    parser = argparse.ArgumentParser(description="Analyze PPTX template visual contract")
    parser.add_argument("pptx", help="Input PPTX template")
    parser.add_argument("--output", "-o", help="Output JSON file")
    args = parser.parse_args()

    contract = analyze_template(args.pptx)

    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            json.dump(contract, f, ensure_ascii=False, indent=2)
        print(f"Contract saved to {args.output}")
    else:
        print(json.dumps(contract, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
