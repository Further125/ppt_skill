#!/usr/bin/env python3
"""Font size guard: warn when text is too small to read comfortably.

Similar to color_guard.py but for typography legibility.
"""

from pptx.util import Pt

# WCAG-inspired readable thresholds (PowerPoint context)
DEFAULT_MIN_READABLE_PT = 12
DEFAULT_MIN_COMFORTABLE_PT = 16


def check_slide_font_sizes(slide, slide_idx=0, thresholds=None):
    """Check all text shapes on a slide for font size issues.

    Returns list of warning dicts:
      [{"slide": N, "shape": role, "font_size": pt, "threshold": pt, "severity": "warn|error", "message": "..."}]
    """
    if thresholds is None:
        thresholds = {
            "error": 8,   # Below this: almost unreadable
            "warn": 12,   # Below this: hard to read
            "notice": 16, # Below this: comfortable but small
        }
    issues = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is None:
                    continue
                pt = run.font.size.pt
                if pt < thresholds.get("error", 8):
                    issues.append({
                        "slide": slide_idx,
                        "shape": getattr(shape, "name", "unknown"),
                        "font_size": round(pt, 1),
                        "threshold": thresholds["error"],
                        "severity": "error",
                        "message": f"Slide {slide_idx}: font size {pt:.1f}pt is below error threshold ({thresholds['error']}pt). Text may be unreadable."
                    })
                elif pt < thresholds.get("warn", 12):
                    issues.append({
                        "slide": slide_idx,
                        "shape": getattr(shape, "name", "unknown"),
                        "font_size": round(pt, 1),
                        "threshold": thresholds["warn"],
                        "severity": "warn",
                        "message": f"Slide {slide_idx}: font size {pt:.1f}pt is below readable threshold ({thresholds['warn']}pt). Consider reducing content or enlarging the area."
                    })
                elif pt < thresholds.get("notice", 16):
                    issues.append({
                        "slide": slide_idx,
                        "shape": getattr(shape, "name", "unknown"),
                        "font_size": round(pt, 1),
                        "threshold": thresholds["notice"],
                        "severity": "notice",
                        "message": f"Slide {slide_idx}: font size {pt:.1f}pt is smaller than comfortable ({thresholds['notice']}pt)."
                    })
    return issues


def check_presentation(prs, thresholds=None):
    """Check all slides in a presentation."""
    all_issues = []
    for i, slide in enumerate(prs.slides, start=1):
        issues = check_slide_font_sizes(slide, slide_idx=i, thresholds=thresholds)
        all_issues.extend(issues)
    return all_issues


def print_report(issues):
    """Print formatted warning report."""
    if not issues:
        print("  Font size guard: All text is comfortably readable.")
        return
    errors = [i for i in issues if i["severity"] == "error"]
    warns = [i for i in issues if i["severity"] == "warn"]
    notices = [i for i in issues if i["severity"] == "notice"]
    if errors:
        print(f"  Font size guard: {len(errors)} ERROR(s) — text may be unreadable!")
        for i in errors[:3]:
            print(f"    {i['message']}")
        if len(errors) > 3:
            print(f"    ... and {len(errors) - 3} more")
    if warns:
        print(f"  Font size guard: {len(warns)} WARNING(s) — text is hard to read.")
        for i in warns[:3]:
            print(f"    {i['message']}")
        if len(warns) > 3:
            print(f"    ... and {len(warns) - 3} more")
    if notices:
        print(f"  Font size guard: {len(notices)} notice(s) — text is smaller than comfortable.")


def get_recommendation(issue):
    """Return human-readable recommendation for a font size issue."""
    recs = {
        "error": "Recommendation: Reduce content, split into multiple slides, or enlarge the text area.",
        "warn": "Recommendation: Consider trimming bullet points or increasing region height/width.",
        "notice": "Recommendation: Fine-tune if presenting to a large audience.",
    }
    return recs.get(issue["severity"], "")
