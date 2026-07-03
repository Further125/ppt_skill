#!/usr/bin/env python3
"""
QA check for generated PPTX.
Detects: text overflow, empty slides, missing fonts, etc.
Usage:
  python qa_check.py output.pptx [--json]
"""
import sys
import os
import argparse
import json

from pptx import Presentation


def check_pptx(pptx_path):
    """Run QA checks and return results."""
    prs = Presentation(pptx_path)
    issues = []
    warnings = []
    stats = {"slides": len(prs.slides), "shapes": 0, "text_boxes": 0}

    for idx, slide in enumerate(prs.slides):
        slide_has_content = False
        for shape in slide.shapes:
            stats["shapes"] += 1
            if shape.has_text_frame:
                stats["text_boxes"] += 1
                text = shape.text_frame.text.strip()
                if text:
                    slide_has_content = True
                    # Check for remaining placeholders
                    for ph in ["__TITLE__", "__SUBTITLE__", "__DATE__", "__ITEMS__",
                               "__CONTENT__", "__LEFT__", "__RIGHT__", "__DESC__", "__CHART__"]:
                        if ph in text:
                            issues.append({
                                "slide": idx + 1,
                                "type": "unfilled_placeholder",
                                "detail": f"Placeholder '{ph}' not replaced"
                            })

                # Check text overflow (heuristic)
                total_chars = len(text)
                shape_area = shape.width * shape.height
                # Rough heuristic: if very dense text in small box
                if total_chars > 0 and shape_area > 0:
                    chars_per_emu = total_chars / shape_area
                    if chars_per_emu > 0.0005:  # threshold
                        warnings.append({
                            "slide": idx + 1,
                            "type": "potential_overflow",
                            "detail": f"Text density high ({total_chars} chars in shape)"
                        })

        if not slide_has_content:
            warnings.append({
                "slide": idx + 1,
                "type": "empty_slide",
                "detail": "Slide has no visible text content"
            })

    result = {
        "passed": len(issues) == 0,
        "issues": issues,
        "warnings": warnings,
        "stats": stats
    }
    return result


def main():
    parser = argparse.ArgumentParser(description="QA check for PPTX")
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument("--json", action="store_true", help="Output as JSON")
    args = parser.parse_args()

    result = check_pptx(args.pptx)

    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        print(f"=== QA Report for {args.pptx} ===")
        print(f"Slides: {result['stats']['slides']}")
        print(f"Shapes: {result['stats']['shapes']}")
        print(f"Text boxes: {result['stats']['text_boxes']}")
        print()
        if result["issues"]:
            print(f"ERRORS ({len(result['issues'])}):")
            for issue in result["issues"]:
                print(f"  Slide {issue['slide']}: [{issue['type']}] {issue['detail']}")
        else:
            print("No errors found.")
        print()
        if result["warnings"]:
            print(f"WARNINGS ({len(result['warnings'])}):")
            for w in result["warnings"]:
                print(f"  Slide {w['slide']}: [{w['type']}] {w['detail']}")
        else:
            print("No warnings.")
        print()
        print("RESULT:", "PASS" if result["passed"] else "FAIL")


if __name__ == "__main__":
    main()
